#!/usr/bin/env python3
"""Build the 40-min PhD viva deck on the user's pushed template.pptx.

The template (example-phd-passes/template.pptx, 15 layout archetypes) is used as
the visual start point. We clone its layouts, rewrite the text in Roca Two (already
applied to the template), and fill it with the full thesis content, keeping the same
section structure and the per-contribution RQ / Objectives / Methodology / Results /
Findings framework used in the earlier example-style deck.

Output: viva/MOUAD_LOUHICHI_VIVA_40min.pptx
"""
import copy, os, re, zipfile
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

HERE = os.path.dirname(os.path.abspath(__file__))
SRC = os.path.join(HERE, "..", "example-phd-passes", "template.pptx")
OUT = os.path.join(HERE, "MOUAD_LOUHICHI_VIVA_40min.pptx")

MONTS = "Roca Two"   # font already applied throughout the template

# ---- read the template's accent colour from its theme ----
def _template_accent():
    with zipfile.ZipFile(SRC) as z:
        theme = [n for n in z.namelist() if n.startswith('ppt/theme/theme')][0]
        td = z.read(theme).decode('utf-8')
    m = re.search(r'<a:accent1>.*?<a:srgbClr val="([0-9A-Fa-f]{6})"', td)
    if m:
        return RGBColor.from_string(m.group(1))
    return RGBColor(0x0E, 0x7C, 0x7B)
ACCENT = _template_accent()
BODY   = RGBColor(0x22, 0x30, 0x3C)
SUBTLE = RGBColor(0x55, 0x5E, 0x6B)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)

from build_colleague_style import _ensure_figs, add_figure, add_table  # reuse figure gen + table

def _walk(shapes):
    """Yield every shape, descending into groups (template titles often live inside groups)."""
    for sh in shapes:
        yield sh
        if sh.shape_type == 6:  # group
            yield from _walk(sh.shapes)

def _find(shapes, name):
    for sh in _walk(shapes):
        if sh.name == name:
            return sh
    return None

_next_shape_id = [1000]  # unique cNvPr id counter for cloned shapes

def clone_slide(prs, source):
    dest = prs.slides.add_slide(source.slide_layout)
    for shp in list(dest.shapes):
        shp._element.getparent().remove(shp._element)
    for shp in source.shapes:
        new_el = copy.deepcopy(shp._element)
        # reassign every cNvPr id in the (possibly nested) shape so cloned
        # slides never collide with the originals or with each other
        for cNvPr in new_el.iter(qn('p:cNvPr')):
            _next_shape_id[0] += 1
            cNvPr.set('id', str(_next_shape_id[0]))
        dest.shapes._spTree.append(new_el)
    for shp in dest.shapes:
        if shp.shape_type == 13:
            for blip in shp._element.findall('.//' + qn('a:blip')):
                rid = blip.get(qn('r:embed'))
                if rid is None:
                    continue
                img_part = source.part.related_part(rid)
                new_rid = dest.part.relate_to(img_part, 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/image')
                blip.set(qn('r:embed'), new_rid)
    return dest

def set_notes(slide, text):
    if text:
        slide.notes_slide.notes_text_frame.text = text

# ---------------------------------------------------------------------------
def add_textbox(slide, x, y, w, h, lines, size=18, color=BODY, bold=False,
                align=PP_ALIGN.LEFT, spacing=1.05, space_after=6, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    for i, text in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.line_spacing = spacing; p.space_after = Pt(space_after)
        r = p.add_run(); r.text = text
        r.font.name = MONTS; r.font.size = Pt(size); r.font.bold = bold
        r.font.color.rgb = color
    return box

def add_bullets(slide, x, y, w, h, items, size=18, color=BODY, accent=None,
                space=9, anchor=MSO_ANCHOR.TOP, spacing=1.05):
    if accent is None:
        accent = ACCENT
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    for i, it in enumerate(items):
        if isinstance(it, str):
            text, level, bold = it, 0, False
        elif len(it) == 2:
            text, level = it; bold = False
        else:
            text, level, bold = it
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(space if level == 0 else max(space - 3, 2))
        p.line_spacing = spacing
        if level == 0:
            r0 = p.add_run(); r0.text = "\u2022  "
            r0.font.name = MONTS; r0.font.size = Pt(size); r0.font.bold = True
            r0.font.color.rgb = accent
            r = p.add_run(); r.text = text
            r.font.name = MONTS; r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color
        else:
            p.level = 1
            r0 = p.add_run(); r0.text = "\u2013  "
            r0.font.name = MONTS; r0.font.size = Pt(size - 1); r0.font.color.rgb = accent
            r = p.add_run(); r.text = text
            r.font.name = MONTS; r.font.size = Pt(size - 1); r.font.bold = bold; r.font.color.rgb = color
    return box

def _is_junk(t):
    return (t.startswith("Lorem") or "Borcelle" in t or "reallygreatsite" in t
            or "Presentation By Donna" in t or "www." in t
            or t in ("Website :", "Social Media :", "Phone Number :", "Email Address :",
                     "@reallygreatsite", "+123-456-7890"))

def _find_text(shapes, text, exact=True):
    for sh in _walk(shapes):
        if sh.has_text_frame:
            t = sh.text_frame.text.strip()
            if (t == text) if exact else (text in t):
                return sh
    return None

def _text_shapes(shapes):
    return [sh for sh in _walk(shapes) if sh.has_text_frame and sh.text_frame.text.strip()]

def _find_title(shapes):
    """Title = text shape with the largest font (template titles are always the biggest)."""
    best, bestsz = None, -1
    for sh in _text_shapes(shapes):
        if _is_junk(sh.text_frame.text.strip()):
            continue
        sz = 0
        for p in sh.text_frame.paragraphs:
            for r in p.runs:
                if r.font.size:
                    sz = max(sz, r.font.size.pt)
        if sz > bestsz:
            bestsz, best = sz, sh
    return best

def _title_size(t):
    n = len(t)
    if n <= 16: return 36
    if n <= 28: return 30
    if n <= 42: return 26
    if n <= 60: return 22
    return 18

def _set_title(tsh, title):
    if tsh is None or not tsh.has_text_frame:
        return
    tf = tsh.text_frame
    tf.word_wrap = True
    paras = list(tf.paragraphs)
    if paras and paras[0].runs:
        r = paras[0].runs[0]
        r.text = title
        for rr in paras[0].runs[1:]:
            rr.text = ""
        r.font.size = Pt(_title_size(title))
    else:
        tf.text = title
    for p in paras[1:]:
        p._p.getparent().remove(p._p)

def _set_text(slide, old, new):
    sh = _find_text(slide.shapes, old, exact=True)
    if sh is not None:
        sh.text_frame.text = new
    return sh

def _set_text_if(slide, pred, new):
    for sh in _walk(slide.shapes):
        if sh.has_text_frame and pred(sh.text_frame.text.strip()):
            sh.text_frame.text = new
            return sh
    return None

def _strip_except(slide, keep):
    # use stable cNvPr id (python-pptx hands out fresh wrapper objects each access,
    # so object identity / id() cannot be used to match shapes)
    keep_ids = set(x.shape_id for x in keep if x is not None)
    for sh in list(_walk(slide.shapes)):
        if sh.shape_id in keep_ids:
            continue
        if sh.has_text_frame and sh.text_frame.text.strip():
            sh._element.getparent().remove(sh._element)
        elif sh.shape_type == 13:
            sh._element.getparent().remove(sh._element)

# ---------------------------------------------------------------------------
# clone an archetype, find its title by largest font, rewrite it, drop all other
# text/pictures, then let `builder` add content.
def content_slide(prs, arch_idx, title_shape, title, builder, note=None):
    s = clone_slide(prs, prs.slides[arch_idx])
    tsh = _find_title(s.shapes)
    _set_title(tsh, title)
    tsh_id = tsh.shape_id if tsh is not None else None
    for sh in list(_walk(s.shapes)):
        if tsh_id is not None and sh.shape_id == tsh_id:
            continue
        if sh.has_text_frame and sh.text_frame.text.strip():
            sh._element.getparent().remove(sh._element)
        elif sh.shape_type == 13:
            sh._element.getparent().remove(sh._element)
    builder(s)
    set_notes(s, note)
    return s

def section_slide(prs, title, note=None):
    s = clone_slide(prs, prs.slides[3])  # slide 3 = section divider
    tsh = _find_title(s.shapes)
    _set_title(tsh, title)
    tsh_id = tsh.shape_id if tsh is not None else None
    for sh in list(_walk(s.shapes)):
        if tsh_id is not None and sh.shape_id == tsh_id:
            continue
        if sh.has_text_frame and sh.text_frame.text.strip():
            sh._element.getparent().remove(sh._element)
        elif sh.shape_type == 13:
            sh._element.getparent().remove(sh._element)
    set_notes(s, note)
    return s

def main():
    _ensure_figs()
    prs = Presentation(SRC)

    # ---------- 1. TITLE ----------
    s = clone_slide(prs, prs.slides[0])
    ths = _find_title(s.shapes)
    _set_title(ths, "Cooperative Game Theory for Explainable AI in Recommendation Systems:\nA Shapley Framework for Actionable Insight")
    uni = _set_text(s, "Borcelle University", "ENSIAS, Mohammed V University, Rabat")
    name = _set_text(s, "Presentation By Donna Stroupe", "PhD Viva \u2014 Presented by Mouad LOUHICHI")
    sub = _set_text_if(s, lambda t: t.startswith("Lorem"), "Supervisor: Pr. Mohamed LAZAAR")
    keep = [x for x in (ths, uni, name, sub) if x is not None]
    _strip_except(s, keep)
    set_notes(s, "Good morning. Thank you for the time to present and discuss my thesis. "
                  "My thesis is titled Cooperative Game Theory for Explainable AI in Recommendation "
                  "Systems: A Shapley Framework for Actionable Insight, supervised by Professor "
                  "Mohamed Lazaar. It advances one idea: that Shapley-value attribution is not just a "
                  "post-hoc explanation, but a single formal mechanism that explains black-box "
                  "clustering, stays coherent under hierarchical scale, and finally acts as an "
                  "in-training signal inside a recommender.")

    # ---------- 2. ABSTRACT ----------
    s = clone_slide(prs, prs.slides[1])
    ths = _find_title(s.shapes)
    _set_title(ths, "Abstract")
    _strip_except(s, [ths])
    add_textbox(s, 0.7, 2.3, 11.9, 4.0, [
        "Modern recommenders and clustering pipelines have grown remarkably accurate yet remain "
        "opaque, which collides with trust, debugging and regulation (EU AI Act, OECD, GDPR).",
        "This thesis proposes cooperative game theory \u2014 Shapley values \u2014 as a single, shared "
        "attribution perspective spanning clustering explanation and recommendation learning.",
        "C1 explains black-box clustering faithfully (PCA\u2013KMeans\u2013LightGBM\u2013TreeSHAP); C2 scales it to "
        "multi-level, large-scale clustering with cross-level consistency; C3 (DyHuCoG) turns attribution "
        "into an in-training signal inside a hypergraph recommender.",
        "Across wine, Beijing air quality, MovieLens-1M and Amazon-Book, ranking, coverage and diversity "
        "improve together \u2014 evidence that the accuracy\u2013diversity trade-off is negotiable when attribution is "
        "part of the modelling logic.",
    ], size=18, spacing=1.12, space_after=10)
    set_notes(s, "Here is the one-paragraph version: accurate but opaque recommenders meet trustworthy-AI "
                  "requirements; Shapley values provide one principled attribution language for both explaining "
                  "clustering and learning recommendations.")

    # ---------- 3. AGENDA ----------
    s = clone_slide(prs, prs.slides[2])
    agenda = ["1.  Introduction", "2.  Context & Problematic", "3.  Experimental Protocol",
              "4.  Contribution I \u2014 Explainable Clustering", "5.  Contribution II \u2014 Multi-level XAI",
              "6.  Contribution III \u2014 DyHuCoG", "7.  Conclusion & Perspectives"]
    cands = [sh for sh in _text_shapes(s.shapes) if not _is_junk(sh.text_frame.text)]
    cands.sort(key=lambda sh: (sh.top, sh.left))
    keep = []
    for i, sh in enumerate(cands):
        if i < len(agenda):
            sh.text_frame.text = agenda[i]
            keep.append(sh)
        else:
            sh.text_frame.text = ""
    _strip_except(s, keep)
    set_notes(s, "I will walk through five parts: introduction, context and problematic, the shared "
                  "experimental protocol, the three contributions, and the conclusion.")

    # ---------- 4. SECTION: INTRODUCTION ----------
    section_slide(prs, "Introduction",
                  note="Let us begin with the introduction: why explainability is a first-class requirement "
                  "for recommender systems, and what I mean by an actionable insight.")

    # Motivation
    content_slide(prs, 7, "TextBox 2", "Motivation: Three Questions", lambda sl: add_bullets(sl, 0.7, 2.3, 11.9, 4.3, [
        ("Ubiquity \u2014 how do opaque AI systems shape what billions of users see, buy and watch every day?", 0, True),
        ("The Black Box \u2014 why do state-of-the-art recommenders and clustering pipelines remain opaque to users and designers?", 0, True),
        ("Toward Trust \u2014 how can transparency be built into the model, instead of being bolted on afterwards?", 0, True),
        ("Core tension: as models gain expressive power, they lose the transparency needed for trustworthy deployment.", 1, True),
    ], size=19, space=12), note="Three questions frame the work: ubiquity, the black box, and trust. The core "
        "tension is that expressive power trades away transparency.")

    # Actionable insight (two-column builder, defined just below)
    # Motivation already added; Actionable:
    def _actionable(sl):
        add_bullets(sl, 0.7, 2.3, 6.3, 4.3, [
            ("Definition 1.1 (Actionable insight)", 0, True),
            ("An explanation is actionable when it identifies at least one modifiable factor whose change is associated with a specifiable change in model output\u2026", 1),
            ("\u2026 and that factor is expressible in the semantic vocabulary of the task domain.", 1),
            ("Accessibility means the domain's own terms: a physicochemical variable, a pollution indicator, or a preference signal \u2014 not an opaque latent code.", 1),
            ("Why: an explanation that identifies a modifiable driver supports intervention, not merely description.", 0, True),
        ], size=18)
        add_figure(sl, os.path.join(HERE, "_figs", "actionable.png"), 7.1, 2.4, w=5.5)

    content_slide(prs, 4, "TextBox 25", "Actionable Insight \u2014 the Definition", _actionable,
                  note="We frame the whole thesis around actionable insight: an explanation that identifies a modifiable, "
                  "semantically meaningful driver supports intervention, not just description.")

    # Research context + evolution figure
    def _research(sl):
        add_bullets(sl, 0.7, 2.3, 6.3, 4.3, [
            ("Recommenders evolved from similarity filters to complex representation learning on sparse, high-dimensional, dynamic data.", 0),
            ("\u2192 MF \u2192 neural CF \u2192 graph CNNs \u2192 hypergraph: each step improved ranking but intensified the interpretability deficit.", 1),
            ("Undermines user trust; constrains debugging; collides with EU AI Act, OECD, GDPR.", 1),
            ("EU AI Act (Art. 13) requires high-risk systems to provide explanations in human-understandable terms.", 1, True),
        ], size=18)
        add_figure(sl, os.path.join(HERE, "_figs", "evolution.png"), 7.1, 2.4, w=5.5)
    content_slide(prs, 5, "TextBox 29", "Research Context", _research,
                  note="The context is the progression from simple to hypergraph recommenders; each step raised "
                  "expressiveness while lowering transparency, which regulation now penalises.")

    # AI recsys examples + figure
    def _examples(sl):
        add_bullets(sl, 0.7, 2.3, 6.3, 4.3, [
            ("AI recommendation mediates what billions see, buy and listen to every day:", 0, True),
            ("Streaming & video: Netflix, Prime Video, YouTube.", 1),
            ("Music: Spotify, YouTube Music, Deezer.", 1),
            ("Shopping & e-commerce: Amazon, AliExpress, Noon.", 1),
            ("Social & feeds: TikTok, Instagram, LinkedIn.", 1),
            ("Maps, ads, search: Google Maps, Uber, Google Ads.", 1),
            ("Accurate yet opaque \u2014 the very gap our thesis targets.", 0, True),
        ], size=18)
        add_figure(sl, os.path.join(HERE, "_figs", "recsys_examples.png"), 7.1, 2.4, w=5.5)
    content_slide(prs, 5, "TextBox 29", "AI-Powered Recommendation Systems Around Us", _examples,
                  note="Recommendation is infrastructure for everyday decisions; its ubiquity plus opacity is what "
                  "makes explainability a first-class requirement.")

    # ---------- 5. SECTION: CONTEXT & PROBLEMATIC ----------
    section_slide(prs, "Context & Problematic",
                  note="Now the problem this thesis addresses: the paradigm landscape, the structuring limitations, "
                  "the research questions, and how the three contributions map to them.")

    # Recommendation paradigms + content-based figure
    def _paradigms(sl):
        add_bullets(sl, 0.7, 2.3, 6.3, 4.3, [
            ("Collaborative filtering \u2013 users who behaved similarly will value similar items (user-/item-based).", 0),
            ("Content-based \u2013 recommends items sharing attributes with a user profile.", 0),
            ("Hybrid \u2013 combines collaborative and content signals.", 0),
            ("Matrix factorisation \u2013 R \u2248 PQ\u1d40: compact but opaque latent factors.", 0),
            ("Graph-based \u2013 interaction graph with neighbourhood propagation (LightGCN, hypergraph).", 0),
        ], size=18)
        add_figure(sl, os.path.join(HERE, "_figs", "content_based.png"), 7.1, 2.4, w=5.5)
    content_slide(prs, 5, "TextBox 29", "Recommendation & Clustering Paradigms", _paradigms,
                  note="Quick orientation: collaborative, content-based, hybrid, matrix factorisation, and graph/hypergraph "
                  "methods \u2014 each strengthens modelling but complicates interpretation in a specific way.")

    # Limitations
    content_slide(prs, 4, "TextBox 25", "Limitations of Classical Recommenders & Unsupervised Models",
                  lambda sl: add_bullets(sl, 0.7, 2.3, 11.9, 4.3, [
        ("Data sparsity & scalability \u2014 the user\u2013item matrix is overwhelmingly empty.", 0),
        ("Cold-start \u2014 new users and items are structurally disadvantaged.", 0),
        ("Popularity bias & lack of diversity \u2014 exposure begets interaction, begets exposure.", 0),
        ("Absence of interpretability \u2014 the most fundamental limit.", 0, True),
        ("For clustering: methods favour local OR global explanation, not both; struggle to scale; explanations rarely stay coherent across resolutions.", 1),
    ], size=19, space=12),
                  note="Four classical limitations; the most fundamental \u2014 and the one this thesis targets \u2014 is the "
                  "absence of interpretability.")

    # Problem statement
    content_slide(prs, 4, "TextBox 25", "Three Structuring Limitations (Problem Statement)",
                  lambda sl: add_bullets(sl, 0.7, 2.3, 11.9, 4.3, [
        ("1. Lack of explainability \u2014 complex models remain hard to interpret faithfully and actionably.", 0, True),
        ("2. Difficulty of scaling \u2014 local explanations do not transfer naturally to hierarchical structures or large datasets.", 0, True),
        ("3. Weak integration into learning \u2014 most explanations stay post-hoc and do not shape model dynamics or the accuracy\u2013diversity\u2013context trade-off.", 0, True),
        ("Thesis gap: the literature still lacks a single cooperative-attribution framework that explains clustering faithfully, stays coherent under hierarchy, and operates as an in-training signal in recommendation.", 1, True),
    ], size=19, space=12),
                  note="Three structuring problems; the thesis gap is a single cooperative-attribution framework spanning "
                  "explanation, hierarchy and learning.")

    # Research questions
    content_slide(prs, 7, "TextBox 2", "Research Questions (RQ1\u2013RQ5) and Overall Aim",
                  lambda sl: add_bullets(sl, 0.7, 2.3, 11.9, 4.3, [
        ("Aim: develop, justify and evaluate a cooperative-game-theoretic perspective for XAI in clustering and recommendation.", 0, True),
        ("RQ1 \u00b7 How can Shapley values explain black-box clustering faithfully at instance and cluster level?", 0),
        ("RQ2 \u00b7 How can this extend to large-scale, hierarchical clustering without losing tractability or consistency?", 0),
        ("RQ3 \u00b7 Can cooperative attribution move beyond post-hoc and enter the learning dynamics of graph recommenders?", 0),
        ("RQ4 \u00b7 Can a recommender jointly optimise ranking accuracy, context and diversity when importance is estimated by a cooperative-game utility?", 0),
        ("RQ5 \u00b7 What emerges when clustering explanation and recommendation learning are two stages of one cooperative-game perspective?", 0),
    ], size=18, space=10),
                  note="The overall aim is a cooperative-game-theoretic perspective for XAI in clustering and recommendation; "
                  "five research questions form the spine.")

    # Contributions overview
    content_slide(prs, 7, "TextBox 2", "The Three Contributions",
                  lambda sl: add_bullets(sl, 0.7, 2.3, 11.9, 4.3, [
        ("C1 \u2013 Explainable black-box clustering: PCA\u2013KMeans\u2013LightGBM\u2013TreeSHAP pipeline.", 0, True),
        ("\u2192 Wine Quality: faithful instance- and cluster-level feature attribution.", 1),
        ("C2 \u2013 Enhanced multi-level XAI for large-scale clustering with cross-level SHAP aggregation.", 0, True),
        ("\u2192 Beijing Air Quality: hierarchical attribution consistency under scale.", 1),
        ("C3 \u2013 DyHuCoG: Dynamic Hypergraph Cooperative Game for preference-aware recommendation.", 0, True),
        ("\u2192 MovieLens-1M & Amazon-Book: preference-aware Monte Carlo Shapley as an in-training signal.", 1),
        ("Thesis claim: cooperative game theory functions as a shared attribution perspective for explanation, optimisation and intervention.", 2, True),
    ], size=18, space=10),
                  note="Three contributions, one thread: C1 establishes Shapley explanation for clustering; C2 scales it to "
                  "hierarchy; C3 turns it into an in-training signal. The thesis claim is that cooperative game theory is a "
                  "shared attribution perspective.")

    # ---------- 6. SECTION: EXPERIMENTAL PROTOCOL ----------
    section_slide(prs, "Experimental Protocol",
                  note="Before the contributions, the shared experimental setup: datasets, splitting, baselines, metrics, hardware.")

    # Datasets table
    content_slide(prs, 12, "TextBox 6", "Datasets Used Throughout the Thesis",
                  lambda sl: add_table(sl, 0.7, 2.3, 11.9, 3.9,
             [["Dataset", "Scale", "Type", "Role"],
              ["Wine Quality (vinho verde)", "4,898 \u00d7 11", "Tabular, numeric", "C1 \u2013 single-level clustering"],
              ["Beijing Multi-Site Air Quality", "383,585 \u00d7 11", "Tabular, pollutant + meteorology", "C2 \u2013 multi-level clustering"],
              ["MovieLens-1M", "6,040 u / 3,706 i / 1.0M int", "Implicit feedback (0.0447)", "C3 \u2013 DyHuCoG"],
              ["Amazon-Book", "52,643 u / 91,599 i / 3.0M int", "Implicit feedback (0.0006)", "C3 \u2013 DyHuCoG"]],
             col_ratios=[4, 2.6, 2.8, 3.4], header_bg=ACCENT),
                  note="Two clustering datasets and two recommendation datasets, each chosen so attribution lands in a "
                  "vocabulary a domain expert can act on.")

    # Dataset characteristics
    content_slide(prs, 12, "TextBox 6", "Dataset Characteristics & Why Each Was Chosen",
                  lambda sl: add_bullets(sl, 0.7, 2.3, 11.9, 4.3, [
        ("Wine Quality \u2013 11 interpretable physicochemical features; small, dense, chemically correlated \u2192 tests attribution in a trusted semantic space.", 0),
        ("Beijing Air Quality \u2013 11 pollutant + meteorological variables, 383,585 hourly records \u2192 tests scale, noise and temporal structure.", 0),
        ("MovieLens-1M \u2013 ~1.0M interactions, density 0.0447 \u2192 benchmark-standard recommendation with established baselines.", 0),
        ("Amazon-Book \u2013 ~3.0M interactions, density 0.0006 \u2192 deliberately sparse to stress-test Shapley weighting when signal is weak.", 0),
        ("The 0.0447 vs 0.0006 density contrast is intentional: gains are largest where data are weakest.", 0, True),
    ], size=18, space=11),
                  note="Each dataset was chosen so attribution lands in a vocabulary a domain expert can act on; the density "
                  "contrast isolates whether cooperative attribution helps most when interaction data are scarce.")

    # Splitting
    content_slide(prs, 12, "TextBox 6", "Data Splitting & Preprocessing",
                  lambda sl: add_bullets(sl, 0.7, 2.3, 11.9, 4.3, [
        ("Clustering: five-fold cross-validation for surrogate/attribution stability.", 0),
        ("Recommendation: user-level, temporal split \u2014 70% train / 10% val / 20% test.", 0),
        ("Leave-one-out: the latest test positive per user is the target, ranked against negatives.", 0),
        ("Implicit conversion: MovieLens-1M ratings > 3 treated as positive.", 0),
        ("Popularity-aware negative sampling: q(i) \u221d f_i^\u03b7 for harder contrasts.", 0),
        ("Reproducibility: seeds {42,43,44,45,46}; early-stopping patience 20.", 0),
    ], size=18, space=10),
                  note="Splitting is designed against leakage; recommendation uses a user-level temporal holdout with "
                  "leave-one-out evaluation.")

    # Baselines & metrics + figure
    def _metrics(sl):
        add_bullets(sl, 0.7, 2.3, 6.3, 4.3, [
            ("Clustering benchmark: LIME-based surrogate explanation pipeline.", 0, True),
            ("Recommendation benchmarks: MF, NCF, LightGCN, RecDCL, HCCF, HPCF (strongest reference).", 0, True),
            ("Ranking: Precision@K, Recall@K, NDCG@20 (principal).", 0),
            ("Beyond-accuracy: Coverage, ILD, Novelty, Fairness.", 0),
            ("Clustering quality: Silhouette, Davies\u2013Bouldin.", 0),
        ], size=18)
        add_figure(sl, os.path.join(HERE, "_figs", "metrics.png"), 7.1, 2.4, w=5.5)
    content_slide(prs, 12, "TextBox 6", "Baselines & Evaluation Metrics", _metrics,
                  note="Baselines span classical, neural, graph and hypergraph methods; metrics cover ranking, beyond-accuracy "
                  "and clustering validity. ILD is built into the DyHuCoG coalition utility.")

    # Hardware
    content_slide(prs, 12, "TextBox 6", "Hardware & Software",
                  lambda sl: add_bullets(sl, 0.7, 2.3, 11.9, 4.3, [
        ("CPU: Intel Core i9-14900K, 24 cores \u2014 clustering, preprocessing, data loading.", 0),
        ("GPU: NVIDIA GeForce RTX 4090, 24 GB \u2014 DyHuCoG training & inference.", 0),
        ("RAM: 48 GB; Storage: 2 TB SSD.", 0),
        ("Python 3.8; scikit-learn, LightGBM, SHAP, PyTorch 2.0.1, NumPy/SciPy/pandas.", 0),
        ("Altair for interactive SHAP visualisation; metrics at K \u2208 {5, 10, 20}.", 0),
    ], size=18, space=11),
                  note="The hardware explains some runtime figures; everything stays within ordinary academic compute.")

    # ---------- 7. SECTION: CONTRIBUTION I ----------
    section_slide(prs, "Contribution I \u2014 Explainable Black-Box Clustering",
                  note="Contribution I explains black-box clustering with Shapley values. It answers RQ1 in four steps: "
                  "objectives, methodology, results, findings.")

    def c1_obj(sl):
        add_bullets(sl, 0.7, 2.3, 11.9, 4.3, [
            ("Gap: Shapley explanation is dominant in supervised tasks, but clustering remains comparatively under-explained.", 0, True),
            ("Existing clustering-interpretability methods privilege local or global explanation, not both.", 1),
            ("They often fail to scale or preserve coherence across clusters.", 1),
            ("Objectives: build a pipeline yielding cluster-level explanation while preserving feature-level attribution.", 0, True),
            ("Preserve the semantics of the original feature space; justify why Shapley is better than LIME.", 1),
        ], size=18, space=10)
    content_slide(prs, 7, "TextBox 2", "Contribution I \u2014 Objectives", c1_obj,
                  note="Why start with clustering? It is the hardest test bed for attribution because the model creates its "
                  "own structure. Objectives: cluster-level + feature-level explanation, in the original semantics, justified vs LIME.")

    def c1_rq(sl):
        add_bullets(sl, 0.7, 2.3, 11.9, 4.3, [
            ("RQ1 \u00b7 How can Shapley values explain black-box clustering faithfully at instance and cluster level?", 0, True),
            ("Objective O1 \u00b7 A pipeline that yields cluster-level explanation while keeping feature-level attribution.", 1),
            ("Objective O2 \u00b7 Preserve the semantics of the original feature space, not a reduced latent space.", 1),
            ("Objective O3 \u00b7 Justify Shapley over an ad-hoc surrogate such as LIME.", 1),
        ], size=18, space=10)
    content_slide(prs, 7, "TextBox 2", "Contribution I \u2014 Research Question & Objectives", c1_rq,
                  note="RQ1 and its three objectives: pipeline, semantics, and a justificatory argument for Shapley over LIME.")

    def c1_vf(sl):
        add_bullets(sl, 0.7, 2.3, 11.9, 4.3, [
            ("Player set N = F \u2014 each feature is a player.", 0, True),
            ("Value function v(S) = Silhouette( KMeans(X_S, k*) ) \u2014 how well data cluster using only features in S.", 0),
            ("A feature's Shapley value = its expected marginal contribution to clustering quality over all coalition orders.", 0),
            ("Why Silhouette: bounded, normalised, semantically intuitive.", 0),
            ("Direct evaluation for every coalition is intractable \u2192 we need a bridge.", 0, True),
        ], size=18, space=10)
    content_slide(prs, 5, "TextBox 29", "Contribution I \u2014 Value Function (Cooperative Game)", c1_vf,
                  note="Features are players; the value of a coalition is the Silhouette of K-Means on that feature subset. "
                  "A feature's Shapley value is its expected marginal contribution, but evaluating every subset is intractable.")

    def c1_sur(sl):
        add_bullets(sl, 0.7, 2.3, 6.3, 4.3, [
            ("Once K-Means produces cluster labels, train a LightGBM multiclass surrogate to predict them from original features.", 0, True),
            ("Apply TreeSHAP to the surrogate \u2014 fast, exact tree-based attribution in the original semantic feature space.", 0),
            ("Direct TreeSHAP on K-Means is impossible (it explains tree models, not centroids).", 1),
            ("Validity condition: surrogate fidelity is high (macro-F1 \u2248 0.82).", 1, True),
        ], size=18)
        add_figure(sl, os.path.join(HERE, "_figs", "c1_pipeline_paper.png"), 7.1, 2.4, w=5.4)
    content_slide(prs, 4, "TextBox 25", "Contribution I \u2014 Surrogate + TreeSHAP", c1_sur,
                  note="The bridge: a LightGBM surrogate predicts cluster labels from original features, then TreeSHAP attributes "
                  "in the original semantic space. Fidelity (macro-F1 ~0.82) is the validity condition.")

    def c1_pipe(sl):
        add_bullets(sl, 0.7, 2.3, 6.3, 4.3, [
            ("Stage 1 \u2013 PCA: stabilise geometry + visual diagnostic (NOT the explanatory space).", 0),
            ("Stage 2 \u2013 K-Means++ with multi-criteria k selection (elbow, Silhouette, Davies\u2013Bouldin).", 0),
            ("Stage 3 \u2013 LightGBM surrogate trained on original features to predict cluster labels.", 0),
            ("Stage 4 \u2013 TreeSHAP attribution in the original feature space.", 0),
            ("Stage 5 \u2013 Global importance, cluster-specific profiles, local force plots.", 0),
        ], size=18)
        add_figure(sl, os.path.join(HERE, "_figs", "c1_pipeline_paper.png"), 7.1, 2.4, w=5.4)
    content_slide(prs, 4, "TextBox 25", "Contribution I \u2014 Pipeline (5 stages)", c1_pipe,
                  note="Five stages; PCA is only a visual/computational diagnostic, not the explanatory space. TreeSHAP scales "
                  "with tree count, not exponentially in features.")

    def c1_k(sl):
        add_bullets(sl, 0.7, 2.3, 11.9, 4.3, [
            ("Multi-criteria evaluation across k \u2208 {2..10} using elbow, Silhouette, Davies\u2013Bouldin.", 0),
            ("We select k* = 3 \u2014 even though it is NOT geometrically optimal:", 0, True),
            ("k = 2: Silhouette 0.214, Davies\u2013Bouldin 1.775 (better separation).", 1),
            ("k = 3: Silhouette 0.144, Davies\u2013Bouldin 2.097 (weaker separation).", 1),
            ("Why: three clusters give a semantically richer oenological partition \u2192 more actionable.", 1, True),
            ("Note: the Silhouette ~0.63 belongs to Beijing (C2), not this wine partition.", 1),
        ], size=18, space=9)
    content_slide(prs, 4, "TextBox 25", "Contribution I \u2014 Choosing k", c1_k,
                  note="We select k=3 deliberately though it is not geometrically optimal, because three clusters support three "
                  "distinct chemically meaningful narratives \u2014 far more actionable.")

    def c1_global(sl):
        add_bullets(sl, 0.7, 2.3, 11.9, 4.3, [
            ("Global SHAP ranking (high \u2192 low):", 0, True),
            ("density \u2192 pH \u2192 fixed acidity \u2192 sulfur-dioxide \u2192 alcohol", 1),
            ("Dominant drivers relate to structure, preservation, and sensory balance.", 1),
            ("This is NOT an arbitrary classifier artefact \u2014 it recovers a chemically interpretable hierarchy.", 1, True),
        ], size=18, space=10)
    content_slide(prs, 4, "TextBox 25", "Contribution I \u2014 Global SHAP Ranking (Wine)", c1_global,
                  note="The global ranking is dominated by density, pH, fixed acidity, sulfur dioxide and alcohol \u2014 precisely the "
                  "variables an oenologist would point to. This is the strongest evidence the pipeline is faithful.")

    def c1_prof(sl):
        add_bullets(sl, 0.7, 2.3, 11.9, 4.3, [
            ("Three clusters show distinct explanatory signatures.", 0, True),
            ("Cluster 0 \u2013 density + sulfur-dioxide-related variables.", 1),
            ("Cluster 1 \u2013 acidity and pH-related effects.", 1),
            ("Cluster 2 \u2013 a different balance of acidity, alcohol, and related chemical attributes.", 1),
            ("The same small set of variables recurs across clusters, with different relative weights within each.", 1),
        ], size=18, space=10)
    content_slide(prs, 4, "TextBox 25", "Contribution I \u2014 Cluster-Specific Profiles", c1_prof,
                  note="Each cluster exhibits a distinct signature; the same drivers re-weigh differently. This is the actionable "
                  "insight: a cluster is a distinct, domain-meaningful combination of drivers.")

    def c1_tbl(sl):
        add_table(sl, 0.7, 2.3, 11.9, 3.8,
             [["Criterion", "SHAP (cooperative)", "LIME (local surrogate)"],
              ["Basis", "Cooperative-game marginal contribution", "Local surrogate approximation"],
              ["Local / global", "Both", "Primarily local"],
              ["Theoretical guarantee", "Efficiency, symmetry, null player, additivity", "None equivalent"],
              ["Stability", "Higher when surrogate faithful", "Sensitive to perturbation design"],
              ["Cluster comparison", "Strong", "Limited"]],
             col_ratios=[3, 4.4, 4.4], header_bg=ACCENT)
    content_slide(prs, 12, "TextBox 6", "Contribution I \u2014 SHAP vs LIME", c1_tbl,
                  note="SHAP grounds attribution in a cooperative-game allocation rule satisfying efficiency, symmetry, null-player "
                  "and additivity; LIME has no equivalent guarantee and is sensitive to perturbation design.")

    def c1_ans(sl):
        add_bullets(sl, 0.7, 2.3, 11.9, 4.3, [
            ("Answers to the questions and objectives", 0, True),
            ("RQ1 answered in the affirmative \u2014 faithful, coherent cluster-level explanation from Shapley values.", 1),
            ("O1 met \u00b7 cluster-level explanation while preserving feature-level attribution.", 1),
            ("O2 met \u00b7 attribution returned to the original chemical variables.", 1),
            ("O3 met \u00b7 Shapley grounded in four axioms; LIME has no equivalent guarantee.", 1),
        ], size=18, space=11)
    content_slide(prs, 7, "TextBox 2", "Contribution I \u2014 Answers", c1_ans,
                  note="RQ1 is answered yes; each objective is met. Efficiency holds with respect to the LightGBM output, which is "
                  "why surrogate fidelity is the critical validity condition.")

    def c1_find(sl):
        add_bullets(sl, 0.7, 2.3, 11.9, 4.3, [
            ("Key findings", 0, True),
            ("Cluster-level explanation anchored to individual feature contributions.", 1),
            ("Explanations returned to the original chemical variables, not a latent space.", 1),
            ("Theoretically grounded case for Shapley over LIME.", 1),
            ("Recovers an oenologically interpretable ranking \u2014 density, pH, acidity, sulfur dioxide, alcohol.", 1),
        ], size=18, space=11)
    content_slide(prs, 7, "TextBox 2", "Contribution I \u2014 Key Findings", c1_find,
                  note="C1 answers RQ1: Shapley explains the black-box partition at cluster and instance level, in the original "
                  "chemical variables, grounded in axioms.")

    def c1_lim(sl):
        add_bullets(sl, 0.7, 2.3, 11.9, 4.3, [
            ("Limitations", 0, True),
            ("Fidelity depends on the LightGBM surrogate \u2014 not a direct mechanism of the K-Means geometry.", 1),
            ("Confined to tabular data; no structured, graph, or image input.", 1),
            ("Single-level structure: cannot explain how importance reconfigures between a partition and its sub-partitions.", 1),
            ("The surrogate approximation compresses observation-level variation.", 1),
        ], size=18, space=11)
    content_slide(prs, 7, "TextBox 2", "Contribution I \u2014 Limitations", c1_lim,
                  note="C1 is single-level and surrogate-dependent; these limits define the departure to Contribution II.")

    def c1_take(sl):
        add_bullets(sl, 0.7, 2.3, 11.9, 4.3, [
            ("Takeaways", 0, True),
            ("Shapley attribution is a single, principled lens for explaining an unsupervised partition.", 1),
            ("Keeping attribution in the original feature space is what makes it actionable.", 1),
            ("But real data are rarely single-level: broad regimes contain nested sub-groups.", 1, True),
            ("So the next question is whether this logic survives scale and hierarchy.", 0, True),
        ], size=18, space=11)
    content_slide(prs, 7, "TextBox 2", "Contribution I \u2014 Takeaways", c1_take,
                  note="Shapley attribution works as a single principled lens; but real data are hierarchical, motivating C2.")

    # ---------- 8. SECTION: CONTRIBUTION II ----------
    section_slide(prs, "Contribution II \u2014 Enhanced Multi-Level XAI for Large-Scale Clustering",
                  note="Contribution II scales the explanation logic to multi-level, large-scale clustering. It answers RQ2; "
                  "the central concern is multi-granularity, not merely scale.")

    def c2_obj(sl):
        add_bullets(sl, 0.7, 2.3, 11.9, 4.3, [
            ("Gap: once clustering is multi-level, feature importance must stay interpretable within a cluster, across sub-clusters, and across the hierarchy.", 0, True),
            ("Large-scale data make exact explanation computationally burdensome.", 1),
            ("Flat explanation may be true yet incomplete \u2014 it cannot show how importance changes inside a cluster.", 1),
            ("Objectives: a genuinely multi-level workflow; a formal cross-level consistency argument (Prop. 6.1); validation on a structurally different large-scale dataset.", 0, True),
        ], size=18, space=10)
    content_slide(prs, 7, "TextBox 2", "Contribution II \u2014 Objectives", c2_obj,
                  note="C2 asks whether the C1 logic survives scale and hierarchy. It adds a multi-level workflow, a formal "
                  "cross-level consistency argument, and validation on Beijing air quality.")

    def c2_rq(sl):
        add_bullets(sl, 0.7, 2.3, 11.9, 4.3, [
            ("RQ2 \u00b7 How can this extend to large-scale, hierarchical clustering without losing tractability or consistency?", 0, True),
            ("Objective O1 \u00b7 A genuinely multi-level workflow, not a rerun of the single-level pipeline.", 1),
            ("Objective O2 \u00b7 A formal cross-level consistency argument (Proposition 6.1).", 1),
            ("Objective O3 \u00b7 Validation on a structurally different large-scale dataset.", 1),
        ], size=18, space=10)
    content_slide(prs, 7, "TextBox 2", "Contribution II \u2014 Research Question & Objectives", c2_rq,
                  note="RQ2 and its three objectives: multi-level workflow, Proposition 6.1, and a generality check on Beijing.")

    def c2_arch(sl):
        add_bullets(sl, 0.7, 2.3, 6.3, 4.3, [
            ("Recursive/nested: coarse clustering on the full dataset, then subdivide each cluster.", 0, True),
            ("For each level, train a level-specific surrogate and compute SHAP in the SAME original feature space.", 0),
            ("Cross-level aggregation is NOT a naive average \u2014 it respects cluster size and nesting structure.", 0, True),
            ("Parent-level attribution = an expectation over the explanatory structure of its descendants.", 1),
            ("The hierarchy is a pragmatic analytical device, not a claim of true ontological hierarchy.", 1),
        ], size=18)
        add_figure(sl, os.path.join(HERE, "_figs", "c2_hierarchy.png"), 7.1, 2.4, w=5.5)
    content_slide(prs, 4, "TextBox 25", "Contribution II \u2014 Multi-Level Architecture", c2_arch,
                  note="Recursive coarse-to-fine clustering; per-level surrogate + SHAP; cross-level aggregation respects size and "
                  "nesting. Proposition 6.1 makes differences interpretable, not inconsistent.")

    def c2_prop(sl):
        add_bullets(sl, 0.7, 2.3, 11.9, 4.3, [
            ("Let \u03a6^(l,c)_j = E_{x~c}[|\u03c6_j^(l)(x)|] be expected absolute SHAP importance of feature j at level l in cluster c.", 0),
            ("Let w_c' = |c'| / |c| be the relative size of child c' within parent c.", 0),
            ("For a strict nested hierarchy on a consistent feature space:", 0, True),
            ("\u03a6^(l,c)_j = \u03a3_{c'\\u2208child(c)} w_c' \u00b7 \u03a6^(l+1,c')_j + \u03b5_j", 1, True),
            ("\u03b5_j vanishes under perfect surrogate fidelity. Derived via law of total expectation.", 1),
            ("Does NOT imply explanations are identical across levels \u2014 differences can be interpreted, not dismissed.", 1),
        ], size=18, space=9)
    content_slide(prs, 5, "TextBox 29", "Contribution II \u2014 Proposition 6.1", c2_prop,
                  note="Proposition 6.1: parent expected absolute importance equals the size-weighted sum of children's, plus a "
                  "residual from surrogate mismatch. It makes hierarchical explanation self-consistent and auditable.")

    def c2_conv(sl):
        add_bullets(sl, 0.7, 2.3, 11.9, 4.3, [
            ("Full dataset, k = 3 (strong convergence on multi-criteria evaluation).", 0, True),
            ("Silhouette \u2248 0.63 \u2014 materially stronger separation than wine.", 1),
            ("Davies\u2013Bouldin \u2248 0.55 \u2014 low between-cluster ambiguity.", 1),
            ("PCA projection (2 components) used only for visual inspection.", 1),
            ("Sensitivity: robust to modest variation in k, projection dim, surrogate depth; only low-ranked variables shift.", 1),
        ], size=18, space=10)
    content_slide(prs, 5, "TextBox 29", "Contribution II \u2014 Convergence & Stability", c2_conv,
                  note="On Beijing the partition converges strongly: Silhouette ~0.63, Davies-Bouldin ~0.55. This is where the 0.63 "
                  "figure belongs \u2014 not the wine partition.")

    def c2_global(sl):
        add_bullets(sl, 0.7, 2.3, 11.9, 4.3, [
            ("Global SHAP ranking (high \u2192 low):", 0, True),
            ("temperature \u2192 dew point \u2192 pressure \u2192 CO \u2192 NO2 \u2192 PM10 \u2192 PM2.5", 1),
            ("It is NOT simply pollutant concentrations that matter \u2014 meteorological variables play a structurally central role.", 1),
            ("Temperature, dew point and pressure condition dispersion, trapping, and photochemical behaviour.", 1),
            ("This is the kind of insight flat descriptive summaries often fail to make explicit.", 1),
        ], size=18, space=10)
    content_slide(prs, 5, "TextBox 29", "Contribution II \u2014 Global Ranking (Beijing)", c2_global,
                  note="The dominant features are temperature, dew point and pressure, then pollutants \u2014 meteorology sets the regime. "
                  "Attribution in the original variable space reveals this.")

    def c2_regime(sl):
        add_bullets(sl, 0.7, 2.3, 11.9, 4.3, [
            ("Regime A \u2013 warm photochemical: ozone, temperature, dew point prominent (summer photochemical smog).", 0),
            ("Regime B \u2013 wintertime smog: CO, SO2, PM dominate; low wind speed suppresses dispersion.", 0),
            ("Regime C \u2013 comparatively clean air: favourable meteorology, weak pollutant pushes.", 0),
            ("The framework shows not only that these regimes exist, but which variable combinations define them.", 1, True),
        ], size=18, space=10)
    content_slide(prs, 5, "TextBox 29", "Contribution II \u2014 Regimes", c2_regime,
                  note="Force plots reveal three regimes; the interpretation is which variable combinations define each \u2014 actionable "
                  "for an air-quality analyst.")

    def c2_hier(sl):
        add_bullets(sl, 0.7, 2.3, 11.9, 4.3, [
            ("At the coarse level, temperature and dew point dominate \u2014 they differentiate broad atmospheric regimes.", 0),
            ("Within individual clusters, CO, SO2, PM10, wind speed, pressure, or ozone become more discriminative.", 0),
            ("This change is NOT contradictory \u2014 it is exactly what a multi-level explanation should reveal.", 0, True),
            ("Parent-level story = regime selection. Cluster-level story = variation within a regime.", 1),
            ("A variable can be globally important yet locally uninformative within a sub-cluster.", 1),
        ], size=18, space=9)
    content_slide(prs, 5, "TextBox 29", "Contribution II \u2014 Cross-Level Story", c2_hier,
                  note="At the coarse level temperature/dew point dominate; within clusters the discriminative variables shift. "
                  "Proposition 6.1 lets us interpret these shifts as meaningful structure, not noise.")

    def c2_gen(sl):
        add_bullets(sl, 0.7, 2.3, 11.9, 4.3, [
            ("Wine: small, dense, chemically correlated. Beijing: large, noisy, temporally and meteorologically variable.", 0),
            ("The same explanatory logic remains productive in both \u2192 not tied to one domain-specific peculiarity.", 0, True),
            ("vs SHAP-based clustering literature: Beijing Silhouette \u2248 0.63 vs Gramegna & Giudici credit-risk 0.37.", 1),
            ("LIME comparator: weaker structural coherence, less stable local narratives for hierarchical reasoning.", 1),
        ], size=18, space=10)
    content_slide(prs, 5, "TextBox 29", "Contribution II \u2014 Generalisation", c2_gen,
                  note="The same logic works on both a small correlated dataset and a large noisy one; vs the SHAP-clustering "
                  "literature, Beijing's Silhouette ~0.63 beats credit-risk ~0.37.")

    def c2_ans(sl):
        add_bullets(sl, 0.7, 2.3, 11.9, 4.3, [
            ("Answers to the questions and objectives", 0, True),
            ("RQ2 answered in the affirmative, with bounds \u2014 coherence retained under scale and hierarchy.", 1),
            ("O1 met \u00b7 a multi-granular explanation that does not collapse into a single flat summary.", 1),
            ("O2 met \u00b7 Proposition 6.1 provides a formal cross-level consistency argument.", 1),
            ("O3 met \u00b7 validated on a structurally different large-scale dataset (Beijing air quality).", 1),
        ], size=18, space=10)
    content_slide(prs, 7, "TextBox 2", "Contribution II \u2014 Answers", c2_ans,
                  note="RQ2 answered yes with bounds; O1/O2/O3 met. The honest bound: the model is still an explanation of a "
                  "pre-computed partition, not yet influencing learning \u2014 the bridge to C3.")

    def c2_find(sl):
        add_bullets(sl, 0.7, 2.3, 11.9, 4.3, [
            ("Key findings", 0, True),
            ("Scalable, multi-granular explanation that does not collapse into a single flat summary.", 1),
            ("Formal cross-level consistency argument (Proposition 6.1).", 1),
            ("Validated on a structurally different large-scale dataset (Beijing air quality).", 1),
            ("Differences across levels are interpretable, not dismissed as inconsistency.", 1),
        ], size=18, space=11)
    content_slide(prs, 7, "TextBox 2", "Contribution II \u2014 Key Findings", c2_find,
                  note="C2 answers RQ2: scalable, coherent, multi-granular explanation, formalised by Proposition 6.1.")

    def c2_lim(sl):
        add_bullets(sl, 0.7, 2.3, 11.9, 4.3, [
            ("Limitations", 0, True),
            ("Clustering remains static, even though the Beijing data are temporal.", 1),
            ("Surrogate-based SHAP plus representative-instance reporting compress observation-level variation.", 1),
            ("Confined to tabular data.", 1),
            ("Still an explanation of a pre-computed partition \u2014 it does not influence learning.", 1, True),
        ], size=18, space=11)
    content_slide(prs, 7, "TextBox 2", "Contribution II \u2014 Limitations", c2_lim,
                  note="C2 is still static and post-hoc; that motivates Contribution III.")

    def c2_take(sl):
        add_bullets(sl, 0.7, 2.3, 11.9, 4.3, [
            ("Takeaways", 0, True),
            ("Shapley attribution stays coherent across granularity \u2014 when the hierarchy is explicit.", 1),
            ("Explanations become interpretable against scale, not just against a single flat partition.", 1),
            ("But the attribution is still post-hoc: it explains a partition that was already computed.", 1, True),
            ("So the next step is to let attribution shape the learning itself.", 0, True),
        ], size=18, space=11)
    content_slide(prs, 7, "TextBox 2", "Contribution II \u2014 Takeaways", c2_take,
                  note="C2 shows Shapley attribution stays coherent across granularity, but remains post-hoc \u2014 motivating C3.")

    # ---------- 9. SECTION: CONTRIBUTION III ----------
    section_slide(prs, "Contribution III \u2014 DyHuCoG: A Dynamic Hypergraph Cooperative Game",
                  note="Contribution III introduces DyHuCoG, a Dynamic Hypergraph Cooperative Game, where Shapley attribution "
                  "becomes an in-training signal. It answers RQ3 and RQ4.")

    def c3_obj(sl):
        add_bullets(sl, 0.7, 2.3, 11.9, 4.3, [
            ("Gap: graph and hypergraph recommenders treat message importance as uniform or attention-weighted, without a principled marginal-contribution account.", 0, True),
            ("Diversity is often a secondary objective or a re-ranking heuristic.", 1),
            ("Interpretability is added after prediction, not integrated into the learning objective.", 1),
            ("Objectives: formulate recommendation as a cooperative game; embed preference-aware Monte Carlo Shapley into message passing; improve ranking, coverage, and diversity jointly.", 0, True),
        ], size=18, space=10)
    content_slide(prs, 7, "TextBox 2", "Contribution III \u2014 Objectives", c3_obj,
                  note="C3 is the flagship: formulate recommendation as a cooperative game, embed Shapley into message passing, and "
                  "improve ranking, coverage and diversity jointly.")

    def c3_rq(sl):
        add_bullets(sl, 0.7, 2.3, 11.9, 4.3, [
            ("RQ3 \u00b7 Can cooperative attribution move beyond post-hoc and enter the learning dynamics of graph recommenders?", 0, True),
            ("RQ4 \u00b7 Can a recommender jointly optimise ranking accuracy, context and diversity when importance is estimated by a cooperative-game utility?", 1),
            ("Objective O1 \u00b7 Formulate recommendation as a cooperative game with users, items and contexts as players.", 1),
            ("Objective O2 \u00b7 Embed preference-aware Monte Carlo Shapley into hypergraph message passing.", 1),
            ("Objective O3 \u00b7 Improve ranking, coverage and diversity jointly.", 1),
        ], size=18, space=9)
    content_slide(prs, 7, "TextBox 2", "Contribution III \u2014 Research Questions & Objectives", c3_rq,
                  note="C3 answers RQ3 and RQ4: cooperative attribution becomes an in-training signal; ranking, coverage and "
                  "diversity improve together.")

    def c3_play(sl):
        add_bullets(sl, 0.7, 2.3, 11.9, 4.3, [
            ("Player set N = U \u222a I \u222a C (users, items, contexts).", 0, True),
            ("Hypergraph H = (V, E, W); V = U \u222a I \u222a C; W = dynamic edge weights from Shapley estimates.", 0),
            ("Coalition S \u2286 N represents entities participating in a recommendation episode.", 0),
            ("Coalition value v(S) measures the quality of the recommendation outcome achievable by S.", 0),
            ("Top-N task: produce a ranked list L_u balancing relevance, diversity, and contextual fit.", 0),
        ], size=18, space=9)
    content_slide(prs, 5, "TextBox 29", "Contribution III \u2014 Players, Hypergraph & Coalition Value", c3_play,
                  note="Users, items and contexts are players; the hypergraph's dynamic edge weights come from Shapley estimates; "
                  "the coalition value measures recommendation quality. The parallelism with clustering is deliberate.")

    def c3_util(sl):
        add_bullets(sl, 0.7, 2.3, 11.9, 4.3, [
            ("v(S) = \u03b1 \u00b7 NDCG@20(S) + \u03b2 \u00b7 Diversity(S) + \u03b3 \u00b7 ContextScore(S), with \u03b1 + \u03b2 + \u03b3 = 1.", 0, True),
            ("The same trade-off the recommender must satisfy is the trade-off from which attribution is computed \u2014 explanatory game and predictive objective aligned by design.", 1),
            ("Preference-weighted: v_pref(S) = v(S) + \u03bb_pref \u00b7 \u03a3_{(u,i)\u2208S} sim(u,i).", 0, True),
            ("\u03b1 = 0.60, \u03b2 = 0.25, \u03b3 = 0.15; \u03bb_pref = 0.20 \u2014 grid-searched, stable (<1.5% variance in NDCG@20).", 1),
            ("Coalition evaluation scoped to the interaction episode (a few dozen players), not the full catalogue.", 1),
        ], size=18, space=8)
    content_slide(prs, 5, "TextBox 29", "Contribution III \u2014 Coalition Utility", c3_util,
                  note="The coalition utility combines ranking, diversity and context; this is the same trade-off the recommender "
                  "optimises and from which attribution is computed, so explanation and objective are aligned by design.")

    def c3_mc(sl):
        add_bullets(sl, 0.7, 2.3, 11.9, 4.3, [
            ("Exact Shapley is combinatorial and infeasible for realistic systems.", 0),
            ("Monte Carlo estimator: \u03c6\u0302_j = (1/M) \u03a3_m [ v(S_m \u222a {j}) \u2212 v(S_m) ].", 0, True),
            ("Preference-aware: \u03c6\u0302_j^pref = (1/M) \u03a3_m [ v_pref(S_m \u222a {j}) \u2212 v_pref(S_m) ].", 0, True),
            ("Unbiased; variance = \u03c3\u00b2/M \u2192 MSE decays O(1/M), absolute error O(1/\u221aM).", 1),
            ("M = 50 selected: MSE \u2248 1.4\u00d710\u207b\u2075, ~99% accuracy on MovieLens-1M.", 1, True),
            ("Refreshed every 10 batches (~49 updates/epoch), smoothed by exponential moving average.", 1),
        ], size=18, space=8)
    content_slide(prs, 5, "TextBox 29", "Contribution III \u2014 Monte Carlo Shapley", c3_mc,
                  note="Exact Shapley is infeasible; the Monte Carlo estimator is unbiased with variance decaying as 1/M. M=50 gives "
                  "~99% accuracy, refreshed every 10 batches.")

    def c3_arch(sl):
        add_bullets(sl, 0.7, 2.3, 6.3, 4.3, [
            ("Base propagation: e^(l+1) = \u03c3( D^-1/2 A D^-1/2 e^(l) ).", 0),
            ("Shapley-weighted: e_j^(l+1) = \u03c3( W^(l) e_j^(l) + \u03a3_{k\u2208N(j)} w_jk e_k^(l) ).", 0, True),
            ("Normalised weights: w_jk = \u03c6\u0302_jk / \u03a3_{k'\u2208N(j)} \u03c6\u0302_jk'.", 0, True),
            ("Clipped + exponentially smoothed before normalisation (stabilises sparse regimes).", 1),
            ("Attention gate: a_ui = \u03c3( W_a[ e_u, e_i, l_i ] ); y_ui = (1 + a_ui) \u27e8e_u, e_i\u27e9.", 1),
            ("Context-aware score: f(u,i,c) = y_ui + \u03bb_c \u27e8g(c_ui), e_cui\u27e9.", 1),
        ], size=17)
        add_figure(sl, os.path.join(HERE, "_figs", "c3_architecture_paper.png"), 7.1, 2.4, w=5.5)
    content_slide(prs, 5, "TextBox 29", "Contribution III \u2014 DyHuCoG Architecture", c3_arch,
                  note="The decisive move: Shapley-weighted message passing, normalised and smoothed, plus an attention gate and a "
                  "context-aware score. Attribution directly governs how information propagates.")

    def c3_loss(sl):
        add_bullets(sl, 0.7, 2.3, 11.9, 4.3, [
            ("L = L_rec + \u03bb_div L_div + \u03bb_ctx L_ctx + \u03bb_reg L_reg.", 0, True),
            ("L_rec \u2013 Bayesian Personalised Ranking (pairwise, implicit feedback).", 1),
            ("L_div \u2013 Intra-List Diversity regulariser: penalises redundant ranked lists.", 1),
            ("L_ctx \u2013 Context alignment: match context embedding to context-node representation.", 1),
            ("L_reg \u2013 L2 weight decay.", 1),
            ("The learning objective and coalition value are aligned: DyHuCoG trains to optimise the same balance that later determines attribution.", 1, True),
        ], size=18, space=9)
    content_slide(prs, 5, "TextBox 29", "Contribution III \u2014 Composite Loss", c3_loss,
                  note="The loss combines BPR ranking, intra-list diversity and context alignment. The objective and the coalition "
                  "value are aligned, so the explanation is a direct read-out of what the model optimises.")

    def c3_tab1(sl):
        add_table(sl, 0.7, 2.3, 11.9, 3.9,
             [["Dataset", "Model", "NDCG@20", "Recall@20", "Coverage", "Diversity"],
              ["MovieLens-1M", "HPCF", "0.2528", "0.2098", "0.342", "0.461"],
              ["MovieLens-1M", "DyHuCoG", "0.2775", "0.2362", "0.397", "0.516"],
              ["Amazon-Book", "HPCF", "0.0270", "0.0359", "0.259", "0.535"],
              ["Amazon-Book", "DyHuCoG", "0.0306", "0.0417", "0.336", "0.602"]],
             col_ratios=[3.2, 2, 2.2, 2.2, 2.2, 2.2], header_bg=ACCENT, hl_rows=[1, 3])
    content_slide(prs, 12, "TextBox 6", "Contribution III \u2014 Headline Results", c3_tab1,
                  note="DyHuCoG improves NDCG, recall, coverage and diversity together vs the strongest baseline HPCF. The sparser "
                  "the data, the larger the gain \u2014 Shapley-driven weighting helps most when signal is weak.")

    def c3_tab2(sl):
        add_table(sl, 0.7, 2.3, 11.8, 3.6,
             [["Metric", "MovieLens-1M", "Amazon-Book"],
              ["NDCG@20", "+9.77%", "+13.33%"],
              ["Recall@20", "+12.58%", "+16.16%"],
              ["Coverage", "+16.1%", "+29.7%"],
              ["Intra-List Diversity", "+11.9%", "+12.5%"]],
             col_ratios=[4, 3.2, 3.2], header_bg=ACCENT)
    content_slide(prs, 12, "TextBox 6", "Contribution III \u2014 Relative Gains", c3_tab2,
                  note="DyHuCoG improves ranking accuracy, coverage and diversity simultaneously \u2014 it does not sacrifice one for "
                  "the others. The largest gains appear on the sparser Amazon-Book.")

    def c3_div(sl):
        add_bullets(sl, 0.7, 2.3, 11.9, 4.3, [
            ("MovieLens-1M: Coverage 0.342 \u2192 0.397 (+16.1%); ILD 0.461 \u2192 0.516 (+11.9%).", 0, True),
            ("Amazon-Book: Coverage 0.259 \u2192 0.336 (+29.7%); ILD 0.535 \u2192 0.602 (+12.5%).", 0, True),
            ("Reduced filter-bubble effect and greater discovery opportunity \u2014 while NDCG/Recall also improve.", 1),
        ], size=18, space=12)
    content_slide(prs, 5, "TextBox 29", "Contribution III \u2014 Diversity Gains", c3_div,
                  note="Both diversity metrics improve while NDCG/recall also improve \u2014 so we are not trading accuracy for "
                  "diversity. The model surfaces more of the catalogue and recommends less redundant lists.")

    def c3_ab(sl):
        add_table(sl, 0.7, 2.3, 11.9, 3.9,
             [["Variant", "ML-1M NDCG@20", "% Drop", "Amazon NDCG@20", "% Drop"],
              ["Full DyHuCoG", "0.2775", "\u2013", "0.0306", "\u2013"],
              ["w/o Shapley Value", "0.2647", "4.6%", "0.0287", "6.1%"],
              ["w/o Hypergraph", "0.2586", "6.8%", "0.0279", "8.9%"],
              ["w/o Attention", "0.2678", "3.5%", "0.0295", "3.5%"],
              ["w/o Context", "0.2547", "8.2%", "0.0272", "11.0%"],
              ["w/o Diversity", "0.2614", "5.8%", "0.0288", "5.8%"]],
             col_ratios=[3.6, 2.4, 1.8, 2.4, 1.8], header_bg=ACCENT, hl_rows=[1, 5])
    content_slide(prs, 12, "TextBox 6", "Contribution III \u2014 Ablation", c3_ab,
                  note="Every component contributes. Removing Shapley drops NDCG by 4.6% (ML-1M) / 6.1% (Amazon); context removal "
                  "causes the largest single loss. Shapley weighting is load-bearing, not decorative.")

    def c3_cost(sl):
        add_bullets(sl, 0.7, 2.3, 11.9, 4.3, [
            ("Training: DyHuCoG ~2000 s vs HPCF ~1125 s on MovieLens-1M (\u2248 1.78\u00d7).", 0),
            ("Inference: 1.84 ms/query (ML-1M), 8.52 ms (Amazon) \u2014 suitable for real-time deployment.", 0),
            ("Memory: 4.4 vs 4.1 GB (ML-1M); 17.9 vs 16.8 GB (Amazon).", 0),
            ("Per-epoch cost: O((L+1)md) + O((M/f)m).", 0),
            ("Shapley convergence: M=50 \u2192 MSE 1.4\u00d710\u207b\u2075, ~99% accuracy; M=100 \u2192 MSE 3.5\u00d710\u207b\u2076 (diminishing returns).", 0, True),
        ], size=18, space=9)
    content_slide(prs, 5, "TextBox 29", "Contribution III \u2014 Computational Cost", c3_cost,
                  note="The attribution cost is proportionate: ~1.78\u00d7 training time, milliseconds per query at inference. M=50 is the "
                  "right operating point.")

    def c3_sig(sl):
        add_bullets(sl, 0.7, 2.3, 11.9, 4.3, [
            ("Paired t-tests on per-user NDCG@20 (n = 6,040 users; df = 6,039).", 0),
            ("DyHuCoG outperforms every baseline with extremely small p-values after Holm\u2013Bonferroni correction.", 0),
            ("vs HPCF: t = 46.38, Cohen\u2019s d_z = 1.3345, p = 1.81\u00d710\u207b\u00b2\u2077\u2070.", 1, True),
            ("Wilcoxon signed-rank test also significant (p < 0.001).", 1),
            ("Effect sizes are large \u2014 improvements are substantively meaningful, not merely statistically visible.", 1),
        ], size=18, space=9)
    content_slide(prs, 5, "TextBox 29", "Contribution III \u2014 Statistical Significance", c3_sig,
                  note="Improvements are statistically significant (p ~10^-270) with large effect sizes, so they are not a "
                  "large-sample artefact. Tabulated paired tests apply to MovieLens-1M; Amazon remains descriptive.")

    def c3_cold(sl):
        add_bullets(sl, 0.7, 2.3, 11.9, 4.3, [
            ("Cold-start (5 or fewer interactions): NDCG@20 \u2248 0.061 (user) and 0.057 (item), improving over HPCF by ~10%.", 0),
            ("Cross-dataset: MovieLens +9.9%, Amazon +14.8%, Yelp2018 +11.8%.", 0),
            ("Interpretability: a SHAP waterfall decomposes a recommendation into ranking, diversity, context and preference contributions.", 0, True),
            ("Popularity bias: Shapley measures marginal utility, not raw frequency \u2014 weak but informative interactions retain influence.", 1),
        ], size=18, space=10)
    content_slide(prs, 5, "TextBox 29", "Contribution III \u2014 Cold-Start & Interpretability", c3_cold,
                  note="DyHuCoG also improves cold-start and produces a SHAP waterfall that decomposes a recommendation into the "
                  "same components used during training \u2014 structurally faithful, not an external approximation.")

    def c3_ans(sl):
        add_bullets(sl, 0.7, 2.3, 11.9, 4.3, [
            ("Answers to the questions and objectives", 0, True),
            ("RQ3 answered in the affirmative \u2014 attribution becomes an in-training signal, not a post-hoc diagnostic.", 1),
            ("RQ4 answered in the affirmative \u2014 ranking, coverage and diversity improve together.", 1),
            ("O1 met \u00b7 recommendation formulated as a cooperative game over users, items and contexts.", 1),
            ("O2 met \u00b7 preference-aware Monte Carlo Shapley embedded in message passing.", 1),
            ("O3 met \u00b7 +9.9% (MovieLens) / +14.8% (Amazon) NDCG, with higher coverage and diversity.", 1),
        ], size=18, space=8)
    content_slide(prs, 7, "TextBox 2", "Contribution III \u2014 Answers", c3_ans,
                  note="RQ3 and RQ4 answered yes: attribution becomes an in-training signal; accuracy, coverage and diversity improve "
                  "together. The strongest claim is that the explanation is structurally faithful.")

    def c3_find(sl):
        add_bullets(sl, 0.7, 2.3, 11.9, 4.3, [
            ("Key findings", 0, True),
            ("Cooperative attribution used as an in-training signal \u2014 a stronger claim than the clustering chapters.", 1),
            ("The accuracy\u2013diversity\u2013context trade-off is not structurally fixed.", 1),
            ("NDCG, recall, coverage and diversity improve together on both datasets.", 1),
            ("Largest gains on the sparsest data (Amazon-Book), consistent with Shapley weighting helping when signal is weak.", 1),
        ], size=18, space=10)
    content_slide(prs, 7, "TextBox 2", "Contribution III \u2014 Key Findings", c3_find,
                  note="C3 answers RQ3 and RQ4, making the strongest claim of the thesis: attribution is an in-training signal and "
                  "the accuracy\u2013diversity trade-off is negotiable.")

    def c3_lim(sl):
        add_bullets(sl, 0.7, 2.3, 11.9, 4.3, [
            ("Limitations", 0, True),
            ("Measurable computational overhead \u2014 roughly 1.78\u00d7 training time over HPCF.", 1),
            ("Depends on availability of meaningful context.", 1),
            ("Monte Carlo Shapley could be improved by variance reduction.", 1),
            ("Ablation is component-wise, so it does not test factorial interactions.", 1),
            ("Baselines finalised through early 2026, so superiority is claimed only against the tested baselines.", 1),
        ], size=18, space=9)
    content_slide(prs, 7, "TextBox 2", "Contribution III \u2014 Limitations", c3_lim,
                  note="Limitations are honest: computational overhead, dependence on context, and component-wise ablation.")

    def c3_take(sl):
        add_bullets(sl, 0.7, 2.3, 11.9, 4.3, [
            ("Takeaways", 0, True),
            ("Attribution is a first-class part of the learning objective, not a post-hoc diagnostic.", 1),
            ("The explanation is a direct read-out of the objective the model already optimises.", 1),
            ("This makes the explanation structurally faithful rather than an external approximation.", 1, True),
            ("Built on a principled, axiomatic basis \u2014 aligned with trustworthy-AI expectations.", 1),
        ], size=18, space=10)
    content_slide(prs, 7, "TextBox 2", "Contribution III \u2014 Takeaways", c3_take,
                  note="The takeaway is the conceptual shift: attribution is part of the modelling logic itself, which makes it "
                  "structurally faithful and aligned with trustworthy-AI regulation.")

    # ---------- 10. SECTION: CONCLUSION & PERSPECTIVES ----------
    section_slide(prs, "Conclusion & Perspectives",
                  note="Let me bring everything together: a synthesis of the three contributions, the published papers, the honest "
                  "limitations, the future directions, and a clear statement of the thesis answer.")

    def synth(sl):
        add_table(sl, 0.7, 2.3, 11.9, 3.9,
             [["Contribution", "Main idea", "Achievement", "Key finding"],
              ["C1", "Explain black-box clustering via Shapley", "PCA\u2013KMeans\u2013LightGBM\u2013TreeSHAP pipeline", "Faithful, chemistry-consistent cluster attribution (wine)"],
              ["C2", "Multi-level, large-scale clustering XAI", "Cross-level SHAP aggregation + Prop. 6.1", "Coherent under hierarchy; interprets differences"],
              ["C3", "DyHuCoG hypergraph cooperative game", "Preference-aware Shapley as in-training signal", "Accuracy + coverage + diversity improve together"]],
             col_ratios=[1.6, 4.2, 4.2, 5.0], header_bg=ACCENT)
    content_slide(prs, 12, "TextBox 6", "Conclusion \u2014 Synthesis", synth,
                  note="Three contributions, one thread: C1 makes hidden structure intelligible; C2 keeps explanation coherent under "
                  "scale and hierarchy; C3 carries the same logic inside the learning dynamics. A shared attribution perspective.")

    def papers(sl):
        add_table(sl, 0.7, 2.3, 11.9, 3.9,
             [["No.", "Title", "Venue", "Status"],
              ["I", "Shapley Values for Explaining the Black Box Nature of ML Model Clustering", "Procedia Computer Science 220, 806\u2013811", "Published, 2023"],
              ["II", "Game Theory Meets Explainable AI: An Enhanced Approach to Understanding Black Box Models Through Shapley Values", "IJACSA 16(7), 716\u2013725", "Published, 2025"],
              ["III", "DyHuCoG: A Dynamic Hypergraph Cooperative Game for Preference-aware Recommendation", "IJIES 19(2), 887\u2013902", "Published, 2026"]],
             col_ratios=[1, 6.5, 3.2, 2], header_bg=ACCENT)
    content_slide(prs, 12, "TextBox 6", "Conclusion \u2014 Publications", papers,
                  note="The thesis synthesises three peer-reviewed publications, mapped to the chapters. What it adds is the "
                  "multi-level formalisation, cross-chapter comparisons, and the explicit mapping to the five research questions.")

    def concl_lim(sl):
        add_bullets(sl, 0.7, 2.3, 11.9, 4.3, [
            ("Computational \u2013 exact Shapley is intractable; every contribution relies on approximation, surrogates, or restricted reporting.", 0),
            ("Methodological \u2013 clustering depends on surrogate fidelity; recommendation depends on stable approximate contributions and adequate context.", 0),
            ("Empirical \u2013 tabular clustering + benchmark recommendation; no multimodal, sequential, or online deployment; no dedicated human-subject actionability study.", 0),
            ("Claim scope \u2013 a coherent and productive perspective, not one fully unified framework eliminating all tension.", 0, True),
        ], size=18, space=10)
    content_slide(prs, 4, "TextBox 25", "Conclusion \u2014 Limitations", concl_lim,
                  note="Honest limitations: computational approximation, surrogate/context dependence, tabular+benchmark scope, and a "
                  "coherent perspective rather than one framework eliminating all tension.")

    def persp(sl):
        add_bullets(sl, 0.7, 2.3, 11.9, 4.3, [
            ("Scalable cooperative attribution \u2013 lower-variance Shapley, learned proposal distributions, adaptive refresh policies.", 0),
            ("Online / streaming recommendation \u2013 truly incremental settings with evolving graphs and delayed feedback.", 0),
            ("Richer human-centred evaluation \u2013 do explanations measurably improve analyst judgement, user trust, intervention quality, or perceived fairness?", 0),
            ("Broader trustworthy-AI evaluation \u2013 exposure fairness, transparency requirements, governance-oriented auditing.", 0),
        ], size=18, space=11)
    content_slide(prs, 4, "TextBox 25", "Perspectives & Future Work", persp,
                  note="Future work turns limitations into an agenda: scalable attribution, online recommendation, human-centred "
                  "evaluation, and broader trustworthy-AI/fairness evaluation.")

    def answer(sl):
        add_bullets(sl, 0.7, 2.3, 11.9, 4.3, [
            ("Thesis answer: cooperative game theory can function as a shared methodological perspective for actionable explanation across clustering and recommendation.", 0, True),
            ("Shapley attribution as a common formal language for feature, interaction, and context importance allocation.", 1),
            ("Faithful clustering explanation, hierarchical explanatory coherence, and contribution-aware recommendation learning.", 1),
            ("Explanation as method, not commentary \u2014 from post-hoc description to in-training guidance.", 1),
            ("Aligned with trustworthy-AI requirements (EU AI Act, OECD principles, GDPR).", 1),
        ], size=18, space=9)
    content_slide(prs, 4, "TextBox 25", "Conclusion \u2014 Thesis Answer", answer,
                  note="The thesis answer: cooperative game theory is a shared methodological perspective for actionable explanation; "
                  "explanation moves from commentary to method, aligned with trustworthy-AI regulation.")

    def refs(sl):
        add_bullets(sl, 0.7, 2.3, 11.9, 4.3, [
            ("[R1] Louhichi, M. & Lazaar, M. Shapley Values for Explaining the Black Box Nature of ML Model Clustering. Procedia Computer Science 220, 806\u2013811 (2023).", 0),
            ("[R2] Louhichi, M. & Lazaar, M. Game Theory Meets Explainable AI. IJACSA 16(7), 716\u2013725 (2025).", 0),
            ("[R3] Louhichi, M. & Lazaar, M. DyHuCoG. IJIES 19(2), 887\u2013902 (2026).", 0),
            ("[R4] Lundberg, S.M. & Lee, S.-I. A Unified Approach to Interpreting Model Predictions (SHAP). NeurIPS 30 (2017).", 0),
            ("[R5] Shapley, L.S. A Value for n-Person Games. Contributions to the Theory of Games II (1953).", 0),
            ("[R6] Wang, X. et al. Hypergraph Learning: Methods and Practices. IEEE TPAMI 44(5) (2022).", 0),
            ("[R7] European Commission. Proposal for a Regulation on Artificial Intelligence (EU AI Act). COM(2021) 206 final (2021).", 0),
        ], size=15, space=8)
    content_slide(prs, 12, "TextBox 6", "References", refs,
                  note="References underpinning the thesis: my three peer-reviewed papers, the SHAP and Shapley foundations, the "
                  "hypergraph-learning literature, and the regulation that motivates actionable explanation.")

    # ---------- 11. THANK YOU ----------
    s = clone_slide(prs, prs.slides[14])
    ths = _find_title(s.shapes)
    _set_title(ths, "Thank You")
    uni = _set_text(s, "Borcelle University", "ENSIAS, Mohammed V University, Rabat")
    name = _set_text(s, "Presentation By Donna Stroupe", "Mouad LOUHICHI")
    email = _set_text(s, "hello@reallygreatsite.com", "mouad.louhichi@um5.ac.ma")
    keep = [x for x in (ths, uni, name, email) if x is not None]
    _strip_except(s, keep)
    set_notes(s, "Thank you very much for your attention. I am now happy to take your questions and comments.")

    # drop the original 15 template slides; keep only our cloned content
    for _ in range(15):
        el = prs.slides._sldIdLst[0]
        rid = el.get(qn('r:id'))
        if rid is not None:
            prs.part.drop_rel(rid)
        prs.slides._sldIdLst.remove(el)

    prs.save(OUT)
    print(f"Saved {len(prs.slides)} slides -> {OUT}")

if __name__ == "__main__":
    main()
