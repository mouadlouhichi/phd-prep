#!/usr/bin/env python3
"""
Build Mouad Louhichi's PhD viva presentation, mirroring the structure of
Redwane Nesmaoui's defense deck (example-phd-passes/Presentation1 (1) (1).pptx)
but filled with the content of Mouad's thesis.

Output: viva/MOUAD_LOUHICHI_VIVA.pptx  (16:9, with speaker notes on every slide)

Run:
    ./.venv/bin/python3 viva/build_presentation.py
"""
import os, shutil, glob
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ----------------------------------------------------------------------------
# Style constants
# ----------------------------------------------------------------------------
OUT = os.path.join(os.path.dirname(os.path.abspath(__file__)), "MOUAD_LOUHICHI_VIVA.pptx")

SW, SH = Inches(13.333), Inches(7.5)

# Colours (academic, clean)
NAVY   = RGBColor(0x0B, 0x2C, 0x4A)   # deep navy
BLUE   = RGBColor(0x1F, 0x6F, 0xB2)   # accent blue
TEAL   = RGBColor(0x0E, 0x7C, 0x7B)   # teal accent
GREEN  = RGBColor(0x2E, 0x7D, 0x32)   # positive
GREY   = RGBColor(0x44, 0x44, 0x44)
LGREY  = RGBColor(0x9A, 0x9A, 0x9A)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
BG     = RGBColor(0xF7, 0xF9, 0xFB)
BORDER = RGBColor(0xD8, 0xDE, 0xE6)

prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH
BLANK = prs.slide_layouts[6]

FONT = "Calibri"

# ----------------------------------------------------------------------------
# Low-level helpers
# ----------------------------------------------------------------------------
def _bg(slide, color=WHITE):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = color

def _rect(slide, x, y, w, h, color):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp

def _tb(slide, x, y, w, h, text, size, color, bold=False, align=PP_ALIGN.LEFT,
        font=FONT, wrap=True, italic=False, anchor=MSO_ANCHOR.TOP, line_spacing=1.0):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = wrap
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return box

def _bullets(slide, x, y, w, h, items, size=18, color=GREY, bullet_color=BLUE,
             space_after=10, line_spacing=1.05, anchor=MSO_ANCHOR.TOP):
    """items: list of strings OR (text, level) OR (text, level, bold)."""
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    first = True
    for it in items:
        if isinstance(it, str):
            text, level, bold = it, 0, False
        elif len(it) == 2:
            text, level = it
            bold = False
        else:
            text, level, bold = it
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = level
        p.space_after = Pt(space_after if level == 0 else space_after - 3)
        p.line_spacing = line_spacing
        marker = ("\u25aa  " if level == 0 else "\u2013  ")
        r0 = p.add_run(); r0.text = marker
        r0.font.name = FONT; r0.font.size = Pt(size); r0.font.bold = True
        r0.font.color.rgb = bullet_color
        r = p.add_run(); r.text = text
        r.font.name = FONT; r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
    return box

def _notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text

def _footer(slide, num, section=""):
    _tb(slide, Inches(0.55), Inches(7.06), Inches(8), Inches(0.35),
        "PhD Viva \u2013 Mouad LOUHICHI \u2013 Cooperative Game Theory & Shapley for XAI in Recommendation",
        size=11, color=LGREY)
    _tb(slide, Inches(12.2), Inches(7.06), Inches(0.6), Inches(0.35),
        str(num), size=11, color=LGREY, align=PP_ALIGN.RIGHT)

def _header(slide, title, subtitle=None):
    _rect(slide, 0, 0, SW, Inches(1.05), NAVY)
    _tb(slide, Inches(0.55), Inches(0.22), Inches(12), Inches(0.62),
        title, size=28, color=WHITE, bold=True)
    if subtitle:
        _tb(slide, Inches(0.55), Inches(0.80), Inches(12), Inches(0.30),
            subtitle, size=13, color=RGBColor(0xC9, 0xD9, 0xEA))
    _topnav(slide, title)

# ----------------------------------------------------------------------------
# Slide constructors
# ----------------------------------------------------------------------------
def title_slide(title_lines, sub_lines, notes):
    s = prs.slides.add_slide(BLANK)
    _bg(s, NAVY)
    _rect(s, 0, Inches(2.55), SW, Inches(0.07), TEAL)
    # banner
    _tb(s, Inches(0.9), Inches(0.7), Inches(11.5), Inches(0.5),
        "Doctoral Studies Center in Information and Engineering Sciences and Technologies (ST2I)",
        size=15, color=RGBColor(0xBF, 0xD4, 0xE8), italic=True)
    y = Inches(1.6)
    for i, ln in enumerate(title_lines):
        _tb(s, Inches(0.9), y, Inches(11.5), Inches(0.7), ln,
            size=34 if i == 0 else 24,
            color=WHITE, bold=True)
        y += Inches(0.75)
    y += Inches(0.1)
    for ln in sub_lines:
        _tb(s, Inches(0.9), y, Inches(11.5), Inches(0.45), ln, size=15,
            color=RGBColor(0xCF, 0xE0, 0xEF))
        y += Inches(0.5)
    _notes(s, notes)
    return s

def section_slide(big, small, notes=None):
    s = prs.slides.add_slide(BLANK)
    _bg(s, BG)
    _rect(s, 0, 0, Inches(0.5), SH, BLUE)
    _tb(s, Inches(1.6), Inches(3.0), Inches(10.5), Inches(1.4), big,
        size=48, color=NAVY, bold=True)
    if small:
        _tb(s, Inches(1.6), Inches(4.5), Inches(10.5), Inches(0.6), small,
            size=20, color=BLUE, italic=True)
    if notes:
        _notes(s, notes)
    return s

def content_slide(title, bullets, notes, subtitle=None, num=None, two_col=None,
                  layout="bullets", image=None):
    s = prs.slides.add_slide(BLANK)
    _bg(s, WHITE)
    _header(s, title, subtitle)
    if layout == "bullets":
        _bullets(s, Inches(0.7), Inches(1.95), Inches(11.9), Inches(4.9), bullets, size=19)
    elif layout == "two_col":
        left, right = bullets
        _bullets(s, Inches(0.7), Inches(1.95), Inches(6.0), Inches(4.9), left, size=17)
        _bullets(s, Inches(6.9), Inches(1.95), Inches(5.9), Inches(4.9), right, size=17)
    elif layout == "image_right":
        _bullets(s, Inches(0.7), Inches(1.95), Inches(6.0), Inches(4.9), bullets, size=18)
        if image:
            _pic_fit(s, image, Inches(7.0), Inches(1.95), Inches(5.7), Inches(4.9), align="center")
    elif layout == "image_left":
        if image:
            _pic_fit(s, image, Inches(0.7), Inches(1.95), Inches(5.7), Inches(4.9), align="center")
        _bullets(s, Inches(7.0), Inches(1.95), Inches(6.0), Inches(4.9), bullets, size=18)
    if num:
        _footer(s, num)
    _notes(s, notes)
    return s

def table_slide(title, headers, rows, notes, subtitle=None, num=None, col_widths=None,
                font_size=13, note_col=None, highlight_rows=None):
    s = prs.slides.add_slide(BLANK)
    _bg(s, WHITE)
    _header(s, title, subtitle)
    ncols = len(headers)
    nrows = len(rows) + 1
    left = Inches(0.7)
    top = Inches(1.95)
    width = Inches(11.9)
    height = Inches(0.42)
    shape = s.shapes.add_table(nrows, ncols, left, top, width, height * nrows)
    table = shape.table
    # column widths
    if col_widths:
        total = sum(col_widths)
        for j, w in enumerate(col_widths):
            table.columns[j].width = Emu(int(width * (w / total)))
    # header
    for j, h in enumerate(headers):
        c = table.cell(0, j)
        c.text = h
        p = c.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.runs[0]
        r.font.size = Pt(font_size); r.font.bold = True; r.font.color.rgb = WHITE
        fill = c.fill; fill.solid(); fill.fore_color.rgb = NAVY
        c.vertical_anchor = MSO_ANCHOR.MIDDLE
    # body
    for i, row in enumerate(rows, start=1):
        is_hl = highlight_rows is not None and (i - 1) in highlight_rows
        for j, val in enumerate(row):
            c = table.cell(i, j)
            c.text = str(val)
            p = c.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT
            r = p.runs[0]
            r.font.size = Pt(font_size)
            r.font.color.rgb = WHITE if is_hl else GREY
            if is_hl or (i % 2 == 0):
                r.font.bold = is_hl
                c.fill.solid()
                c.fill.fore_color.rgb = BLUE if is_hl else RGBColor(0xEF, 0xF3, 0xF8)
            else:
                c.fill.solid()
                c.fill.fore_color.rgb = WHITE
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            c.margin_left = Inches(0.06); c.margin_right = Inches(0.06)
            c.margin_top = Inches(0.02); c.margin_bottom = Inches(0.02)
    if num:
        _footer(s, num)
    _notes(s, notes)
    return s

def result_card_slide(title, cards, notes, subtitle=None, num=None):
    """cards: list of (big_number, label, sublabel) — placed in a row."""
    s = prs.slides.add_slide(BLANK)
    _bg(s, WHITE)
    _header(s, title, subtitle)
    n = len(cards)
    gap = Inches(0.35)
    total_w = Inches(11.9)
    cw = Emu(int((total_w - gap * (n - 1)) / n))
    x = Inches(0.7)
    top = Inches(2.2)
    ch = Inches(3.2)
    for big, label, sub in cards:
        card = _rect(s, x, top, cw, ch, BG)
        # accent
        _rect(s, x, top, cw, Inches(0.12), BLUE)
        _tb(s, x + Inches(0.15), top + Inches(0.45), cw - Inches(0.3), Inches(0.9),
            big, size=30, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
        _tb(s, x + Inches(0.15), top + Inches(1.5), cw - Inches(0.3), Inches(0.6),
            label, size=16, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
        _tb(s, x + Inches(0.15), top + Inches(2.15), cw - Inches(0.3), Inches(1.0),
            sub, size=12, color=GREY, align=PP_ALIGN.CENTER)
        x += cw + gap
    if num:
        _footer(s, num)
    _notes(s, notes)
    return s

# ----------------------------------------------------------------------------
# Assets: copy extracted paper figures + generate clean diagrams (PIL)
# ----------------------------------------------------------------------------
ASSETS = os.path.join(os.path.dirname(os.path.abspath(__file__)), "assets")
ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
os.makedirs(ASSETS, exist_ok=True)

FONT_PATH = "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf"
FONT_BOLD = "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf"

# PIL colour tuples
_P_NAVY = (11, 44, 74); _P_BLUE = (31, 111, 178); _P_TEAL = (14, 124, 123)
_P_GREY = (68, 68, 68); _P_LGREY = (154, 154, 154); _P_WHITE = (255, 255, 255)
_P_LIGHT = (238, 242, 247); _P_ACCENT = (13, 80, 160); _P_RED = (196, 40, 40)

def _pf(sz, bold=False):
    return ImageFont.truetype(FONT_BOLD if bold else FONT_PATH, sz)

def _wrap(draw, text, font, maxw):
    out = []
    for para in text.split("\n"):
        words = para.split(" ")
        line = ""
        for w in words:
            t = (line + " " + w).strip()
            if draw.textlength(t, font=font) <= maxw or not line:
                line = t
            else:
                out.append(line); line = w
        out.append(line)
    return out

def _pbox(d, xy, fill, outline=None, text="", tcolor=_P_WHITE, tsize=20, bold=False,
          radius=14, sub=None, scolor=None, ssize=15):
    d.rounded_rectangle(xy, radius=radius, fill=fill,
                        outline=outline if outline else fill, width=3 if outline else 0)
    fnt = _pf(tsize, bold)
    if text:
        lines = text.split("\n")
        ascent, descent = fnt.getmetrics()
        lh = ascent + descent
        total = lh * len(lines)
        cx = (xy[0] + xy[2]) / 2
        y = (xy[1] + xy[3]) / 2 - total / 2 - (lh * 0.35 if sub else 0)
        for ln in lines:
            bb = d.textbbox((0, 0), ln, font=fnt)
            w = bb[2] - bb[0]
            d.text((cx - w / 2, y), ln, font=fnt, fill=tcolor)
            y += lh
    if sub:
        sf = _pf(ssize, False)
        for i, ln in enumerate(_wrap(d, sub, sf, (xy[2] - xy[0]) - 24)):
            bb = d.textbbox((0, 0), ln, font=sf)
            w = bb[2] - bb[0]
            cx = (xy[0] + xy[2]) / 2
            yy = xy[3] - 14 - (ssize + 6) * (len(_wrap(d, sub, sf, (xy[2]-xy[0])-24)) - i)
            d.text((cx - w / 2, yy), ln, font=sf, fill=scolor or _P_GREY)

def _parrow(d, p1, p2, color=_P_GREY, width=4, head=10):
    d.line([p1, p2], fill=color, width=width)
    import math
    ang = math.atan2(p2[1] - p1[1], p2[0] - p1[0])
    for da in (math.pi * 5 / 6, -math.pi * 5 / 6):
        a = ang + da
        d.polygon([p2,
                   (p2[0] - head * math.cos(a), p2[1] - head * math.sin(a)),
                   (p2[0] - head * 0.6 * math.cos(a - 0.4), p2[1] - head * 0.6 * math.sin(a - 0.4))],
                  fill=color)

def _extract(pdf_name, xref, dst):
    """Extract an embedded raster figure (by xref) directly from a source PDF."""
    try:
        import pymupdf
    except Exception:
        return False
    path = os.path.join(ROOT, pdf_name)
    if not os.path.exists(path):
        return False
    try:
        doc = pymupdf.open(path)
        d = doc.extract_image(xref)
        doc.close()
    except Exception:
        return False
    with open(os.path.join(ASSETS, dst), "wb") as f:
        f.write(d["image"])
    return True

def _d_evolution(path):
    W, H = 1240, 360
    img = Image.new("RGB", (W, H), _P_WHITE); d = ImageDraw.Draw(img)
    boxes = ["Similarity\nfilters", "Matrix\nFactorisation", "Neural CF",
             "Graph CNN\n(LightGCN)", "Hypergraph\nrecommender"]
    bw, bh, gap, x0, y = 218, 86, 26, 24, 60
    for i, b in enumerate(boxes):
        x = x0 + i * (bw + gap)
        fill = _P_NAVY if i == 0 else (_P_TEAL if i == len(boxes) - 1 else _P_ACCENT)
        _pbox(d, (x, y, x + bw, y + bh), fill, text=b, tsize=19, bold=True)
        if i < len(boxes) - 1:
            _parrow(d, (x + bw + 3, y + bh / 2), (x + bw + gap - 3, y + bh / 2),
                    color=_P_BLUE, width=4)
    # interpretability deficit arrow (red, downward, on the right)
    rx = x0 + len(boxes) * (bw + gap) - bw + bw / 2
    _parrow(d, (rx, y + bh + 24), (rx, H - 56), color=_P_RED, width=5, head=12)
    d.text((rx - 70, H - 50), "Interpretability deficit \u2193", fill=_P_RED,
           font=_pf(18, True))
    img.save(path)

def _d_motivation(path):
    W, H = 1240, 400
    img = Image.new("RGB", (W, H), _P_WHITE); d = ImageDraw.Draw(img)
    cards = [("Ubiquity", "Opaque AI mediates what billions see, buy and watch every day."),
             ("Black box", "Top recommenders & clustering pipelines stay opaque to users and designers."),
             ("Toward trust", "Transparency & accountability built into the model, not bolted on.")]
    cw, ch, gap, x0, y = 386, 250, 36, 18, 30
    for i, (t, sub) in enumerate(cards):
        x = x0 + i * (cw + gap)
        _pbox(d, (x, y, x + cw, y + ch), _P_LIGHT, outline=_P_BLUE, text=t,
              tcolor=_P_NAVY, tsize=26, bold=True, sub=sub, scolor=_P_GREY, ssize=16)
    img.save(path)

def _d_actionable(path):
    W, H = 1240, 330
    img = Image.new("RGB", (W, H), _P_WHITE); d = ImageDraw.Draw(img)
    boxes = [("Modifiable factor", "expressed in the\ndomain vocabulary"),
             ("\u0394 model output", "specifiable change\nin prediction"),
             ("Actionable insight", "supports intervention,\nnot just description")]
    bw, bh, gap, x0, y = 300, 110, 70, 30, 60
    for i, (t, sub) in enumerate(boxes):
        x = x0 + i * (bw + gap)
        _pbox(d, (x, y, x + bw, y + bh), _P_BLUE, text=t, tsize=21, bold=True,
              sub=sub, scolor=_P_WHITE, ssize=15)
        if i < len(boxes) - 1:
            _parrow(d, (x + bw + 6, y + bh / 2), (x + bw + gap - 6, y + bh / 2),
                    color=_P_NAVY, width=4)
    # return loop: analyst intervenes
    _parrow(d, (x0 + 2 * (bw + gap) + bw / 2, y + bh + 26),
            (x0 + bw / 2, y + bh + 26), color=_P_TEAL, width=4)
    d.text((W / 2 - 90, y + bh + 34), "analyst intervenes \u2192", fill=_P_TEAL,
           font=_pf(17, True))
    img.save(path)

def _d_c2(path):
    W, H = 1240, 470
    img = Image.new("RGB", (W, H), _P_WHITE); d = ImageDraw.Draw(img)
    # top coarse partition
    _pbox(d, (W / 2 - 230, 20, W / 2 + 230, 96), _P_NAVY,
          text="Full dataset  \u2192  coarse partition (k*)", tsize=20, bold=True)
    kids = ["Cluster A", "Cluster B", "Cluster C"]
    kw, kh, gap, y = 250, 70, 60, 150
    total = len(kids) * kw + (len(kids) - 1) * gap
    x0 = (W - total) / 2
    for i, k in enumerate(kids):
        x = x0 + i * (kw + gap)
        _parrow(d, (W / 2, 96 + 14), (x + kw / 2, y - 14), color=_P_BLUE, width=3)
        _pbox(d, (x, y, x + kw, y + kh), _P_ACCENT, text=k, tsize=18, bold=True)
        # sub-clusters
        subs = ["sub-1", "sub-2"] if i < 2 else ["sub-1", "sub-2", "sub-3"]
        sw = 105; sgap = 14; sy = y + kh + 26
        stot = len(subs) * sw + (len(subs) - 1) * sgap
        sx0 = x + (kw - stot) / 2
        for j, s in enumerate(subs):
            sx = sx0 + j * (sw + sgap)
            _parrow(d, (x + kw / 2, y + kh + 8), (sx + sw / 2, sy - 8), color=_P_LGREY, width=2)
            _pbox(d, (sx, sy, sx + sw, sy + 46), _P_LIGHT, outline=_P_BLUE,
                  text=s, tcolor=_P_NAVY, tsize=14, bold=True)
    # aggregation label
    d.text((W / 2 - 250, H - 70),
           "Cross-level SHAP aggregation (Prop. 6.1):  parent importance =\nsize-weighted average of children  (up to surrogate residual)",
           fill=_P_TEAL, font=_pf(17, True))
    img.save(path)

def _d_metrics(path):
    W, H = 1240, 430
    img = Image.new("RGB", (W, H), _P_WHITE); d = ImageDraw.Draw(img)
    groups = [("Ranking quality", ["Precision@K", "Recall@K", "NDCG@20  (principal)"], _P_BLUE),
              ("Diversity", ["Coverage = |\u222a R|/|I|", "Intra-List Div. (ILD)", "built into utility"], _P_TEAL),
              ("Clustering", ["Silhouette", "Davies\u2013Bouldin", "k-selection (elbow)"], _P_ACCENT)]
    cw, gap, x0, y, ch = 380, 40, 25, 30, 320
    for i, (h, items, col) in enumerate(groups):
        x = x0 + i * (cw + gap)
        _pbox(d, (x, y, x + cw, y + 54), col, text=h, tsize=20, bold=True)
        iy = y + 76
        for it in items:
            _pbox(d, (x + 14, iy, x + cw - 14, iy + 52), _P_LIGHT, outline=col,
                  text=it, tcolor=_P_NAVY, tsize=16, bold=False)
            iy += 64
    img.save(path)

def generate_assets():
    # extract genuine paper figures (architecture / workflows) directly from the source PDFs
    _extract("DyHuCoG A Dynamic Hypergraph Cooperative Game for Preference-aware Recommendation.pdf", 38, "dyhucog_workflow.png")
    _extract("DyHuCoG A Dynamic Hypergraph Cooperative Game for Preference-aware Recommendation.pdf", 52, "dyhucog_results.png")
    _extract("DyHuCoG A Dynamic Hypergraph Cooperative Game for Preference-aware Recommendation.pdf", 62, "dyhucog_covdiv.png")
    _extract("DyHuCoG A Dynamic Hypergraph Cooperative Game for Preference-aware Recommendation.pdf", 67, "dyhucog_waterfall.png")
    _extract("Game Theory Meets Explainable AI- An Enhanced Approach to Understanding Black Box Models Through Shapley Values.pdf", 53, "gamexai_workflow.png")
    _extract("Shapley Values for Explaining the Black Box Nature of Machine Learning Model Clustering.pdf", 136, "shapcluster_summary.png")
    _extract("Shapley Values for Explaining the Black Box Nature of Machine Learning Model Clustering.pdf", 137, "shapcluster_clusters.png")
    # generated diagrams
    _d_evolution(os.path.join(ASSETS, "evolution.png"))
    _d_motivation(os.path.join(ASSETS, "motivation.png"))
    _d_actionable(os.path.join(ASSETS, "actionable.png"))
    _d_c2(os.path.join(ASSETS, "c2_multilevel.png"))
    _d_metrics(os.path.join(ASSETS, "metrics.png"))

def _pic_fit(slide, path, x, y, maxw, maxh, align="center"):
    iw, ih = Image.open(path).size
    ar = iw / ih
    w, h = maxw, maxw / ar
    if h > maxh:
        h, w = maxh, maxh * ar
    if align == "center":
        px = x + (maxw - w) / 2
    elif align == "left":
        px = x
    else:
        px = x + (maxw - w)
    py = y + (maxh - h) / 2
    slide.shapes.add_picture(path, Inches(px), Inches(py), Inches(w), Inches(h))

# ----------------------------------------------------------------------------
# Top navigation bar (mirrors the example deck's tabbed nav)
# ----------------------------------------------------------------------------
def _section_of(title):
    t = title
    if "C1" in t: return ("Contributions", "C1")
    if "C2" in t: return ("Contributions", "C2")
    if "C3" in t: return ("Contributions", "C3")
    if any(k in t for k in ("Motivation", "Actionable", "Research Context")):
        return ("Introduction", None)
    if any(k in t for k in ("Recommendation & Clustering", "Limitations",
                            "Three Structuring", "Research Questions", "Three Contributions")):
        return ("Context & Problematic", None)
    if any(k in t for k in ("Datasets", "Data Splitting", "Baselines", "Hardware")):
        return ("Experimental Protocol", None)
    if any(k in t for k in ("Thesis", "Published", "Future", "Questions")):
        return ("Conclusion & Perspectives", None)
    return ("Conclusion & Perspectives", None)

def _topnav(slide, title):
    parts = ["Introduction", "Context & Problematic", "Experimental Protocol",
             "Contributions", "Conclusion & Perspectives"]
    active_main, active_sub = _section_of(title)
    by = 1.05
    _rect(slide, 0, Inches(by), SW, Inches(0.42), RGBColor(0xEE, 0xF2, 0xF7))
    _rect(slide, 0, Inches(by + 0.42), SW, Pt(1.2), BORDER)
    n = len(parts); gap = 0.16
    tw = (13.333 - gap * (n + 1)) / n
    x = gap
    for p in parts:
        active = (p == active_main)
        if active:
            _rect(slide, Inches(x), Inches(by + 0.03), Inches(tw), Inches(0.36), BLUE)
        _tb(slide, Inches(x), Inches(by + 0.03), Inches(tw), Inches(0.36), p,
            size=12.5, color=WHITE if active else GREY, bold=active,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        x += tw + gap
    if active_main == "Contributions":
        sy = 1.49
        _rect(slide, 0, Inches(sy), SW, Inches(0.32), RGBColor(0xF2, 0xF5, 0xF9))
        sx = 0.55
        for s in ["C1", "C2", "C3"]:
            act = (s == active_sub)
            w = 1.15
            if act:
                _rect(slide, Inches(sx), Inches(sy + 0.03), Inches(w), Inches(0.26), NAVY)
            else:
                _rect(slide, Inches(sx), Inches(sy + 0.03), Inches(w), Inches(0.26),
                      RGBColor(0xDD, 0xE6, 0xF0))
            _tb(slide, Inches(sx), Inches(sy + 0.03), Inches(w), Inches(0.26), s,
                size=11, color=WHITE if act else BLUE, bold=True,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            sx += w + 0.2

generate_assets()

# ============================================================================
# BUILD THE DECK
# ============================================================================
scores = []
Q = 0
def N():
    global Q
    Q += 1
    return Q

# --- 1. TITLE ---------------------------------------------------------------
title_slide(
    ["Cooperative Game Theory for Explainable AI in Recommendation Systems:",
     "A Shapley Framework for Actionable Insight"],
    ["PhD Viva Presented by: Mouad LOUHICHI",
     "Supervisor | Pr. Mohamed LAZAAR | PES | ENSIAS, Mohammed V University, Rabat",
     "Doctoral Studies Center in Information and Engineering Sciences and Technologies (ST2I)"],
    "Good morning. Thank you, President and Professors, for being here to examine my thesis. "
    "My name is Mouad Louhichi. I will present my PhD thesis: Cooperative Game Theory for "
    "Explainable Artificial Intelligence in Recommendation Systems, a Shapley framework for "
    "actionable insight, supervised by Professor Mohamed Lazaar."
)
Q = 1

# --- 2. OUTLINE -------------------------------------------------------------
section_slide("Viva Outline",
    "Introduction \u2192 Context & Problematic \u2192 Protocols \u2192 3 Contributions \u2192 Conclusion & Perspectives",
    "I will walk you through five parts: the introduction and motivation, the context and problematic, "
    "the experimental protocol shared across my three contributions, then the three contributions "
    "themselves, and finally the conclusion and perspectives. Each contribution follows the same "
    "structure: objectives, methodology, results, and findings.")

# --- 3. SECTION: INTRODUCTION ----------------------------------------------
section_slide("Introduction", "Recommenders \u00b7 The black-box problem \u00b7 Why XAI",
    "Let us begin with the introduction.")
content_slide(
    "Motivation: Three Questions",
    [("Ubiquity \u2013 How do opaque AI systems shape what billions of users see, buy, and watch every day?", 0),
     ("The black box \u2013 Why do state-of-the-art recommenders and clustering pipelines remain opaque to the users and designers who rely on them?", 0),
     ("Toward trust \u2013 How can transparency and accountability be built *as part of* the model, not bolted on afterwards?", 0),
     ("", 0),
     ("The central tension: as models gain expressive power they lose the transparency needed for trustworthy deployment.", 1)],
    "Three driving questions frame the thesis. First, ubiquity: opaque systems mediate what people "
    "see and which options become actionable. Second, the black box: even strong recommenders and "
    "clustering pipelines are hard to interrogate. Third, trust: transparency should be built into "
    "the modelling logic rather than added post-hoc. The core tension is that predictive power and "
    "interpretability tend to trade off.",
    subtitle="Where the thesis begins", num=N(),
    layout="image_right", image=os.path.join(ASSETS, "motivation.png"))

content_slide(
    "Actionable Insight \u2014 the Definition We Use",
    [("Definition 1.1 (Actionable insight)", 0, True),
     ("An explanation is actionable when it identifies at least one modifiable factor whose change is associated with a specifiable change in model output, and that factor is accessible to the decision-maker.", 1),
     ("", 0),
     ("Accessibility = expressed in the semantic vocabulary of the task domain:", 0),
     ("a physicochemical variable (wine), a pollution indicator (air quality), or a preference signal (recommendation),", 1),
     ("not an opaque latent code.", 1)],
    "We frame the whole thesis around the notion of actionable insight. An explanation is actionable "
    "when it points to a modifiable driver expressed in the language of the domain. This matters "
    "because we want explanations that support intervention, not merely describe what happened.",
    subtitle="Actionability, not just plausibility", num=N(),
    layout="image_right", image=os.path.join(ASSETS, "actionable.png"))

content_slide(
    "Research Context",
    [("Recommender systems evolved from similarity filters to complex representation-learning systems operating on sparse, high-dimensional, dynamic data.", 0),
     ("Matrix factorisation, neural CF, graph CNNs, then hypergraph recommenders \u2014 each step improved ranking but intensified the interpretability deficit.", 0),
     ("Why that deficit matters:", 0, True),
     ("Undermines user trust in the items they are shown", 1),
     ("Constrains debugging and scientific learning about the model", 1),
     ("Collides with emerging regulation: EU AI Act (2024), OECD AI Principles, GDPR [R1]", 1),
     ("Consequence: explanation becomes an accountability mechanism, not a usability add-on.", 0)],
    "The context is the progression from simple recommenders to hypergraph models: each step raised "
    "expressiveness while lowering transparency. That deficit matters for three reasons: user trust, "
    "debugging, and regulatory expectations. In these settings explanation is an accountability "
    "mechanism, not a usability add-on.",
    subtitle="From similarity filters to hypergraph recommenders", num=N(),
    layout="image_right", image=os.path.join(ASSETS, "evolution.png"))

# --- 4. SECTION: CONTEXT & PROBLEMATIC -------------------------------------
section_slide("Context & Problematic", "Three structuring limitations \u00b7 The research gap",
    "Now let me look more precisely at the problem this thesis addresses.")
content_slide(
    "Recommendation & Clustering Paradigms",
    [("Collaborative filtering \u2013 users who behaved similarly will value similar items (user- or item-based).", 0),
     ("Content-based \u2013 recommends items sharing attributes with a user profile.", 0),
     ("Hybrid \u2013 combines collaborative and content signals.", 0),
     ("Matrix factorisation \u2013 latent factors R \u2248 PQ\u1d40, compact but opaque.", 0),
     ("Graph-based \u2013 interaction graph with neighbourhood propagation (LightGCN, hypergraph).", 0)],
    "Quick orientation across the recommendation paradigms we build on. Each strengthens modelling "
    "but complicates interpretation. Matrix factorisation made latent dimensions opaque; graph models "
    "added structure but kept importance implicit; hypergraph models added higher-order relations but "
    "still assumed uniform message importance.",
    subtitle="The paradigm landscape", num=N())

content_slide(
    "Limitations of Classical Recommenders & Unsupervised Models",
    [("Data sparsity & scalability \u2013 the user\u2013item matrix is overwhelmingly empty.", 0),
     ("Cold-start \u2013 new users and items are structurally disadvantaged.", 0),
     ("Popularity bias & lack of diversity \u2013 exposure begets interaction, begets exposure (filter bubble).", 0),
     ("Absence of interpretability \u2013 the most fundamental limit: weak explanatory traction.", 0, True),
     ("For clustering: methods privilege local OR global explanation, not both; they struggle to scale; explanations rarely preserve coherence across resolutions.", 1)],
    "Four classical limitations. The most fundamental for this thesis is the absence of "
    "interpretability. In clustering specifically, existing approaches tend to privilege either local "
    "or global explanation, struggle to scale, and rarely remain coherent when explanations are "
    "generated at multiple resolutions.",
    subtitle="Why the gap exists", num=N())

content_slide(
    "Three Structuring Limitations (the Problem Statement)",
    [("1. Lack of explainability \u2013 complex models remain hard to interpret faithfully and in an actionable way.", 0, True),
     ("2. Difficulty of scaling \u2013 local explanations do not transfer naturally to hierarchical structures or large datasets.", 0, True),
     ("3. Weak integration into learning \u2013 most explanations stay post-hoc and do not shape model dynamics or the accuracy-diversity-context trade-off.", 0, True),
     ("", 0),
     ("The thesis gap in one sentence: the literature still lacks a single cooperative-attribution framework that can explain clustering faithfully, remain coherent under hierarchy, and then operate as an in-training signal in recommendation rather than as a purely post-hoc device.", 1)],
    "Three structuring problems define the thesis: lack of explainability, poor scalability of local "
    "explanations, and weak integration into the learning loop. The gap is that no existing framework "
    "carries a single cooperative-attribution logic across clustering, hierarchical analysis, and "
    "in-training recommendation.",
    subtitle="Problem statement", num=N())

content_slide(
    "Research Questions (RQ1\u2013RQ5) and the Overall Aim",
    [("Aim: develop, justify, and evaluate a cooperative-game-theoretic perspective for XAI in clustering and recommendation, using Shapley attribution as both an explanatory mechanism and an in-training signal.", 0, True),
     ("RQ1 \u00b7 How can Shapley values explain black-box clustering faithfully at instance and cluster level?", 0),
     ("RQ2 \u00b7 How can this extend to large-scale, hierarchical clustering without losing tractability or consistency?", 0),
     ("RQ3 \u00b7 Can cooperative attribution move beyond post-hoc and enter the learning dynamics of graph recommenders?", 0),
     ("RQ4 \u00b7 Can a recommender jointly optimise ranking accuracy, context, and diversity when importance is estimated by a cooperative-game utility?", 0),
     ("RQ5 \u00b7 What emerges when clustering explanation and recommendation learning are read as two stages of one cooperative-game perspective?", 0)],
    "The aim is to build one cooperative-game perspective for XAI across clustering and "
    "recommendation. Five research questions form the spine: RQ1-RQ2 on clustering, RQ3-RQ4 on "
    "recommendation, RQ5 at thesis level. Each maps to a chapter.",
    subtitle="Research aim and questions", num=N())

content_slide(
    "The Three Contributions",
    [("C1 \u2013 Explainable black-box clustering: PCA\u2013K-Means\u2013LightGBM\u2013TreeSHAP pipeline.", 0, True),
     ("\u2192 Wine Quality dataset: faithful instance- and cluster-level feature attribution.", 1),
     ("C2 \u2013 Enhanced multi-level XAI for large-scale clustering with cross-level SHAP aggregation.", 0, True),
     ("\u2192 Beijing Multi-Site Air Quality: hierarchical attribution consistency under scale.", 1),
     ("C3 \u2013 DyHuCoG: a Dynamic Hypergraph Cooperative Game for preference-aware recommendation.", 0, True),
     ("\u2192 MovieLens-1M & Amazon-Book: preference-aware Monte Carlo Shapley as an in-training signal.", 1),
     ("", 0),
     ("Thesis-level claim: cooperative game theory functions as a shared attribution perspective for explanation, optimisation, and intervention.", 2, True)],
    "Three contributions, one thread. C1 establishes Shapley-based explanation for black-box "
    "clustering. C2 scales that logic to hierarchy and large data. C3 replaces post-hoc attribution "
    "with an in-training signal in a hypergraph recommender. Together they support the thesis-level "
    "claim: cooperative game theory is a shared attribution perspective.",
    subtitle="The research arc", num=N())

# --- 5. SECTION: PROTOCOLS --------------------------------------------------
section_slide("Experimental Protocol", "Datasets \u00b7 Baselines \u00b7 Metrics \u00b7 Hardware",
    "Before the contributions, let me briefly cover the experimental setup shared across all three.")

table_slide(
    "Datasets Used Throughout the Thesis",
    ["Dataset", "Scale", "Type", "Role"],
    [["Wine Quality (Portugal, vinho verde)", "4,898 \u00d7 11", "Tabular, numeric features", "C1 \u2013 single-level clustering"],
     ["Beijing Multi-Site Air Quality", "383,585 \u00d7 11", "Tabular, pollutant + meteorology", "C2 \u2013 multi-level clustering"],
     ["MovieLens-1M", "6,040 u / 3,706 i / 1,000,209 int", "Implicit-feedback", "C3 \u2013 DyHuCoG"],
     ["Amazon-Book", "52,643 u / 91,599 i / 2,984,108 int", "Implicit-feedback (very sparse)", "C3 \u2013 DyHuCoG"]],
    "Two clustering datasets and two recommendation datasets. The clustering datasets are chosen for "
    "semantically interpretable features (so attribution is meaningful), the recommendation datasets "
    "are standard benchmarks with established baselines. Density: MovieLens-1M 0.0447, Amazon-Book 0.0006.",
    subtitle="Four datasets, two regimes", num=N())

content_slide(
    "Dataset Characteristics & Statistics",
    [("Wine Quality (vinho verde) - 4,898 rows x 11 physico-chemical features (acidity, sugar, pH, alcohol, sulphur); quality 3-9.", 0),
     ("Beijing Air Quality - 383,585 hourly records x 11 pollutants + meteorology across 12 sites; strong temporal structure.", 0),
     ("MovieLens-1M - 6,040 users / 3,706 items / 1,000,209 implicit ratings; density 4.47%; evaluated at K in {5,10,20}.", 0),
     ("Amazon-Book - 52,643 users / 91,599 items / 2,984,108 interactions; density 0.06% (very sparse).", 0),
     ("Why these: clustering sets have semantically interpretable features (actionable attribution); rec sets are standard benchmarks with established baselines [R6,R7,R8].", 0, True),
     ("Preprocessing: ratings > 3 as positive; popularity-aware negative sampling q(i) propto f_i^eta; five fixed seeds for reproducibility.", 1)],
    "The four datasets span two regimes. The two clustering datasets are chosen so attribution is "
    "expressed in a meaningful domain vocabulary; the two recommendation datasets are standard benchmarks. "
    "Density differs by two orders of magnitude (MovieLens-1M 4.47% vs Amazon-Book 0.06%), which is why "
    "sparsity is a first-class experimental axis.",
    subtitle="What each dataset brings", num=N())

content_slide(
    "Data Splitting & Preprocessing",
    [("Clustering: five-fold cross-validation for surrogate/attribution stability.", 0),
     ("Recommendation: user-level, temporal split \u2014 70% train / 10% validation / 20% test.", 0),
     ("Leave-one-out evaluation: the latest test positive per user is the target, ranked against negatives.", 0),
     ("Implicit feedback conversion: MovieLens-1M ratings > 3 treated as positive.", 0),
     ("Popularity-aware negative sampling: q(i) \u221d f_i^\u03b7 for harder training contrasts.", 0),
     ("Reproducibility: seeds {42, 43, 44, 45, 46}; early-stopping patience 20; fixed seeds.", 0)],
    "The splitting is designed against leakage. Clustering uses cross-validation; recommendation uses "
    "a temporal per-user holdout with leave-one-out. MovieLens ratings above 3 become implicit "
    "positives, negatives are sampled with a popularity-aware distribution, and five seeds are used "
    "for reproducibility.",
    subtitle="Leakage control and reproducibility", num=N())

content_slide(
    "Baselines & Evaluation Metrics",
    [("Clustering benchmarks: LIME-based surrogate explanation pipeline (interpretability comparator).", 0, True),
     ("Recommendation benchmarks: MF, NCF, LightGCN, RecDCL, HCCF, HPCF (strongest reference).", 0, True),
     ("", 0),
     ("Ranking: Precision@K, Recall@K, NDCG@20 (principal metric).", 0),
     ("System diversity: Catalogue Coverage = |\u22c3 R_u| / |I|.", 0),
     ("List diversity: Intra-List Diversity (ILD) \u2014 average pairwise dissimilarity; built into the utility.", 0),
     ("Clustering quality: Silhouette coefficient, Davies\u2013Bouldin index.", 0)],
    "Baselines span classical, neural, graph, and hypergraph methods to isolate the contribution of "
    "cooperative attribution. Metrics capture ranking quality plus coverage and intra-list diversity "
    "for the recommendation side, and Silhouette and Davies-Bouldin for the clustering side. ILD is "
    "not decorative: it is part of the DyHuCoG coalition utility.",
    subtitle="Comparing under one evaluative frame", num=N(),
    layout="image_right", image=os.path.join(ASSETS, "metrics.png"))

content_slide(
    "Hardware & Software",
    [("CPU: Intel Core i9-14900K, 24 cores \u2014 clustering, preprocessing, data loading.", 0),
     ("GPU: NVIDIA GeForce RTX 4090, 24 GB (Ada Lovelace) \u2014 DyHuCoG training & inference.", 0),
     ("RAM: 48 GB; Storage: 2 TB SSD.", 0),
     ("Python 3.8; scikit-learn (PCA, K-Means), LightGBM (surrogate), SHAP (TreeSHAP), PyTorch 2.0.1 (GNN/HGNN), NumPy/SciPy/pandas (stats).", 0),
     ("Altair for interactive SHAP visualisation; metrics at K \u2208 {5, 10, 20}.", 0)],
    "Clustering runs on standard workstation CPU; DyHuCoG needs a GPU but stays within academic "
    "resources. The stack is scikit-learn + LightGBM + SHAP for clustering, PyTorch for the "
    "recommendation model, and SciPy for significance testing.",
    subtitle="Reproducible environment", num=N())

# --- 6. SECTION: CONTRIBUTION I --------------------------------------------
section_slide("Contribution I \u2014 Explainable Black-Box Clustering",
              "PCA \u2192 K-Means \u2192 LightGBM \u2192 SHAP \u00b7 Wine Quality \u00b7 RQ1",
    "Let us move to the first contribution: explaining black-box clustering with Shapley values.")

content_slide(
    "C1 \u00b7 Research Gap & Objectives",
    [("Gap: Shapley explanation is dominant in supervised tasks, but clustering remains comparatively under-explained.", 0, True),
     ("Existing clustering-interpretability methods privilege local or global explanation, not both.", 1),
     ("They often fail to scale or to preserve coherence across clusters.", 1),
     ("", 0),
     ("Objectives:", 0, True),
     ("Build a pipeline yielding cluster-level explanation from an unsupervised partition while preserving feature-level attribution.", 1),
     ("", 0),
     ("Research Questions (RQ1):", 0, True),
     ("RQ1a - Can Shapley values attribute cluster membership faithfully at instance and cluster level? [R2]", 1),
     ("RQ1b - Is Shapley attribution more stable and coherent than LIME across clusters? [R5]", 1),
     ("Preserve the semantics of the original feature space even though clustering uses a reduced dimensionality.", 1),
     ("Justify, on theoretical and literature grounds, why Shapley-based attribution is better than LIME for stable cluster explanation.", 1)],
    "Clustering is a natural starting point because it strips the explanatory problem to essentials: "
    "the model creates its own structure, so the meaning of a cluster must be inferred after the fact. "
    "The gap is that explainable clustering is fragmented. Our objectives are to build a pipeline that "
    "explains a partition while preserving original feature semantics, and to argue for Shapley over LIME.",
    subtitle="RQ1", num=N())

content_slide(
    "C1 \u00b7 Cooperative-Game Formulation for Clustering",
    [("Player set: N = F \u2014 each feature is a player.", 0, True),
     ("Value function: v(S) = Silhouette( KMeans(X_S, k*) ) \u2014 how well data cluster using only features in S.", 0),
     ("A feature's Shapley value = its expected marginal contribution to clustering quality over all coalition orders [R2].", 0),
     ("", 0),
     ("Why Silhouette: bounded, normalised, semantically intuitive.", 0),
     ("Alternatives (Davies\u2013Bouldin, Calinski\u2013Harabasz) possible but Silhouette is the interpretable choice.", 1),
     ("Direct evaluation for every coalition is intractable \u2192 we need a bridge.", 0)],
    "We frame clustering as a cooperative game where features are players and the value function is "
    "the Silhouette of clustering on the feature subset. A feature gets high attribution when its "
    "presence consistently improves separation. But exact evaluation is combinatorial, so we need a "
    "tractable bridge.",
    subtitle="Features as players", num=N())

content_slide(
    "C1 \u00b7 The LightGBM Bridge for Tractable Attribution",
    [("Once K-Means produces cluster labels, train a LightGBM multiclass surrogate to predict those labels from the original features [R4].", 0),
     ("Apply TreeSHAP to the surrogate \u2014 exact, fast tree-based attribution in the original semantic feature space.", 0),
     ("Why this bridge matters:", 0, True),
     ("Direct TreeSHAP on K-Means is impossible (it explains tree models, not centroids).", 1),
     ("Explaining the PCA representation would move attribution away from interpretable variables.", 1),
     ("The surrogate preserves the chemistry/pollution vocabulary that makes the analysis actionable.", 1),
     ("Validity condition: surrogate fidelity is high (macro-F1 \u2248 0.82 on held-out cluster labels).", 1, True)],
    "The bridge is the heart of the method. We convert an unsupervised partition into a supervised "
    "prediction task by training a LightGBM classifier to reproduce the cluster labels, then run "
    "TreeSHAP on it. This keeps attribution in the original feature space. It is only valid if the "
    "surrogate is faithful, so we treat macro-F1 around 0.82 as the practical floor.",
    subtitle="Unsupervised partition \u2192 supervised explanation", num=N())

content_slide(
    "C1 \u00b7 Proposed Pipeline: PCA\u2013KMeans\u2013LightGBM\u2013SHAP",
    [("Stage 1 \u2013 PCA: stabilise geometry + visual diagnostic. NOT the explanatory space.", 0),
     ("Stage 2 \u2013 K-Means++ with multi-criteria k selection (elbow, Silhouette, Davies\u2013Bouldin).", 0),
     ("Stage 3 \u2013 LightGBM surrogate trained on original features to predict cluster labels.", 0),
     ("Stage 4 \u2013 TreeSHAP attribution in the original feature space.", 0),
     ("Stage 5 \u2013 Aggregate into global importance, cluster-specific profiles, and local force plots.", 0),
     ("", 0),
     ("Complexity: dominated by PCA and repeated K-Means; TreeSHAP scales with tree count/depth, not exponentially in features.", 1)],
    "The five-stage pipeline. PCA is a computational and visual aid, deliberately not the explanatory "
    "space. K-Means defines the partition, the surrogate restores tractable supervised prediction, and "
    "TreeSHAP returns explanation to the original variables. This sequence is explicit because each "
    "stage serves a distinct purpose.",
    subtitle="Five stages", num=N(),
    layout="image_right", image=os.path.join(ASSETS, "gamexai_workflow.png"))

content_slide(
    "C1 \u00b7 Optimal Cluster Selection \u2014 a Deliberate Choice",
    [("Multi-criteria evaluation across k \u2208 {2..10} using elbow, Silhouette, Davies\u2013Bouldin.", 0),
     ("We select k* = 3 \u2014 even though it is NOT geometrically optimal:", 0, True),
     ("k = 2: Silhouette 0.214, Davies\u2013Bouldin 1.775 (better separation).", 1),
     ("k = 3: Silhouette 0.144, Davies\u2013Bouldin 2.097 (weaker separation).", 1),
     ("Why: three clusters give a semantically richer oenological partition \u2192 more discriminative, more actionable explanation.", 1, True),
     ("Note: the higher Silhouette \u2248 0.63 belongs to the Beijing dataset (C2), not the wine partition.", 1)],
    "This is an important, honest point. We select three clusters even though two gives better raw "
    "geometry metrics. The justification is interpretability: three clusters support three distinct, "
    "chemically meaningful narratives rather than one broad split. I want to avoid any confusion: the "
    "0.63 Silhouette that appears later belongs to Beijing, not the retained wine partition.",
    subtitle="Interpretability over geometry", num=N())

content_slide(
    "C1 \u00b7 Global Feature Importance (Wine Quality)",
    [("Global SHAP ranking (highest \u2192 lower):", 0, True),
     ("density \u2192 pH \u2192 fixed acidity \u2192 sulfur-dioxide variables \u2192 alcohol", 1),
     ("", 0),
     ("Significance:", 0, True),
     ("Dominant drivers are density, pH, fixed acidity, sulfur-dioxide, alcohol \u2014 variables intimately related to structure, preservation, and sensory balance.", 1),
     ("This is NOT an arbitrary classifier artefact: it recovers a chemically interpretable hierarchy.", 1),
     ("The explanatory signal is concentrated, not diffuse.", 1)],
    "The global SHAP ranking is dominated by a small set of physicochemical variables: density, pH, "
    "fixed acidity, sulfur dioxide, alcohol. This is important because it matches oenological "
    "knowledge, showing the surrogate-based pipeline recovers a chemically interpretable feature "
    "hierarchy rather than fitting arbitrary artefacts.",
    subtitle="A chemistry-consistent hierarchy", num=N(),
    layout="image_right", image=os.path.join(ASSETS, "shapcluster_summary.png"))

content_slide(
    "C1 \u00b7 Cluster-Specific Profiles (Three Signatures)",
    [("Three clusters show distinct explanatory signatures, not minor variations.", 0, True),
     ("Cluster 0 \u2013 density + sulfur-dioxide-related variables.", 1),
     ("Cluster 1 \u2013 acidity and pH-related effects.", 1),
     ("Cluster 2 \u2013 a different balance of acidity, alcohol, and related chemical attributes.", 1),
     ("", 0),
     ("Interpretation: the same small set of important variables recurs across clusters, but with different relative weights within each cluster.", 1),
     ("The global ranking translates into cluster-level heterogeneity \u2014 the partition is globally interpretable AND internally differentiated.", 1)],
    "The cluster-specific profiles show that the retained solution is not only globally interpretable "
    "but internally differentiated. Each cluster has its own chemical signature \u2014 sulfur and density "
    "for one, acidity and pH for another, a different acid-alcohol balance for the third. This is "
    "exactly the kind of actionable insight the thesis seeks.",
    subtitle="Cluster-level heterogeneity", num=N(),
    layout="image_right", image=os.path.join(ASSETS, "shapcluster_clusters.png"))

content_slide(
    "C1 \u00b7 SHAP vs. LIME \u2014 Why Cooperative Attribution",
    [("SHAP grounds attribution in a cooperative-game allocation rule; LIME fits a local surrogate.", 0),
     ("Stability: SHAP is more stable when the surrogate is faithful; LIME depends on perturbation design [R5].", 0),
     ("Theoretical guarantees: SHAP satisfies efficiency, symmetry, null player, additivity \u2014 LIME has no equivalent.", 0),
     ("Local/global coherence: SHAP supports both; LIME is primarily local.", 0),
     ("Cluster comparison: SHAP is strong; LIME limited.", 0),
     ("Caveat: in the surrogate pipeline, efficiency holds w.r.t. the LightGBM output, not the Silhouette value directly.", 1)],
    "Why Shapley over LIME? The four axioms give SHAP a normative basis: efficiency ensures "
    "completeness, symmetry fairness, the null-player property no spurious credit, and additivity "
    "composition. LIME's local surrogates are sensitive to perturbation design, which becomes a global "
    "problem once we compare across clusters. One honest caveat: efficiency holds w.r.t. the surrogate, "
    "not the Silhouette value itself.",
    subtitle="Axiomatic attribution vs local surrogate", num=N())

content_slide(
    "C1 \u00b7 Findings & Limitations",
    [("Achieved:", 0, True),
     ("A cluster-level explanation anchored to individual feature contributions.", 1),
     ("Explanations returned to the original variables, not to an abstract principal-component basis.", 1),
     ("Theoretically and literature-backed case for Shapley over LIME for stable cluster attribution.", 1),
     ("Limitations:", 0, True),
     ("Fidelity depends on the LightGBM surrogate.", 1),
     ("Tabular data only \u2014 not automatically transferable to image/text clustering.", 1),
     ("Single-level structure \u2014 cannot yet address hierarchical coherence. That is C2\u2019s point of departure.", 1)],
    "Contribution I answers RQ1: Shapley values can explain black-box clustering faithfully to a "
    "high-fidelity surrogate and coherently at cluster level. The LIME comparison is presented as a "
    "theoretical and literature-backed argument rather than a full empirical bake-off. Limits: surrogate "
    "dependence, tabular-only, and single-level structure \u2014 which is exactly what C2 tackles.",
    subtitle="RQ1 answered", num=N())

# --- 7. SECTION: CONTRIBUTION II -------------------------------------------
section_slide("Contribution II \u2014 Enhanced Multi-Level XAI for Large-Scale Clustering",
              "Beijing Air Quality \u00b7 Hierarchical attribution consistency \u00b7 RQ2",
    "This brings us to the second contribution: scaling the explanation logic to multi-level, "
    "large-scale clustering.")

content_slide(
    "C2 \u00b7 Research Gap & Objectives",
    [("Gap: once clustering is multi-level, feature importance must stay interpretable within a cluster, across sub-clusters, and across the hierarchy as a whole.", 0, True),
     ("Large-scale data make exact explanation computationally burdensome.", 1),
     ("Flat explanation may be true yet incomplete \u2014 it cannot show how importance changes inside a cluster.", 1),
     ("", 0),
     ("Objectives:", 0, True),
     ("A genuinely multi-level explanatory workflow (not a rerun of the C1 pipeline).", 1),
     ("", 0),
     ("Research Questions (RQ2):", 0, True),
     ("RQ2a - Can cross-level SHAP aggregation stay consistent as the hierarchy deepens? [R2]", 1),
     ("RQ2b - Does the explanation remain tractable and coherent at large scale? [R9]", 1),
     ("A formal cross-level consistency argument (Proposition 6.1).", 1),
     ("Validation on a structurally different, large-scale dataset (Beijing).", 1)],
    "C2 asks whether the C1 logic survives scale and hierarchy. Large real-world data contain "
    "structure at more than one granularity. A flat explanation cannot show how the importance of a "
    "variable reconfigures as you zoom into a region. Our objectives are a real multi-level workflow, "
    "a formal consistency argument, and validation on a different large-scale dataset.",
    subtitle="RQ2", num=N())

content_slide(
    "C2 \u00b7 Multi-Level Clustering Architecture",
    [("Recursive / nested design: coarse clustering on the full dataset, then subdivide each cluster where appropriate.", 0),
     ("For each level, train a level-specific surrogate and compute SHAP in the SAME original feature space.", 0),
     ("Cross-level aggregation is NOT a naive average \u2014 it respects cluster size and nesting structure.", 0, True),
     ("Parent-level attribution = an expectation over the explanatory structure of its descendants.", 1),
     ("The hierarchy is a pragmatic analytical device, not a claim that the data have a true ontological hierarchy.", 1)],
    "The multi-level architecture proceeds recursively. A coarse clustering is learned first, then each "
    "cluster is subdivided. Each level gets its own surrogate and SHAP values in the same feature space. "
    "The aggregation respects size and nesting structure. Importantly, we treat the hierarchy as an "
    "analytical device, not a metaphysical claim \u2014 this prevents overclaiming.",
    subtitle="Nested clustering as an analytical tool", num=N(),
    layout="image_right", image=os.path.join(ASSETS, "c2_multilevel.png"))

content_slide(
    "C2 \u00b7 Formal Result: Hierarchical Attribution Consistency (Proposition 6.1)",
    [("Let \u03a6^(l,c)_j = E_{x~c}[ |\u03c6_j^(l)(x)| ] be the expected absolute SHAP importance of feature j at level l in cluster c.", 0),
     ("Let w_c' = |c'| / |c| be the relative size of child c' within parent c.", 0),
     ("For a strict nested hierarchy on a consistent feature space:", 0, True),
     ("\u03a6^(l,c)_j = \u03a3_{c'\u2208child(c)} w_c' \u00b7 \u03a6^(l+1,c')_j + \u03b5_j", 1, True),
     ("\u03b5_j is a residual from surrogate mismatch and vanishes under perfect surrogate fidelity.", 1),
     ("Derived via the law of total expectation (children partition the parent).", 1),
     ("Does NOT imply explanations are identical across levels \u2014 it implies differences can be interpreted, not dismissed as inconsistency.", 1)],
    "Proposition 6.1 is the chapter's original formalisation. It says a parent's expected absolute "
    "importance is the size-weighted average of its children's, up to a surrogate residual. It follows "
    "from the law of total expectation because children partition the parent. Importantly, it does not "
    "claim explanations are identical across levels; it lets us interpret differences rather than treat "
    "them as inconsistency.",
    subtitle="The cross-level consistency claim", num=N())

content_slide(
    "C2 \u00b7 Results: Cluster Quality (Beijing)",
    [("Full dataset, k = 3 (strong convergence on the multi-criteria evaluation).", 0, True),
     ("Silhouette \u2248 0.63 \u2014 materially stronger separation than the wine study.", 1),
     ("Davies\u2013Bouldin \u2248 0.55 \u2014 low between-cluster ambiguity.", 1),
     ("PCA projection (2 components) used only for visual inspection.", 1),
     ("", 0),
     ("Sensitivity: conclusions robust to modest variation in k, projection dimensionality, and surrogate depth; only low-ranked variables shift.", 1)],
    "On Beijing, the multi-criteria evaluation converges much more strongly than on wine: three "
    "clusters with a Silhouette of about 0.63 and Davies-Bouldin of 0.55. This is a markedly clearer "
    "separation. The explanatory conclusions are robust to reasonable parameter changes \u2014 the dominant "
    "meteorological and pollutant drivers remain stable.",
    subtitle="Much clearer separation than wine", num=N())

content_slide(
    "C2 \u00b7 Global Feature Importance (Beijing)",
    [("Global SHAP ranking (highest \u2192 lower):", 0, True),
     ("temperature \u2192 dew point \u2192 pressure \u2192 CO \u2192 NO2 \u2192 PM10 \u2192 PM2.5", 1),
     ("", 0),
     ("Significance:", 0, True),
     ("It is NOT simply pollutant concentrations that matter \u2014 meteorological variables play a structurally central role.", 1),
     ("Temperature, dew point, and pressure condition dispersion, trapping, and photochemical behaviour.", 1),
     ("This is exactly the kind of insight flat descriptive summaries often fail to make explicit.", 1)],
    "The global ranking here is analytically rich: temperature, dew point, and pressure dominate, "
    "followed by CO, NO2, PM10, and PM2.5. The point is that meteorological variables are structurally "
    "central because they condition dispersion, trapping, and photochemistry. That is a domain insight "
    "that a flat centroid summary would miss.",
    subtitle="Meteorology conditions pollution regimes", num=N())

content_slide(
    "C2 \u00b7 Three Pollution Regimes (Force Plots)",
    [("Regime A \u2013 warm photochemical events: ozone, temperature, dew point prominent (summer photochemical smog).", 0),
     ("Regime B \u2013 wintertime smog: CO, SO2, particulate matter dominate; low wind speed suppresses dispersion (stagnant cold air).", 0),
     ("Regime C \u2013 comparatively clean air events: favourable meteorology, weak pollutant pushes.", 0),
     ("", 0),
     ("The framework shows not only that these regimes exist, but which variable combinations define them.", 1)],
    "The force plots reveal three representative regimes: warm photochemical events, wintertime smog, "
    "and cleaner dispersed events. The interpretative value is showing not only that these regimes "
    "exist, but which combinations of variables define each one \u2014 that's the actionable insight.",
    subtitle="Three interpretable air-quality regimes", num=N())

content_slide(
    "C2 \u00b7 The Multi-Level Insight",
    [("At the coarse level, temperature and dew point dominate \u2014 they differentiate broad atmospheric regimes.", 0),
     ("Within individual clusters, CO, SO2, PM10, wind speed, pressure, or ozone become more discriminative.", 0),
     ("", 0),
     ("This change is NOT contradictory to the global explanation \u2014 it is exactly what a multi-level explanation should reveal.", 0, True),
     ("Parent-level story = regime selection. Cluster-level story = variation within a regime.", 1),
     ("A variable can be globally important yet locally uninformative within a sub-cluster \u2014 the question is at what level it matters.", 1)],
    "This is the conceptual payoff. At the coarse level the story is regime selection (temperature, "
    "dew point); within clusters the discriminative variables shift to CO, SO2, PM10, wind speed. That "
    "is not inconsistency \u2014 it is what multi-level explanation is for. A variable can matter globally "
    "and not locally, so the real question is at what level it matters.",
    subtitle="Why flat explanation is insufficient", num=N())

content_slide(
    "C2 \u00b7 Cross-Dataset Generalisation",
    [("Wine study: small, dense, chemically correlated.", 0),
     ("Beijing study: large, noisy, strongly affected by temporal and meteorological variation.", 0),
     ("", 0),
     ("The same explanatory logic remains productive in both \u2192 the contribution is not tied to one domain-specific peculiarity.", 0, True),
     ("Comparison with SHAP-based clustering literature:", 1),
     ("Beijing Silhouette \u2248 0.63 vs Gramegna & Giudici credit-risk SHAP-space 0.37 \u2014 comparatively well separated.", 1),
     ("LIME comparator: weaker structural coherence, less stable local narratives for hierarchical reasoning [R9].", 1)],
    "The fact that the same logic works on both a small chemical dataset and a large noisy "
    "environmental one supports generality. Compared with prior SHAP-based clustering work, our Beijing "
    "partition is comparatively well separated, and our claim about LIME is that its local-surrogate "
    "logic is not a stable enough basis for multi-level reasoning.",
    subtitle="Generalisable beyond one domain", num=N())

content_slide(
    "C2 \u00b7 Findings & Limitations",
    [("Achieved:", 0, True),
     ("Scalable, multi-granular explanation without collapsing into a single flat summary.", 1),
     ("A formal consistency argument for cross-level attribution.", 1),
     ("Validation on a structurally different, large-scale dataset.", 1),
     ("Limitations:", 0, True),
     ("Still static clustering, even though Beijing data are temporally structured.", 1),
     ("Surrogate-based SHAP + representative-instance reporting compress observation-level variation.", 1),
     ("Tabular data only.", 1)],
    "C2 answers RQ2: Shapley-based clustering explanation can scale to hierarchical, large-scale "
    "settings without losing interpretive coherence, provided the hierarchy is modelled explicitly and "
    "the approximation is transparent. The main limitation is that the clustering is still static \u2014 which "
    "becomes the contrast with C3, where attribution is refreshed dynamically during training.",
    subtitle="RQ2 answered", num=N())

# --- 8. SECTION: CONTRIBUTION III ------------------------------------------
section_slide("Contribution III \u2014 DyHuCoG: A Dynamic Hypergraph Cooperative Game",
              "Preference-aware Shapley in hypergraph message passing \u00b7 RQ3 & RQ4",
    "The third and principal contribution introduces DyHuCoG, where Shapley attribution becomes an "
    "in-training signal inside a hypergraph recommender.")

content_slide(
    "C3 \u00b7 Research Gap & Objectives",
    [("Gap: graph and hypergraph recommenders treat message importance as either uniform or attention-weighted, without a principled marginal-contribution account.", 0, True),
     ("Diversity is often a secondary objective or a re-ranking heuristic.", 1),
     ("Interpretability is added after prediction, not integrated into the learning objective.", 1),
     ("", 0),
     ("Objectives:", 0, True),
     ("Formulate recommendation as a cooperative game over users, items, and contexts.", 1),
     ("Embed preference-aware Monte Carlo Shapley into hypergraph message passing.", 1),
     ("Improve ranking quality, coverage, and intra-list diversity jointly.", 1),
     ("", 0),
     ("Research Questions (RQ3 & RQ4):", 0, True),
     ("RQ3 - Can cooperative attribution move beyond post-hoc into the learning dynamics? [R6]", 1),
     ("RQ4 - Can a recommender jointly optimise ranking, context and diversity via a cooperative utility? [R7]", 1)],
    "C3 is the flagship. The gap: existing hypergraph recommenders assume importance is implicit in "
    "propagation, diversity is secondary, and interpretability is bolted on. Our objectives: formulate "
    "recommendation as a cooperative game, inject preference-aware Shapley into message passing, and "
    "jointly improve ranking, coverage, and intra-list diversity.",
    subtitle="RQ3 & RQ4", num=N())

content_slide(
    "C3 \u00b7 Problem Formulation \u2014 Recommendation as a Cooperative Game",
    [("Player set: N = U \u222a I \u222a C (users, items, contexts).", 0, True),
     ("Hypergraph H = (V, E, W); V = U \u222a I \u222a C; W = dynamic edge weights from Shapley estimates.", 0),
     ("Coalition S \u2286 N represents entities participating in a recommendation episode.", 0),
     ("Coalition value v(S) measures the quality of the recommendation outcome achievable by S.", 0),
     ("Top-N task: produce a ranked list L_u balancing relevance, diversity, and contextual fit [R6].", 0)],
    "We model recommendation as a cooperative game whose players are users, items, and contexts. The "
    "hypergraph encodes user-item-context interactions as hyperedges. A coalition value measures how "
    "good the recommendation is for that episode. This parallels the clustering formulation but with a "
    "recommendation-oriented value function.",
    subtitle="Users, items, contexts as players", num=N())

content_slide(
    "C3 \u00b7 Multi-Objective Coalition Utility",
    [("v(S) = \u03b1 \u00b7 NDCG@20(S) + \u03b2 \u00b7 Diversity(S) + \u03b3 \u00b7 ContextScore(S)", 0, True),
     ("with \u03b1 + \u03b2 + \u03b3 = 1. The same trade-off the recommender must satisfy is the trade-off from which attribution is computed \u2014 explanatory game and predictive objective are aligned by design.", 1),
     ("Preference-weighted variant: v_pref(S) = v(S) + \u03bb_pref \u00b7 \u03a3_{(u,i)\u2208S} sim(u,i)", 0, True),
     ("\u03b1 = 0.60, \u03b2 = 0.25, \u03b3 = 0.15; \u03bb_pref = 0.20 \u2014 tuned by grid search, stable (< 1.5% variance in NDCG@20).", 1),
     ("Coalition evaluation is scoped to the interaction episode (a few dozen players), not the full catalogue.", 1)],
    "The coalition utility combines ranking quality, diversity, and contextual relevance. Crucially, the "
    "same trade-off the model must satisfy is the trade-off from which attribution is computed, so the "
    "explanatory game and the predictive objective are aligned by design. We add a preference bonus for "
    "user-item agreement. Weights are tuned by grid search, and coalition evaluation is scoped to the "
    "episode.",
    subtitle="The utility the game is built on", num=N())

content_slide(
    "C3 \u00b7 Preference-Aware Monte Carlo Shapley Estimation",
    [("Exact Shapley is combinatorial and infeasible for realistic systems.", 0),
     ("Monte Carlo estimator: \u03c6\u0302_j = (1/M) \u03a3_m [ v(S_m \u222a {j}) \u2212 v(S_m) ]", 0, True),
     ("Preference-aware variant: \u03c6\u0302_j^pref = (1/M) \u03a3_m [ v_pref(S_m \u222a {j}) \u2212 v_pref(S_m) ]", 0, True),
     ("Unbiased; variance = \u03c3\u00b2/M \u2192 MSE decays O(1/M), absolute error O(1/\u221aM).", 1),
     ("M = 50 selected: MSE \u2248 1.4\u00d710\u207b\u2075, ~99% accuracy on MovieLens-1M.", 1, True),
     ("Estimates refreshed every 10 batches (~49 updates/epoch) and smoothed by exponential moving average \u2014 adaptive but not hypersensitive.", 1)],
    "Exact computation is infeasible, so we use a Monte Carlo estimator and its preference-aware "
    "variant. It is unbiased with variance decaying as 1 over M. We choose M = 50 as the practical "
    "balance between accuracy and cost, and refresh every ten batches with exponential smoothing so "
    "attribution is adaptive but training stays stable.",
    subtitle="Dynamic attribution, not a static diagnostic", num=N())

content_slide(
    "C3 \u00b7 Architecture: Shapley-Weighted Hypergraph Message Passing",
    [("Base propagation: e^(l+1) = \u03c3( D^-1/2 A D^-1/2 e^(l) )", 0),
     ("Shapley-weighted: e_j^(l+1) = \u03c3( W^(l) e_j^(l) + \u03a3_{k\u2208N(j)} w_jk e_k^(l) )", 0, True),
     ("Normalised neighbourhood weights: w_jk = \u03c6\u0302_jk / \u03a3_{k'\u2208N(j)} \u03c6\u0302_jk'", 0, True),
     ("Clipped + exponentially smoothed before normalisation (stabilises sparse regimes).", 1),
     ("Attention gate: a_ui = \u03c3( W_a[ e_u, e_i, l_i ] ); intermediate score y_ui = (1 + a_ui) \u27e8e_u, e_i\u27e9.", 0),
     ("Context-aware score: f(u,i,c) = y_ui + \u03bb_c \u27e8g(c_ui), e_cui\u27e9.", 0)],
    "The architecture is the decisive move. Propagation is standard hypergraph message passing, but the "
    "messages are weighted by normalised Shapley coefficients, so the model is told not only who is "
    "connected to whom but how much each coalition is worth. An attention gate interpolates between "
    "Shapley-weighted and uniform propagation as a stabiliser, and the final score is context-aware.",
    subtitle="Relational structure filtered through cooperative importance", num=N(),
    layout="image_right", image=os.path.join(ASSETS, "dyhucog_workflow.png"))

content_slide(
    "C3 \u00b7 Multi-Objective Training",
    [("L = L_rec + \u03bb_div L_div + \u03bb_ctx L_ctx + \u03bb_reg L_reg", 0, True),
     ("L_rec \u2013 Bayesian Personalised Ranking (pairwise, implicit feedback).", 1),
     ("L_div \u2013 Intra-List Diversity regulariser: penalises redundant ranked lists.", 1),
     ("L_ctx \u2013 Context alignment: match context embedding to context-node representation.", 1),
     ("L_reg \u2013 L2 weight decay.", 1),
     ("", 0),
     ("The learning objective and the coalition value are aligned: DyHuCoG trains to optimise the same balance that later determines attribution.", 1, True)],
    "The composite loss combines BPR ranking with explicit diversity and context regularisation plus "
    "weight decay. The conceptual point is alignment: the model is trained to optimise the same balance "
    "that later defines cooperative attribution, so the explanatory mechanism and the predictive "
    "mechanism are no longer separate.",
    subtitle="Accuracy, diversity, context trained together", num=N())

table_slide(
    "C3 \u00b7 Main Results (MovieLens-1M & Amazon-Book)",
    ["Dataset", "Model", "NDCG@20", "Recall@20", "Coverage", "Diversity"],
    [["MovieLens-1M", "HPCF", "0.2528", "0.2098", "0.342", "0.461"],
     ["MovieLens-1M", "DyHuCoG", "0.2775", "0.2362", "0.397", "0.516"],
     ["Amazon-Book", "HPCF", "0.0270", "0.0359", "0.259", "0.535"],
     ["Amazon-Book", "DyHuCoG", "0.0306", "0.0417", "0.336", "0.602"]],
    "The headline result. On MovieLens-1M, DyHuCoG beats the strongest baseline HPCF by +9.77% in "
    "NDCG@20 and +12.58% in Recall@20, and improves coverage and diversity. On the much sparser "
    "Amazon-Book, the gains are larger still: +13.33% NDCG@20 and +16.16% Recall@20, coverage up "
    "29.7%. The larger relative gain on the sparser benchmark supports the claim that Shapley-guided "
    "weighting helps most where interaction data are weak.",
    subtitle="vs the strongest baseline HPCF", num=N())

result_card_slide(
    "C3 \u00b7 Relative Gains over HPCF",
    [("+9.77%", "NDCG@20", "MovieLens-1M"),
     ("+12.58%", "Recall@20", "MovieLens-1M"),
     ("+13.33%", "NDCG@20", "Amazon-Book"),
     ("+16.16%", "Recall@20", "Amazon-Book")],
    "The most important message: DyHuCoG improves ranking accuracy AND coverage AND diversity "
    "simultaneously, with the largest relative gains on the sparser dataset. This is the evidence for "
    "the accuracy-diversity trade-off not being structurally fixed.",
    subtitle="Accuracy + coverage + diversity together", num=N())

content_slide(
    "C3 \u00b7 Coverage & Intra-List Diversity",
    [("MovieLens-1M:", 0, True),
     ("Coverage 0.342 \u2192 0.397 (+16.1%); Intra-List Diversity 0.461 \u2192 0.516 (+11.9%).", 1),
     ("Amazon-Book:", 0, True),
     ("Coverage 0.259 \u2192 0.336 (+29.7%); ILD 0.535 \u2192 0.602 (+12.5%).", 1),
     ("", 0),
     ("These indicate reduced filter-bubble effect and greater opportunity for discovery \u2014 while NDCG and Recall also improve, so accuracy is not traded off for diversity.", 1)],
    "Coverage and intra-list diversity both improve. This means more of the catalogue is surfaced "
    "across users, and the items within a list are less redundant. Because ranking quality also "
    "improves, we are not trading accuracy for diversity \u2014 the cooperative utility lets the model learn "
    "a balance where relevance, coverage, and diversity improve together.",
    subtitle="Both system- and list-level diversity improve", num=N(),
    layout="image_right", image=os.path.join(ASSETS, "dyhucog_covdiv.png"))

table_slide(
    "C3 \u00b7 Ablation Study (Component-wise)",
    ["Variant", "ML-1M NDCG@20", "% Drop", "Amazon NDCG@20", "% Drop"],
    [["Full DyHuCoG", "0.2775", "\u2013", "0.0306", "\u2013"],
     ["w/o Shapley Value", "0.2647", "4.6%", "0.0287", "6.1%"],
     ["w/o Hypergraph", "0.2586", "6.8%", "0.0279", "8.9%"],
     ["w/o Attention", "0.2678", "3.5%", "0.0295", "3.5%"],
     ["w/o Context", "0.2547", "8.2%", "0.0272", "11.0%"],
     ["w/o Diversity", "0.2614", "5.8%", "0.0288", "5.8%"]],
    "Every component contributes. Removing context causes the largest loss, because context provides "
    "the representational substrate on which Shapley weighting operates. Removing Shapley weighting "
    "and reverting to uniform propagation degrades performance noticeably, which directly supports the "
    "argument that marginal-contribution estimation is not decorative. Removing hierarchy-style "
    "structure (hypergraph) also matters.",
    subtitle="Each block is essential", num=N(), highlight_rows=[1])

content_slide(
    "C3 \u00b7 Computational Efficiency & Shapley Convergence",
    [("Training: DyHuCoG ~2000 s vs HPCF ~1125 s on MovieLens-1M (\u2248 1.78\u00d7).", 0),
     ("Inference: 1.84 ms/query (ML-1M), 8.52 ms (Amazon) \u2014 suitable for real-time deployment.", 0),
     ("Memory: 4.4 GB vs 4.1 GB (ML-1M); 17.9 vs 16.8 GB (Amazon).", 0),
     ("Per-epoch cost: O((L+1)md) + O((M/f)m).", 0),
     ("Shapley convergence: M=50 \u2192 MSE 1.4\u00d710\u207b\u2075, ~99% accuracy; higher M gives diminishing returns (M=100 \u2192 MSE 3.5\u00d710\u207b\u2076).", 0, True)],
    "The cost is bounded. Training is about 1.78 times HPCF, inference stays under a couple of "
    "milliseconds, and memory is only modestly higher. The Monte Carlo budget at M = 50 gives ~99% "
    "accuracy at acceptable cost, and the runtime overhead is proportionate to the performance and "
    "interpretability gains.",
    subtitle="The attribution cost is proportionate", num=N())

content_slide(
    "C3 \u00b7 Statistical Validation (MovieLens-1M)",
    [("Paired t-tests on per-user NDCG@20 (n = 6,040 users; df = 6,039).", 0),
     ("DyHuCoG outperforms every baseline with extremely small p-values after Holm\u2013Bonferroni correction.", 0),
     ("vs HPCF: t = 46.38, Cohen\u2019s d_z = 1.3345, p = 1.81\u00d710\u207b\u00b2\u2077\u2070.", 1, True),
     ("Wilcoxon signed-rank test also significant (p < 0.001) as a non-parametric robustness check.", 1),
     ("", 0),
     ("Effect sizes are large \u2014 the improvements are substantively meaningful, not merely statistically visible on a large user base.", 1)],
    "The improvements are statistically significant and large. The paired tests against HPCF give a t "
    "of 46.38 with Cohen's d_z of 1.33, all surviving Holm-Bonferroni. Effect sizes are large, which "
    "matters because large user counts can make tiny differences appear significant. The fully "
    "tabulated paired tests apply to MovieLens-1M; Amazon-Book results are reported descriptively.",
    subtitle="Significant AND substantively meaningful", num=N())

content_slide(
    "C3 \u00b7 Cold-Start Robustness & Interpretability",
    [("Cold-start (5 or fewer training interactions): NDCG@20 \u2248 0.061 (user) and 0.057 (item), improving over HPCF by ~10%.", 0),
     ("Cross-dataset improvement on the same protocol: MovieLens +9.9%, Amazon +14.8%, Yelp2018 +11.8%.", 0),
     ("", 0),
     ("Interpretability: a SHAP waterfall decomposes a single recommendation into ranking, diversity, context, and preference contributions \u2014 structurally tied to the same components that drove training.", 0, True),
     ("Popularity bias: Shapley measures marginal utility, not raw frequency \u2014 weak but informative interactions retain influence even for less popular items.", 1)],
    "DyHuCoG also improves the regimes where recommenders are usually most brittle: cold-start users "
    "and items, and it generalises across datasets of very different sparsity. Interpretability is "
    "direct: a waterfall plot decomposes a recommendation into the same utility components used during "
    "training, so the explanation is structurally faithful. Because Shapley values measure marginal "
    "utility rather than raw frequency, they help de-concentrate popularity.",
    subtitle="Robust in sparse regimes, and interpretable", num=N(),
    layout="image_right", image=os.path.join(ASSETS, "dyhucog_waterfall.png"))

content_slide(
    "C3 \u00b7 Findings & Limitations",
    [("Achieved:", 0, True),
     ("Cooperative attribution used as an in-training signal \u2014 a stronger claim than the clustering chapters.", 1),
     ("The accuracy\u2013diversity trade-off is not structurally fixed; accuracy, coverage, and diversity improve together.", 1),
     ("Limitations:", 0, True),
     ("Measurable computational overhead; depends on availability of meaningful context.", 1),
     ("Monte Carlo Shapley could be improved by variance reduction.", 1),
     ("Ablation is component-wise, not a factorial interaction design.", 1),
     ("Baselines finalised through early 2026; later models (e.g. LLM-augmented recommenders) not audited.", 1)],
    "C3 answers RQ3 and RQ4: cooperative attribution can move into the learning dynamics of a "
    "recommender, and doing so lets the model balance accuracy, diversity, and context more effectively. "
    "Limits: overhead, dependence on context, no factorial ablation, and the baseline set is capped at "
    "early 2026. The study also evaluates explanation indirectly rather than through a large-scale user "
    "study.",
    subtitle="RQ3 & RQ4 answered", num=N())

# --- 9. SECTION: CONCLUSION ------------------------------------------------
section_slide("Conclusion & Perspectives", "Synthesis \u00b7 Limitations \u00b7 Future work",
    "Let me now bring everything together.")

table_slide(
    "Thesis Synthesis",
    ["Contribution", "Main idea", "Achievement", "Key finding"],
    [["C1", "Explain black-box clustering via Shapley", "PCA\u2013KMeans\u2013LightGBM\u2013TreeSHAP pipeline", "Faithful, chemistry-consistent cluster attribution (wine)"],
     ["C2", "Multi-level, large-scale clustering XAI", "Cross-level SHAP aggregation + Prop. 6.1", "Coherent under hierarchy; interprets differences, not inconsistency"],
     ["C3", "DyHuCoG hypergraph cooperative game", "Preference-aware Shapley as in-training signal", "Accuracy + coverage + diversity improve together (ML & Amazon)"]],
    "Three contributions, one thread. C1 makes hidden structure intelligible; C2 keeps it coherent "
    "under scale and hierarchy; C3 carries the same attribution logic inside the learning dynamics of a "
    "recommender. Together they support the claim that cooperative game theory is a shared attribution "
    "perspective for explanation, optimisation, and intervention.",
    subtitle="The common thread", num=N())

table_slide(
    "Published Research Papers",
    ["No.", "Title", "Venue", "Status"],
    [["I", "Shapley Values for Explaining the Black Box Nature of Machine Learning Model Clustering", "Procedia Computer Science 220, 806\u2013811", "Published, 2023"],
     ["II", "Game Theory Meets Explainable AI: An Enhanced Approach to Understanding Black Box Models Through Shapley Values", "IJACSA 16(7), 716\u2013725", "Published, 2025"],
     ["III", "DyHuCoG: A Dynamic Hypergraph Cooperative Game for Preference-aware Recommendation", "IJIES 19(2), 887\u2013902", "Published, 2026"]],
    "The thesis synthesises three peer-reviewed publications. Paper I is Chapter 5, Paper II is Chapter "
    "6, and Paper III is Chapter 7. The thesis adds the multi-level formalisation and the thesis-level "
    "synthesis on top of these papers.",
    subtitle="Three publications \u2192 Chapters 5\u20137", num=N())

content_slide(
    "Thesis-Level Limitations",
    [("Computational \u2013 exact Shapley computation is intractable; every contribution relies on approximation, surrogates, or restricted reporting.", 0),
     ("Methodological \u2013 clustering depends on surrogate fidelity; recommendation depends on stable approximate contributions and adequate context.", 0),
     ("Empirical \u2013 tabular clustering + benchmark recommendation; no multimodal, sequential, or online deployment; no dedicated human-subject actionability study.", 0),
     ("Claim scope \u2013 a coherent and productive perspective, not one fully unified mathematical framework eliminating all tension between explanation, optimisation, diversity, fairness, and governance.", 0)],
    "Honestly stated limits at the programme level. Computationally, we approximate everywhere. "
    "Methodologically, the clustering contributions depend on surrogate fidelity and the recommendation "
    "on stable approximate contributions. Empirically, the work is tabular and offline, with no "
    "human-subject actionability study. And the claim is a shared perspective, not one grand unified "
    "framework.",
    subtitle="Scope is deliberate", num=N())

content_slide(
    "Future Research Directions",
    [("Scalable cooperative attribution \u2013 lower-variance Shapley, learned proposal distributions, adaptive refresh policies.", 0),
     ("Online / streaming recommendation \u2013 truly incremental settings with continuously evolving interaction graphs and delayed feedback.", 0),
     ("Richer human-centred evaluation \u2013 do explanations measurably improve analyst judgement, user trust, intervention quality, or perceived fairness?", 0),
     ("Broader trustworthy-AI evaluation \u2013 exposure fairness, transparency requirements, governance-oriented auditing.", 0)],
    "Future work turns the limitations into an agenda. The most concrete is more scalable cooperative "
    "attribution, then extending to online and streaming settings, richer human-centred evaluation of "
    "actionability, and broader trustworthy-AI and fairness auditing.",
    subtitle="Where the work goes next", num=N())

content_slide(
    "Thesis Answer & Key Outcomes",
    [("Thesis answer: cooperative game theory can function as a shared methodological perspective for actionable explanation across clustering and recommendation.", 0, True),
     ("", 0),
     ("Key outcomes:", 0, True),
     ("Shapley attribution as a common formal language for feature, interaction, and context importance allocation.", 1),
     ("Faithful clustering explanation, hierarchical explanatory coherence, and contribution-aware recommendation learning.", 1),
     ("Explanation as method, not commentary \u2014 from post-hoc description to in-training guidance.", 1),
     ("Aligned with trustworthy-AI requirements (EU AI Act, OECD principles, GDPR).", 1)],
    "The thesis answer: cooperative game theory is a shared perspective for explanation, optimisation, "
    "and intervention. Key outcomes: Shapley attribution as a common formal language, faithful "
    "clustering explanation, hierarchical coherence, and contribution-aware recommendation learning. "
    "This moves explanation from commentary to method, aligning with the transparency and accountability "
    "requirements of emerging AI regulation.",
    subtitle="The central claim, stated", num=N())

# --- 10. Q&A + CLOSING ------------------------------------------------------
content_slide(
    "References",
    [("R1  EU AI Act (Reg. 2024/1689) - transparency and explainability duties for high-risk AI.", 0),
     ("R2  Shapley, L.S. (1953). A value for n-person games. Annals of Mathematics Studies 28.", 0),
     ("R3  Lundberg & Lee (2017). A Unified Approach to Interpreting Model Predictions (SHAP / TreeSHAP). NeurIPS.", 0),
     ("R4  Ke et al. (2017). LightGBM: A Highly Efficient Gradient Boosting Decision Tree. NeurIPS.", 0),
     ("R5  Ribeiro et al. (2016). Why Should I Trust You? (LIME). ACM SIGKDD.", 0),
     ("R6  He et al. (2020). LightGCN. SIGIR.   R7  Xia et al. (2021) HCCF; Yu et al. (2021) HPCF. SIGIR/AAAI.", 0),
     ("R8  He et al. (2017). Neural Collaborative Filtering (NCF). WWW.", 0),
     ("R9  Gramegna & Giudici (2022). Shapley-value regression for credit-risk clustering.", 0),
     ("Own: Louhichi (2023, 2025, 2026) - the three published papers behind C1, C2, C3 (see Published Research Papers).", 0, True)],
    "A small set of the most load-bearing references. Shapley (1953) is the cooperative-game foundation; "
    "Lundberg & Lee (2017) the TreeSHAP engine used in C1; He (2020) LightGCN and the hypergraph baselines "
    "HCCF/HPCF (2021) are the recommendation comparators; Gramegna & Giudici (2022) the SHAP-clustering "
    "comparator for C2.",
    subtitle="Key sources", num=N())

content_slide(
    "Questions & Discussion",
    [("Thank you for your attention.", 0, True),
     ("I welcome your questions and comments.", 0)],
    "Thank you very much for your attention. I am now happy to take your questions.",
    subtitle="Jury discussion", num=N())

title_slide(
    ["Cooperative Game Theory for Explainable AI in Recommendation Systems:",
     "A Shapley Framework for Actionable Insight"],
    ["Mouad LOUHICHI \u2014 PhD Viva",
     "Supervisor: Pr. Mohamed LAZAAR \u2014 ENSIAS, Mohammed V University, Rabat"],
    "Thank you once again, President and members of the jury. This concludes my presentation. I am "
    "ready for your questions."
)
Q += 1

# --- Save -------------------------------------------------------------------
prs.save(OUT)
print(f"Saved {len(prs.slides)} slides -> {OUT}")
