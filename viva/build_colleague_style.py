#!/usr/bin/env python3
"""
Build Mouad Louhichi's 40-minute PhD viva presentation in the EXACT visual
design and framework of the colleague's deck.

The colleague's chrome (header title bar, 7-square tab row with icons, ENSIAS
footer bar + logo, corner pager icons, subtitle chips, big section-header block,
title-slide crest/logos, 3-question card layout, 4-chip contribution layout)
lives as *per-slide* shapes.  We therefore reproduce it by CLONING real
colleague archetype slides (carrying pixel-identical chrome, theme fills and
Montserrat fonts) and rewriting the textual content to Mouad's thesis, deleting
the colleague's body shapes and inserting our own.

Output: viva/MOUAD_LOUHICHI_VIVA_40min.pptx
"""
import copy, hashlib, io, os, subprocess, sys
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

HERE = os.path.dirname(os.path.abspath(__file__))
OUT = os.path.join(HERE, "MOUAD_LOUHICHI_VIVA_40min.pptx")
SRC = os.path.join(HERE, "colleague_template.pptx")

# The colleague's original 80-slide deck is NOT committed under viva/ (it is a
# 13 MB duplicate of a file already tracked in git history). If the template is
# missing, recreate it from origin/main so the build stays reproducible without
# bloating the repository.
_GIT_COLLEAGUE = "example-phd-passes/Presentation1 (1) (1).pptx"
def _ensure_template():
    if os.path.exists(SRC):
        return
    try:
        data = subprocess.check_output(
            ["git", "show", f"origin/main:{_GIT_COLLEAGUE}"], stderr=subprocess.DEVNULL)
    except subprocess.CalledProcessError:
        print("ERROR: colleague_template.pptx missing and could not be rebuilt "
              "from origin/main. Run:\n"
              "  git show 'origin/main:example-phd-passes/Presentation1 (1) (1).pptx' "
              "> viva/colleague_template.pptx", file=sys.stderr)
        sys.exit(1)
    with open(SRC, "wb") as f:
        f.write(data)
    print(f"Rebuilt {SRC} ({len(data)} bytes) from {_GIT_COLLEAGUE}")

RT_IMAGE = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/image"

MONTS   = "Montserrat"
MONTS_M = "Montserrat Medium"
# Teal + charcoal palette -- replaces the colleague's theme accent1 blue.
#   ACCENT  #0E7C7B  primary teal  -> 5.01:1 on white, white-on-teal 5.01:1 (AA)
#   DEEP    #0A5F5E  deep teal     -> 7.47:1, used for large bars / table highlight
#   DARK    #22303C  charcoal      -> 13.50:1, headings
#   BODY    #2F3B49  body text     -> 11.40:1
ACCENT  = RGBColor(0x0E, 0x7C, 0x7B)
DEEP    = RGBColor(0x0A, 0x5F, 0x5E)
ACCT    = ACCENT                  # generated accents (chips, bullets, tables, strips)
DARK    = RGBColor(0x22, 0x30, 0x3C)   # charcoal headings
BODY    = RGBColor(0x2F, 0x3B, 0x49)   # body text
GREY    = RGBColor(0x6B, 0x72, 0x80)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
MINT    = RGBColor(0xE8, 0xF1, 0xF0)   # soft mint-grey panel
HDRBG   = RGBColor(0xE1, 0xEC, 0xEB)

# Authoritative md5-prefix(PNG) -> digit map for the colleague's corner pager icons
# (bottom-right, top > 6.9", absolute 1-based slide number). These are the SOLID
# digits (the top 7-tab row uses OUTLINED digits -- left untouched).
SOLID_DIGIT = {
    'ebf64fb6': '1', '2a787842': '2', '3513e8dc': '3', '6ce466f5': '4',
    '3844ae35': '5', 'dac0f9a2': '6', 'fd8131e6': '7', '51c5fe8c': '8',
    '7d2fba37': '9', 'fa3a5ada': '0',
}
OUTLINE_DIGIT = {  # top 7-tab row OUTLINE glyphs (inactive tabs)
    'b51c2b4c': '1', '403a7a78': '2', '381e8bf3': '3', '8ddc531b': '4',
    '6635d9f4': '5', '0b4905c7': '6', '7ded1579': '7',
}
# SOLID glyphs for the ACTIVE tab (visually distinct from the pager solid digits)
TAB_SOLID_DIGIT = SOLID_DIGIT  # same solid digit PNGs, per-digit
PAGER_TOP = 7.0            # pager pictures top (inches)
PAGER_SIZE = 0.42          # pager picture square edge (inches)
PAGER_TENS_X = 12.502      # tens digit picture left (inches)
PAGER_UNITS_X = 12.806     # units digit picture left (inches)

FOOTER_TEXT = ("PhD Viva \u2013 Mouad LOUHICHI \u2013 Cooperative Game Theory & "
               "Shapley for XAI in Recommendation Systems")

# ---------------------------------------------------------------------------
# generic shape helpers
# ---------------------------------------------------------------------------
def _find(shapes, name):
    for sh in shapes:
        if sh.name == name:
            return sh
    return None

def set_slide_notes(slide, text):
    if text:
        slide.notes_slide.notes_text_frame.text = text

def set_lines(slide, bound_shape, lines):
    """Replace paragraphs of a shape with `lines` while preserving its first
    run's formatting (font/size/color) from the colleague archetype."""
    sh = _find(slide.shapes, bound_shape)
    if sh is None or not sh.has_text_frame:
        return
    tf = sh.text_frame
    proto_rpr = None
    for p in tf._txBody.findall(qn('a:p')):
        for r in p.findall(qn('a:r')):
            rPr = r.find(qn('a:rPr'))
            if rPr is not None:
                proto_rpr = copy.deepcopy(rPr)
                break
        if proto_rpr is not None:
            break
    bodyPr = tf._txBody.find(qn('a:bodyPr'))
    for p in list(tf._txBody.findall(qn('a:p'))):
        tf._txBody.remove(p)
    if bodyPr is not None:
        bodyPr.set('wrap', 'square')
        for tag in ('a:normAutofit', 'a:spAutoFit'):
            e = bodyPr.find(qn(tag))
            if e is not None:
                bodyPr.remove(e)
    for idx, text in enumerate(lines):
        p = tf._txBody.makeelement(qn('a:p'), {})
        tf._txBody.append(p)
        if idx > 0:
            p.append(p.makeelement(qn('a:pPr'), {'algn': 'l'}))
        r = p.makeelement(qn('a:r'), {})
        p.append(r)
        if proto_rpr is not None:
            r.append(copy.deepcopy(proto_rpr))
        else:
            r.append(r.makeelement(qn('a:rPr'), {}))
        t = r.makeelement(qn('a:t'), {})
        t.text = text
        r.append(t)

def remove_shapes(slide, names):
    for nm in names:
        sh = _find(slide.shapes, nm)
        if sh is not None:
            sh._element.getparent().remove(sh._element)

def clone_slide(prs, source):
    dest = prs.slides.add_slide(source.slide_layout)
    for shp in list(dest.shapes):
        shp._element.getparent().remove(shp._element)
    for shp in source.shapes:
        dest.shapes._spTree.append(copy.deepcopy(shp._element))
    for shp in dest.shapes:
        if shp.shape_type == 13:  # picture
            for blip in shp._element.findall('.//' + qn('a:blip')):
                rid = blip.get(qn('r:embed'))
                if rid is None:
                    continue
                img_part = source.part.related_part(rid)
                new_rid = dest.part.relate_to(img_part, RT_IMAGE)
                blip.set(qn('r:embed'), new_rid)
    return dest


# ---------------------------------------------------------------------------
# corner pager (absolute 1-based slide number)
# ---------------------------------------------------------------------------
def _is_pager_picture(sh):
    """True for the bottom-right pager digit pictures (top>6.9", x in 12.4–12.95")."""
    if sh.shape_type != 13:
        return False
    return (sh.top / 914400 > 6.9 and 12.4 < sh.left / 914400 < 12.95)

def collect_digit_blobs(prs):
    """Return {digit_char: PNG blob} for every solid pager glyph in the source deck."""
    blobs = {}
    for slide in prs.slides:
        for sh in slide.shapes:
            if not _is_pager_picture(sh):
                continue
            h = hashlib.md5(sh.image.blob).hexdigest()[:8]
            d = SOLID_DIGIT.get(h)
            if d is not None and d not in blobs:
                blobs[d] = sh.image.blob
    return blobs

def remap_pager(prs, digit_blobs):
    """Rewrite every bottom-right pager to show the slide's true 1-based number."""
    for i, slide in enumerate(prs.slides):
        number = i + 1
        pics = [sh for sh in slide.shapes if _is_pager_picture(sh)]
        if not pics:
            continue
        # remove the frozen glyphs inherited from the archetype clone
        for sh in pics:
            sh._element.getparent().remove(sh._element)
        digits = [int(c) for c in str(number)]
        if len(digits) == 1:
            _add_pager_pic(slide, digits[0], PAGER_UNITS_X, digit_blobs)
        else:
            _add_pager_pic(slide, digits[0], PAGER_TENS_X, digit_blobs)
            _add_pager_pic(slide, digits[1], PAGER_UNITS_X, digit_blobs)

def _add_pager_pic(slide, digit, x_in, digit_blobs):
    blob = digit_blobs.get(str(digit))
    if blob is None:
        return
    pic = slide.shapes.add_picture(io.BytesIO(blob), Inches(x_in), Inches(PAGER_TOP),
                                   Inches(PAGER_SIZE), Inches(PAGER_SIZE))
    return pic


# ---------------------------------------------------------------------------
# top 7-tab navigation row (active-tab highlight)
# ---------------------------------------------------------------------------
def _is_tab_picture(sh):
    """A digit glyph in the top 7-tab row (top<0.7\", x in 0.2–4.2\")."""
    if sh.shape_type != 13:
        return False
    return (sh.top / 914400 < 0.7 and 0.2 < sh.left / 914400 < 4.2)

def _tab_rect_fill_ref(sh):
    """Return the <a:fillRef> element that controls a tab-square's colour."""
    for fr in sh._element.iter(qn('a:fillRef')):
        return fr
    return None

def collect_tab_glyphs(prs):
    """Return ({digit: solid_blob}, {digit: outline_blob}) for the tab row."""
    solid, outline = {}, {}
    for slide in prs.slides:
        for sh in slide.shapes:
            if not _is_tab_picture(sh):
                continue
            h = hashlib.md5(sh.image.blob).hexdigest()[:8]
            d = SOLID_DIGIT.get(h)
            if d is not None and d not in solid:
                solid[d] = sh.image.blob
            d = OUTLINE_DIGIT.get(h)
            if d is not None and d not in outline:
                outline[d] = sh.image.blob
    return solid, outline

def set_active_tab(slide, active, solid_glyphs, outline_glyphs):
    """Highlight tab `active` (1–7) in the 7-tab row and un-highlight the rest.

    The active tab is accent (fillRef accent1 -> teal) with a SOLID digit glyph;
    inactive tabs are white (lt1) with OUTLINE digit glyphs.  Tab order runs left→right
    by picture x-position, so the tab whose digit == position index is active.
    """
    pics = sorted([sh for sh in slide.shapes if _is_tab_picture(sh)],
                  key=lambda sh: sh.left)
    if not pics:
        return
    # The archetype uses x-offset picture names but keeps left→right order = 1..7.
    # Recompute the tab index each digit glyph should display by its slot (1..7).
    for slot, pic in enumerate(pics, start=1):
        digit = str(slot)
        want_solid = (slot == active)
        blob = (solid_glyphs if want_solid else outline_glyphs).get(digit)
        if blob is None:
            continue
        _swap_pic_image(slide, pic, blob)
    # colour the tab squares: active accent, inactive white
    rects = [sh for sh in slide.shapes
             if sh.name.startswith('Rectangle')
             and sh.top / 914400 < 0.7 and 0.2 < sh.left / 914400 < 4.2
             and 0.3 < sh.width / 914400 < 0.7]
    # match rectangle slots to picture slots by x position
    for slot, pic in enumerate(pics, start=1):
        # find the rect whose left is closest to this pic's left
        best = None; bestd = 1e9
        for rect in rects:
            d = abs(rect.left - pic.left)
            if d < bestd:
                bestd = d; best = rect
        if best is None:
            continue
        col = 'accent1' if slot == active else 'lt1'
        _set_fill_ref_color(best, col)

def _swap_pic_image(slide, pic, blob):
    """Point `pic` at a new image part created from `blob`."""
    img_part, rId = slide.part.get_or_add_image_part(io.BytesIO(blob))
    for blip in pic._element.findall('.//' + qn('a:blip')):
        blip.set(qn('r:embed'), rId)

def _set_fill_ref_color(shape, scheme_val):
    """Set a shape's style/fillRef schemeClr to `scheme_val`."""
    fr = _tab_rect_fill_ref(shape)
    if fr is None:
        return
    sc = fr.find(qn('a:schemeClr'))
    if sc is not None:
        sc.set('val', scheme_val)


# ---------------------------------------------------------------------------
# per-section sub-navigation chips  (active chip filled, others outlined)
# ---------------------------------------------------------------------------
def _style_ref(shape, name):
    """Return the <a:style> child element named `name` (fillRef/fontRef/lnRef)."""
    for sty in shape._element.iter(qn('p:style')):
        for ch in sty:
            if ch.tag.endswith('}%s' % name):
                return ch
    return None

def set_chip_active(shape, active):
    """Fill the chip as ACTIVE (accent fill, white text) or INACTIVE (white fill,
    dark text), mirroring the colleague's chip styling in the archetype."""
    fr = _style_ref(shape, 'fillRef')
    if fr is not None:
        sc = fr.find(qn('a:schemeClr'))
        if sc is not None:
            sc.set('val', 'accent1' if active else 'lt1')
    fnt = _style_ref(shape, 'fontRef')
    if fnt is not None:
        sc = fnt.find(qn('a:schemeClr'))
        if sc is not None:
            sc.set('val', 'lt1' if active else 'dk1')

def apply_chip_states(slide, chip_names, active_idx):
    """Highlight the chip at `active_idx` (0-based) in `chip_names`; outline the rest."""
    for i, nm in enumerate(chip_names):
        sh = _find(slide.shapes, nm)
        if sh is None:
            continue
        set_chip_active(sh, i == active_idx)


# ---------------------------------------------------------------------------
# Teal + charcoal accent conversion
# ---------------------------------------------------------------------------
def recolor_accent_teal(slide):
    """Convert the colleague's accent1-blue chrome fills to teal.

    Colleague shapes express the accent via <p:style><a:fillRef><a:schemeClr
    val="accent1"/></a:fillRef>.  Rewrite that reference to a literal #0E7C7B, and
    also swap any a:solidFill srgbClr 4472C4 runs to teal.  Text colour (fontRef)
    is left untouched -- it resolves to white for high contrast.
    """
    for sh in slide.shapes:
        el = sh._element
        # <a:fillRef><a:schemeClr val="accent1"/> -> srgbClr 0E7C7B
        for fr in el.iter(qn('a:fillRef')):
            sc = fr.find(qn('a:schemeClr'))
            if sc is not None and sc.get('val') == 'accent1':
                srgb = sc.makeelement(qn('a:srgbClr'), {'val': '0E7C7B'})
                fr.replace(sc, srgb)
        # explicit solidFill 4472C4 -> teal
        for c in el.iter(qn('a:srgbClr')):
            if c.get('val') == '4472C4':
                c.set('val', '0E7C7B')
        # vestigial green default run-colour (endParaRPr) -> teal, for a consistent palette
        for c in el.iter(qn('a:srgbClr')):
            if c.get('val') == '98ECB7':
                c.set('val', '0E7C7B')
        # text that sits directly on a teal fill must be white for high contrast
        if _effective_fill_is_accent(sh):
            for rp in el.iter(qn('a:rPr')):
                sf = rp.find(qn('a:solidFill'))
                if sf is None:
                    continue
                sc = sf.find(qn('a:schemeClr'))
                if sc is not None and sc.get('val') in ('tx1', 'dk1'):
                    sc.set('val', 'bg1')
                s = sf.find(qn('a:srgbClr'))
                if s is not None and s.get('val') in ('000000', '1F2A44', '333F50'):
                    s.set('val', 'FFFFFF')

def _effective_fill_is_accent(sh):
    spPr = sh._element.find(qn('p:spPr'))
    if spPr is not None:
        for c in spPr:
            if c.tag == qn('a:solidFill'):
                s = c.find(qn('a:srgbClr'))
                if s is not None:
                    return s.get('val').upper() in ('0E7C7B', '0A5F5E')
                return False
            if c.tag == qn('a:noFill'):
                return False
    for sty in sh._element.iter(qn('p:style')):
        fr = sty.find(qn('a:fillRef'))
        if fr is not None:
            s = fr.find(qn('a:srgbClr'))
            if s is not None:
                return s.get('val').upper() in ('0E7C7B', '0A5F5E')
            sc = fr.find(qn('a:schemeClr'))
            if sc is not None:
                return sc.get('val') == 'accent1'
    return False

# ---------------------------------------------------------------------------
# fresh content shapes (colleague visual language)
# ---------------------------------------------------------------------------
def add_textbox(slide, x, y, w, h, lines, size=16, color=BODY, bold=False,
                font=MONTS, align=PP_ALIGN.LEFT, spacing=1.0, space_after=6,
                anchor=MSO_ANCHOR.TOP, italic=False):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    for i, text in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.line_spacing = spacing; p.space_after = Pt(space_after)
        r = p.add_run(); r.text = text
        r.font.name = font; r.font.size = Pt(size); r.font.bold = bold
        r.font.italic = italic; r.font.color.rgb = color
    return box

def add_bullets(slide, x, y, w, h, items, size=17, color=BODY, accent=None,
                space=9, anchor=MSO_ANCHOR.TOP, spacing=1.05):
    """items: str => bullet; (text, level); (text, level, bold). Wingdings § bullet."""
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    if accent is None:
        accent = ACCT
    for i, it in enumerate(items):
        if isinstance(it, str):
            text, level, bold = it, 0, False
        elif len(it) == 2:
            text, level = it
            bold = False
        else:
            text, level, bold = it
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(space if level == 0 else max(space - 3, 2))
        p.line_spacing = spacing
        if level == 0:
            r0 = p.add_run(); r0.text = "\u00a7  "
            r0.font.name = "Wingdings"; r0.font.size = Pt(size); r0.font.color.rgb = accent
            r = p.add_run(); r.text = text
            r.font.name = MONTS_M if bold else MONTS; r.font.size = Pt(size)
            r.font.bold = bold; r.font.color.rgb = color
        else:
            p.level = 1
            r0 = p.add_run(); r0.text = "\u2013  "
            r0.font.name = MONTS; r0.font.size = Pt(size-1); r0.font.color.rgb = accent
            r = p.add_run(); r.text = text
            r.font.name = MONTS; r.font.size = Pt(size-1); r.font.bold = bold
            r.font.color.rgb = color
    return box

def add_figure(slide, path, x, y, w=None, h=None):
    """Embed a generated PNG figure on the slide at (x,y) with optional size."""
    if not os.path.exists(path):
        return None
    from PIL import Image as _Im
    with _Im.open(path) as im:
        iw, ih = im.size
    if w is None and h is None:
        raise ValueError("add_figure needs at least one of w/h")
    if w is None:
        w = h * iw / ih
    if h is None:
        h = w * ih / iw
    pic = slide.shapes.add_picture(path, Inches(x), Inches(y), Inches(w), Inches(h))
    return pic

def add_table(slide, x, y, w, h, headers, rows, font_size=12, col_ratios=None,
              header_bg=None, body_color=BODY, hl_rows=None, hl_bg=None,
              hl_color=WHITE, row_h=None):
    if header_bg is None: header_bg = ACCT
    if hl_bg is None: hl_bg = DEEP
    ncols = len(headers); nrows = len(rows) + 1
    shape = slide.shapes.add_table(nrows, ncols, Inches(x), Inches(y), Inches(w), Inches(h))
    table = shape.table
    table.first_row = False; table.horz_banding = False
    if col_ratios:
        total = sum(col_ratios)
        for j, rt in enumerate(col_ratios):
            table.columns[j].width = Emu(int(w * (rt / total) * 914400))
    row_ht = row_h if row_h else (h / nrows)
    for r in range(nrows):
        table.rows[r].height = Emu(int(row_ht * 914400))
    for j, htxt in enumerate(headers):
        c = table.cell(0, j); c.text = htxt
        p = c.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        run = p.runs[0]; run.font.size = Pt(font_size); run.font.bold = True
        run.font.name = MONTS_M; run.font.color.rgb = WHITE
        c.fill.solid(); c.fill.fore_color.rgb = header_bg
        c.vertical_anchor = MSO_ANCHOR.MIDDLE
        c.margin_left = Inches(0.06); c.margin_right = Inches(0.06)
        c.margin_top = Inches(0.02); c.margin_bottom = Inches(0.02)
    for i, row in enumerate(rows, start=1):
        is_hl = hl_rows is not None and (i - 1) in hl_rows
        for j, val in enumerate(row):
            c = table.cell(i, j); c.text = str(val)
            p = c.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
            run = p.runs[0]; run.font.size = Pt(font_size); run.font.name = MONTS
            run.font.color.rgb = hl_color if is_hl else body_color
            run.font.bold = is_hl
            if is_hl:
                c.fill.solid(); c.fill.fore_color.rgb = hl_bg
            elif i % 2 == 0:
                c.fill.solid(); c.fill.fore_color.rgb = RGBColor(0xF0, 0xF4, 0xFB)
            else:
                c.fill.solid(); c.fill.fore_color.rgb = WHITE
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            c.margin_left = Inches(0.06); c.margin_right = Inches(0.06)
            c.margin_top = Inches(0.02); c.margin_bottom = Inches(0.02)
    return table

# ---------------------------------------------------------------------------
# Archetype-aware builders
# ---------------------------------------------------------------------------
def build_content(prs, arch, title, chips, body, note, chip_names,
                  remove_content, active=None, tab_solid=None, tab_outline=None,
                  chip_active=None):
    s = clone_slide(prs, arch)
    set_lines(s, 'Rectangle 5', [FOOTER_TEXT])
    # header title bar (name differs between archetypes)
    header = 'Rectangle 27' if _find(s.shapes, 'Rectangle 27') else 'Rectangle 30'
    set_lines(s, header, [title])
    remove_shapes(s, list(remove_content))
    for i, nm in enumerate(chip_names):
        if i < len(chips) and chips[i]:
            set_lines(s, nm, [chips[i]])
        else:
            remove_shapes(s, [nm])
    if chip_active is not None:
        apply_chip_states(s, chip_names, chip_active)
    builders = body if isinstance(body, list) else [body]
    for fn in builders:
        fn(s)
    if active is not None and tab_solid and tab_outline:
        set_active_tab(s, active, tab_solid, tab_outline)
    set_slide_notes(s, note)
    return s

def build_section(prs, source_sec, title_small=None, note=None):
    s = clone_slide(prs, source_sec)
    set_lines(s, 'Rectangle 5', [FOOTER_TEXT])
    # big block title
    set_lines(s, 'Rectangle 27', [title_small] if title_small else
              [sh.text_frame.text for sh in source_sec.shapes if sh.name == 'Rectangle 27'])
    if note:
        set_slide_notes(s, note)
    return s

def build_outline(prs, cards):
    s = clone_slide(prs, OUTLINE_ARCH)
    set_lines(s, 'Rectangle 5', [FOOTER_TEXT])
    mapping = {'intro': 'Rectangle 27', 'context': 'Rectangle 41',
               'protocols': 'Rectangle 42', 'contrib': 'Rectangle 43',
               'conclusion': 'Rectangle 44'}
    for key, lines in cards.items():
        set_lines(s, mapping[key], lines)
    return s

# ---------------------------------------------------------------------------
def main():
    global OUTLINE_ARCH
    _ensure_template()
    prs = Presentation(SRC)

    TITLE_ARCH    = prs.slides[0]
    OUTLINE_ARCH  = prs.slides[1]
    SEC_INTRO     = prs.slides[2]
    CARDS_ARCH    = prs.slides[3]
    SEC_CONTEXT   = prs.slides[7]
    SEC_PROTO     = prs.slides[16]
    CONTENT_ARCH  = prs.slides[20]   # header + 7 tab squares + 3 chips
    SEC_CI        = prs.slides[27]
    CONTRIB_ARCH  = prs.slides[40]   # header + 4 chips (Objectives/Meth/Results/Findings)
    SEC_CII       = prs.slides[39]
    SEC_CIII      = prs.slides[52]
    SEC_CONCL     = prs.slides[68]
    QUES_ARCH     = prs.slides[75]
    END_ARCH      = prs.slides[76]

    orig_ids = list(prs.slides._sldIdLst)
    CONTENT_CHIPS = ('Rectangle 29', 'Rectangle 30', 'Rectangle 33')
    CONTENT_REMOVE = ('TextBox 36',)
    CONTRIB_CHIPS = ('Rectangle 33', 'Rectangle 36', 'Rectangle 45', 'Rectangle 46')
    CONTRIB_REMOVE = ('TextBox 31', 'Table 7', 'TextBox 41')

    # 7-tab navigation: 1 Introduction, 2 Context, 3 Protocols,
    # 4 Contribution I, 5 Contribution II, 6 Contribution III, 7 Conclusion
    TAB_SOLID, TAB_OUTLINE = collect_tab_glyphs(prs)
    CURRENT_ACTIVE = {'n': 1}

    def content(title, chips, body, note, active=None, chip=None):
        a = CURRENT_ACTIVE['n'] if active is None else active
        return build_content(prs, CONTENT_ARCH, title, chips, body, note,
                             CONTENT_CHIPS, CONTENT_REMOVE, active=a,
                             tab_solid=TAB_SOLID, tab_outline=TAB_OUTLINE,
                             chip_active=chip)
    def contrib(title, chips, body, note, active=None, chip=None):
        a = CURRENT_ACTIVE['n'] if active is None else active
        return build_content(prs, CONTRIB_ARCH, title, chips, body, note,
                             CONTRIB_CHIPS, CONTRIB_REMOVE, active=a,
                             tab_solid=TAB_SOLID, tab_outline=TAB_OUTLINE,
                             chip_active=chip)

    def section(source, title, note, active=None):
        if active is not None:
            CURRENT_ACTIVE['n'] = active
        return build_section(prs, source, title, note)

    # ---- 1. TITLE SLIDE ----
    s = clone_slide(prs, TITLE_ARCH)
    set_lines(s, 'Rectangle 8', [
        "Cooperative Game Theory for Explainable AI in Recommendation Systems:",
        "A Shapley Framework for Actionable Insight"])
    set_lines(s, 'TextBox 6', ["PhD Viva Presented by: Mouad LOUHICHI"])
    # jury table: roles + known supervisor, others to be completed
    tbl_frame = _find(s.shapes, 'Table 2')
    tbl = tbl_frame.table if tbl_frame is not None and tbl_frame.has_table else None
    if tbl is not None:
        data = [
            ["President",  "[President to complete]", "", "[Institution]"],
            ["Supervisor", "Pr. Mohamed LAZAAR", "PES", "ENSIAS, Mohammed V University, Rabat"],
            ["Reviewer",   "[Reviewer to complete]", "", "[Institution]"],
            ["Reviewer",   "[Reviewer to complete]", "", "[Institution]"],
            ["Reviewer",   "[Reviewer to complete]", "", "[Institution]"],
            ["Examiner",   "[Examiner to complete]", "", "[Institution]"],
            ["Examiner",   "[Examiner to complete]", "", "[Institution]"],
            ["Examiner",   "[Examiner to complete]", "", "[Institution]"],
            ["Guest",      "[Guest to complete]", "", "[Institution]"],
        ]
        for i, row in enumerate(data):
            for j, val in enumerate(row):
                cell = tbl.cell(i, j); cell.text = val
                p = cell.text_frame.paragraphs[0]
                p.alignment = PP_ALIGN.LEFT if j in (0, 1) else PP_ALIGN.CENTER
                for run in p.runs:
                    run.font.name = MONTS; run.font.size = Pt(10)
                    run.font.color.rgb = DARK if i == 1 else BODY
                    run.font.bold = (j == 0)
    set_slide_notes(s, "Good morning. Thank you, President and Professors, for the time "
        "you are giving me to present and discuss my thesis. My name is Mouad Louhichi. "
        "My thesis is titled Cooperative Game Theory for Explainable AI in Recommendation "
        "Systems: A Shapley Framework for Actionable Insight, supervised by Professor "
        "Mohamed Lazaar. It advances one idea: that Shapley-value attribution is not just a "
        "post-hoc explanation, but a single formal mechanism that can explain black-box "
        "clustering, stay coherent under hierarchical scale, and finally act as an in-training "
        "signal inside a recommender.")
    orig_ids  # keep reference; clone order preserved

    # ---- 2. OUTLINE ----
    build_outline(prs, {
        'intro': ["Introduction", "Motivation", "Actionable insight", "Why XAI"],
        'context': ["Context & Problematic", "Limitations", "Research gap"],
        'protocols': ["Protocols", "Datasets", "Baselines", "Metrics", "Hardware"],
        'contrib': ["Contributions", "C1 Explainable clustering", "C2 Multi-level XAI", "C3 DyHuCoG"],
        'conclusion': ["Conclusion & Perspectives", "Summary", "Limitations", "Future work"]})
    set_slide_notes(prs.slides[-1], "I will walk you through five parts: introduction, context "
        "and problematic, the shared experimental protocol, then my three contributions, and "
        "finally the conclusion and perspectives. Each contribution follows the same structure: "
        "objectives, methodology, results, and findings. So you always know where we are.")

    # ---- 3. SECTION: Introduction ----
    section(SEC_INTRO, "Introduction", active=1, note="Let us begin with the introduction. I will spend a few minutes establishing why "
            "explainability is not a nice-to-have but a first-class requirement for recommender "
            "systems, and I will define carefully what I mean by an actionable insight.")
    # 4. Motivation cards (3-question layout)
    s = clone_slide(prs, CARDS_ARCH)
    set_lines(s, 'Rectangle 5', [FOOTER_TEXT])
    if _find(s.shapes, 'Rectangle 27'):
        set_lines(s, 'Rectangle 27', ["Motivation: Three Questions"])
    set_lines(s, 'Rectangle 29', ["Motivation"])
    if _find(s.shapes, 'Rectangle 28'):
        set_lines(s, 'Rectangle 28', ["Actionable Insight"])
    apply_chip_states(s, ('Rectangle 29', 'Rectangle 28'), 0)
    set_lines(s, 'Rectangle 47', ["Ubiquity", "How do opaque AI systems shape what billions of users see, buy and watch every day?"])
    set_lines(s, 'Rectangle 48', ["The Black Box", "Why do state-of-the-art recommenders and clustering pipelines remain opaque to users and designers?"])
    set_lines(s, 'Rectangle 49', ["Toward Trust", "How can transparency be built as part of the model, instead of being bolted on afterwards?"])
    remove_shapes(s, ['TextBox 50'])
    set_active_tab(s, 1, TAB_SOLID, TAB_OUTLINE)
    set_slide_notes(s, "Three questions frame the whole work. First, ubiquity: opaque systems "
        "mediate what billions of people see, buy and watch every day \u2014 not only on streaming "
        "platforms but in e-commerce, news feeds and search. Second, the black box: even strong "
        "recommenders and clustering pipelines are hard to interrogate \u2014 we observe their "
        "outputs without understanding why a particular item was chosen. Third, trust: transparency "
        "should be built into the modelling logic, not bolted on after the fact. The core tension of "
        "the thesis is that as models gain expressive power, they lose the transparency needed for "
        "trustworthy deployment. Throughout this work I hold accuracy and interpretability as "
        "objectives to be reconciled, not traded against one another.")
    # add a small "core tension" strip
    add_textbox(s, 0.3, 4.2, 12.7, 0.8,
                ["Core tension: as models gain expressive power, they lose the transparency needed for trustworthy deployment."],
                size=16, color=ACCT, bold=True, font=MONTS_M, align=PP_ALIGN.CENTER)

    # 5. Actionable insight
    content("Actionable Insight \u2014 the Definition",
        ["Motivation", "Actionable Insight", "Research Context"],
        [lambda sl: add_bullets(sl, 0.4, 1.9, 12.5, 4.6, [
            ("Definition 1.1 (Actionable insight)", 0, True),
            ("An explanation is actionable when it identifies at least one modifiable factor whose change is associated with a specifiable change in model output\u2026", 1),
            ("\u2026 and that factor is expressible in the semantic vocabulary of the task domain.", 1),
            ("Accessibility means the domain's own terms:", 0),
            ("a physicochemical variable (wine), a pollution indicator (air quality), or a preference signal (recommendation).", 1),
            ("not an opaque latent code.", 1),
            ("Why: an explanation that identifies a modifiable driver supports intervention, not merely description.", 0, True),
        ])],
        "We frame the whole thesis around actionable insight. An explanation is actionable when it "
        "identifies at least one modifiable factor whose change is associated with a specifiable change "
        "in model output, and when that factor is expressible in the semantic vocabulary of the task "
        "domain \u2014 a physicochemical variable for wine, a pollution indicator for air quality, a "
        "preference signal for recommendation. The crucial constraint is that it must not be an opaque "
        "latent code. This distinction separates explanations that let a designer intervene and act from "
        "those that merely describe what happened. By this definition, a credible explanation is "
        "evaluated by whether it supports a downstream decision, not by whether it looks plausible.",
        chip=1)

    # 6. Research context
    content("Research Context",
        ["Motivation", "Actionable Insight", "Research Context"],
        [lambda sl: add_bullets(sl, 0.4, 1.9, 12.5, 4.6, [
            "Recommenders evolved from similarity filters to complex representation-learning systems on sparse, high-dimensional, dynamic data.",
            ("\u2192 MF \u2192 neural CF \u2192 graph CNNs \u2192 hypergraph: each step improved ranking but intensified the interpretability deficit.", 1),
            ("Why the deficit matters:", 0),
            ("Undermines user trust.", 1),
            ("Constrains debugging and scientific learning.", 1),
            ("Collides with regulatory expectations (EU AI Act, OECD principles, GDPR).", 1),
        ])],
        "The context is the progression from simple to hypergraph recommenders. Matrix factorisation "
        "showed that latent factors could capture preferences, but those factors were immediately "
        "uninterpretable. Graph and graph-neural-network models improved ranking by exploiting "
        "connectivity, yet replaced an opaque latent code with an opaque message-passing mechanism. "
        "Finally, hypergraph models added higher-order user-item-context relations, but typically assume "
        "that every message contributes uniformly. Each step raised expressiveness while lowering "
        "transparency. The deficit matters for trust, debugging and regulation \u2014 the EU AI Act, the "
        "OECD principles and the GDPR all put a premium on meaningful explanation. That is why a "
        "principled attribution mechanism is needed, and why I will argue it should be part of the "
        "modelling logic rather than an afterthought.",
        chip=2)

    # ---- 7. SECTION: Context & Problematic ----
    section(SEC_CONTEXT, "Context & Problematic", active=2, note="Now let us look more precisely at the problem this thesis addresses. I will quickly "
            "survey the paradigm landscape, then identify the structuring limitations and the five "
            "research questions, and finally map the three contributions to those questions.")
    # 8. Recommendation paradigms
    content("Recommendation & Clustering Paradigms",
        ["Paradigms", "Problematic", "Contributions"],
        [lambda sl: add_bullets(sl, 0.4, 1.9, 12.5, 4.6, [
            ("Collaborative filtering \u2013 users who behaved similarly will value similar items (user-/item-based).", 0),
            ("Content-based \u2013 recommends items sharing attributes with a user profile.", 0),
            ("Hybrid \u2013 combines collaborative and content signals.", 0),
            ("Matrix factorisation \u2013 R \u2248 PQ\u1d40: compact but opaque latent factors.", 0),
            ("Graph-based \u2013 interaction graph with neighbourhood propagation (LightGCN, hypergraph).", 0),
        ])],
        "Quick orientation across the paradigms we build on. Collaborative filtering recommends on the "
        "principle that similar users will value similar items, computed from the user side or the item "
        "side. Content-based filtering recommends items sharing attributes with a profile, while hybrid "
        "methods combine both. Matrix factorisation factorises the interaction matrix into latent "
        "factors \u2014 compact and effective, but immediately opaque. Graph-based methods, including "
        "LightGCN and hypergraph extensions, propagate information over an interaction structure and "
        "capture multi-hop or higher-order relations. Each strengthens the modelling, but each "
        "complicates interpretation in a specific way: matrix factorisation made latent dimensions "
        "opaque; graph models kept importance implicit; hypergraph models added higher-order relations "
        "but assumed uniform message importance. That uniformity assumption is one of the things I will "
        "challenge.",
        chip=0)

    # 9. Limitations of classical models
    content("Limitations of Classical Recommenders & Unsupervised Models",
        ["Paradigms", "Problematic", "Contributions"],
        [lambda sl: add_bullets(sl, 0.4, 1.9, 12.5, 4.6, [
            "Data sparsity & scalability \u2013 the user\u2013item matrix is overwhelmingly empty.",
            "Cold-start \u2013 new users and items are structurally disadvantaged.",
            "Popularity bias & lack of diversity \u2013 exposure begets interaction, begets exposure.",
            ("Absence of interpretability \u2013 the most fundamental limit.", 0, True),
            ("For clustering: methods favour local OR global explanation, not both.", 1),
            ("They struggle to scale, and explanations rarely stay coherent across resolutions.", 1),
        ])],
        "Four classical limitations. Data sparsity and scalability are the most cited: the user-item "
        "matrix is overwhelmingly empty, so learning is starved of signal. Cold-start is the structural "
        "consequence, because a new user or item has no history to learn from. Popularity bias and the "
        "resulting lack of diversity create a filter-bubble loop, where exposure begets interaction. And "
        "the most fundamental limitation, and the one this thesis targets, is the absence of "
        "interpretability. For clustering specifically the situation is even harder: methods privilege a "
        "local explanation or a global one but not both, they struggle to scale, and their explanations "
        "rarely remain coherent across resolutions. That combination \u2014 no faithful local and global "
        "explanation, poor scaling, and incoherence across levels \u2014 is precisely the gap I work on.",
        chip=1)

    # 10. Problem statement
    content("Three Structuring Limitations (Problem Statement)",
        ["Paradigms", "Problematic", "Contributions"],
        [lambda sl: add_bullets(sl, 0.4, 1.9, 12.5, 4.6, [
            ("1. Lack of explainability \u2013 complex models remain hard to interpret faithfully and actionably.", 0, True),
            ("2. Difficulty of scaling \u2013 local explanations do not transfer naturally to hierarchical structures or large datasets.", 0, True),
            ("3. Weak integration into learning \u2013 most explanations stay post-hoc and do not shape model dynamics or the accuracy-diversity-context trade-off.", 0, True),
            ("Thesis gap: the literature still lacks a single cooperative-attribution framework that explains clustering faithfully, stays coherent under hierarchy, and then operates as an in-training signal in recommendation.", 1, True),
        ])],
        "Three structuring problems. First, lack of explainability: complex models remain hard to "
        "interpret faithfully and actionably. Second, difficulty of scaling: local explanations do not "
        "transfer naturally to hierarchical structures or large datasets, so a method that works on a "
        "toy partition may collapse on hundreds of thousands of records with nested structure. Third, "
        "weak integration into learning: most explanations are post-hoc and do not shape model "
        "dynamics, nor do they participate in the accuracy-diversity-context trade-off the model is "
        "trying to satisfy. The thesis gap follows directly: the literature still lacks a single "
        "cooperative-attribution framework that explains clustering faithfully, stays coherent under "
        "hierarchy, and then operates as an in-training signal in recommendation. My claim is that "
        "Shapley-value attribution can be that single framework.",
        chip=1)

    # 11. Research questions
    content("Research Questions (RQ1\u2013RQ5) and Overall Aim",
        ["Paradigms", "Problematic", "Contributions"],
        [lambda sl: add_bullets(sl, 0.4, 1.9, 12.5, 4.6, [
            ("Aim: develop, justify and evaluate a cooperative-game-theoretic perspective for XAI in clustering and recommendation.", 0, True),
            ("RQ1 \u00b7 How can Shapley values explain black-box clustering faithfully at instance and cluster level?", 0),
            ("RQ2 \u00b7 How can this extend to large-scale, hierarchical clustering without losing tractability or consistency?", 0),
            ("RQ3 \u00b7 Can cooperative attribution move beyond post-hoc and enter the learning dynamics of graph recommenders?", 0),
            ("RQ4 \u00b7 Can a recommender jointly optimise ranking accuracy, context and diversity when importance is estimated by a cooperative-game utility?", 0),
            ("RQ5 \u00b7 What emerges when clustering explanation and recommendation learning are two stages of one cooperative-game perspective?", 0),
        ])],
        "The overall aim is to develop, justify and evaluate a cooperative-game-theoretic perspective "
        "for explainable AI in clustering and recommendation, using Shapley attribution as both an "
        "explanatory mechanism and an in-training signal. Five research questions form the spine. RQ1 "
        "asks how Shapley values can explain black-box clustering faithfully at instance and cluster "
        "level. RQ2 asks how that extends to large-scale hierarchical clustering without losing "
        "tractability or consistency. RQ3 asks whether cooperative attribution can move beyond post-hoc "
        "analysis into the learning dynamics of graph recommenders. RQ4 asks whether a recommender can "
        "jointly optimise ranking accuracy, context and diversity when importance is estimated by a "
        "cooperative-game utility. RQ5 is the thesis-level question: what emerges when clustering "
        "explanation and recommendation learning are read as two stages of one shared perspective.",
        chip=1)

    # 12. Three contributions overview
    content("The Three Contributions",
        ["Paradigms", "Problematic", "Contributions"],
        [lambda sl: add_bullets(sl, 0.4, 1.9, 12.5, 4.6, [
            ("C1 \u2013 Explainable black-box clustering: PCA\u2013K-Means\u2013LightGBM\u2013TreeSHAP pipeline.", 0, True),
            ("\u2192 Wine Quality: faithful instance- and cluster-level feature attribution.", 1),
            ("C2 \u2013 Enhanced multi-level XAI for large-scale clustering with cross-level SHAP aggregation.", 0, True),
            ("\u2192 Beijing Air Quality: hierarchical attribution consistency under scale.", 1),
            ("C3 \u2013 DyHuCoG: Dynamic Hypergraph Cooperative Game for preference-aware recommendation.", 0, True),
            ("\u2192 MovieLens-1M & Amazon-Book: preference-aware Monte Carlo Shapley as an in-training signal.", 1),
            ("Thesis claim: cooperative game theory functions as a shared attribution perspective for explanation, optimisation and intervention.", 2, True),
        ])],
        "Three contributions, one thread. C1 establishes Shapley-based explanation for black-box "
        "clustering through a PCA-K-Means-LightGBM-TreeSHAP pipeline, and validates it on the wine "
        "quality dataset. C2 scales it to hierarchy and large data via multi-level clustering with "
        "cross-level SHAP aggregation, validated on the Beijing air-quality dataset. C3 is the "
        "strongest claim: DyHuCoG replaces post-hoc attribution with an in-training signal inside a "
        "hypergraph recommender, validated on MovieLens-1M and Amazon-Book. Each contribution answers a "
        "distinct research question, but they are deliberately designed to be read together as a "
        "cumulative argument. The thesis claim is that cooperative game theory functions as a shared "
        "attribution perspective for explanation, optimisation and intervention \u2014 not that these "
        "are three unrelated papers.",
        chip=2)

    # ---- 13. SECTION: Protocols ----
    section(SEC_PROTO, "Experimental Protocol", active=3, note="Before the contributions, let me briefly cover the shared experimental setup. The "
            "purpose of this section is to establish that my comparisons are fair: the datasets, the "
            "splitting and preprocessing, the baselines, the evaluation metrics, and the hardware "
            "environment.")
    # 14. Datasets
    content("Datasets Used Throughout the Thesis",
        ["Datasets", "Metrics", "Hardware"],
        [lambda sl: add_table(sl, 0.4, 1.9, 12.6, 3.9,
             ["Dataset", "Scale", "Type", "Role"],
             [["Wine Quality (vinho verde)", "4,898 \u00d7 11", "Tabular, numeric", "C1 \u2013 single-level clustering"],
              ["Beijing Multi-Site Air Quality", "383,585 \u00d7 11", "Tabular, pollutant + meteorology", "C2 \u2013 multi-level clustering"],
              ["MovieLens-1M", "6,040 u / 3,706 i / 1.0M int", "Implicit feedback (0.0447)", "C3 \u2013 DyHuCoG"],
              ["Amazon-Book", "52,643 u / 91,599 i / 3.0M int", "Implicit feedback (0.0006)", "C3 \u2013 DyHuCoG"]],
             col_ratios=[4, 2.6, 2.8, 3.4])],
        "Two clustering datasets and two recommendation datasets. For clustering I deliberately picked "
        "datasets with semantically interpretable features, because the whole point of the method is to "
        "return attribution to the original variables. Wine Quality is a small, dense, chemically "
        "correlated dataset of almost five thousand samples with eleven features. Beijing Multi-Site "
        "Air Quality is a large, noisy, temporally and meteorologically variable dataset of over "
        "three hundred and eighty thousand hourly records. For recommendation I used benchmark-standard "
        "datasets with established baselines: MovieLens-1M, with roughly a million interactions and a "
        "density of 0.0447, and Amazon-Book, which is far sparser at 0.0006. The sparsity contrast is "
        "deliberate, because it stress-tests whether a Shapley-guided model helps most precisely when "
        "interaction data are weak.",
        chip=0)

    # 15. Splitting & preprocessing
    content("Data Splitting & Preprocessing",
        ["Datasets", "Metrics", "Hardware"],
        [lambda sl: add_bullets(sl, 0.4, 1.9, 12.5, 4.6, [
            ("Clustering: five-fold cross-validation for surrogate/attribution stability.", 0),
            ("Recommendation: user-level, temporal split \u2014 70% train / 10% val / 20% test.", 0),
            ("Leave-one-out: the latest test positive per user is the target, ranked against negatives.", 0),
            ("Implicit conversion: MovieLens-1M ratings > 3 treated as positive.", 0),
            ("Popularity-aware negative sampling: q(i) \u221d f_i^\u03b7 for harder contrasts.", 0),
            ("Reproducibility: seeds {42,43,44,45,46}; early-stopping patience 20.", 0),
        ])],
        "Splitting is designed against leakage, because recommendation metrics are especially sensitive "
        "to how you split. For clustering I use five-fold cross-validation to test the stability of the "
        "surrogate and the attribution. For recommendation I use a user-level, temporal holdout: "
        "interactions sorted by time per user, split seventy percent train, ten percent validation, "
        "twenty percent test. Evaluation uses leave-one-out, where the latest test positive per user is "
        "the target, ranked against sampled negatives. Because MovieLens-1M contains explicit ratings, I "
        "convert ratings greater than three into positive implicit feedback. Negatives are drawn from a "
        "popularity-aware distribution to produce harder contrasts. Results are reported across five "
        "seeds, with early stopping on validation NDCG at twenty.",
        chip=0)

    # 16. Baselines & metrics
    content("Baselines & Evaluation Metrics",
        ["Datasets", "Metrics", "Hardware"],
        [lambda sl: add_bullets(sl, 0.4, 1.9, 12.5, 4.6, [
            ("Clustering benchmark: LIME-based surrogate explanation pipeline.", 0, True),
            ("Recommendation benchmarks: MF, NCF, LightGCN, RecDCL, HCCF, HPCF (strongest reference).", 0, True),
            ("Ranking: Precision@K, Recall@K, NDCG@20 (principal).", 0),
            ("System diversity: Catalogue Coverage = |\u22c3 R_u| / |I|.", 0),
            ("List diversity: Intra-List Diversity (ILD), built into the coalition utility.", 0),
            ("Clustering quality: Silhouette, Davies\u2013Bouldin.", 0),
        ])],
        "Baselines span classical, neural, graph and hypergraph methods, precisely so that I can "
        "isolate the contribution of cooperative attribution rather than a favourable model choice. For "
        "recommendation I compare against matrix factorisation, NCF, LightGCN, RecDCL, HCCF and HPCF, "
        "treating HPCF as the strongest reference. For clustering interpretability I compare against a "
        "LIME-based surrogate pipeline. Ranking is measured with Precision at K, Recall at K and NDCG at "
        "twenty as the principal measure. On the diversity side I measure system-level catalogue "
        "coverage, and list-level intra-list diversity, defined as the average pairwise dissimilarity "
        "inside a ranked list. I want to emphasise that ILD is not decorative: it is deliberately built "
        "into the DyHuCoG coalition utility. Clustering quality is measured with Silhouette and "
        "Davies-Bouldin.",
        chip=1)

    # 17. Hardware & software
    content("Hardware & Software",
        ["Datasets", "Metrics", "Hardware"],
        [lambda sl: add_bullets(sl, 0.4, 1.9, 12.5, 4.6, [
            ("CPU: Intel Core i9-14900K, 24 cores \u2014 clustering, preprocessing, data loading.", 0),
            ("GPU: NVIDIA GeForce RTX 4090, 24 GB \u2014 DyHuCoG training & inference.", 0),
            ("RAM: 48 GB; Storage: 2 TB SSD.", 0),
            ("Python 3.8; scikit-learn, LightGBM, SHAP, PyTorch 2.0.1, NumPy/SciPy/pandas.", 0),
            ("Altair for interactive SHAP visualisation; metrics at K \u2208 {5, 10, 20}.", 0),
        ])],
        "The hardware matters mainly because it explains some runtime figures I will quote later. All "
        "clustering, preprocessing and data loading run on an Intel Core i9-14900K with twenty-four "
        "cores. DyHuCoG training and inference run on an NVIDIA RTX 4090 with twenty-four gigabytes, "
        "alongside forty-eight gigabytes of RAM and a two-terabyte SSD. The software stack is Python "
        "3.8, with scikit-learn for PCA and K-Means, LightGBM for the surrogate, SHAP for TreeSHAP, "
        "and PyTorch 2.0.1 for the recommendation model. I use Altair for interactive SHAP "
        "visualisation and report metrics at K equal to five, ten and twenty. Everything stays within "
        "ordinary academic compute; nothing here needs industrial infrastructure.",
        chip=2)

    # ---- 18. SECTION: Contribution I ----
    section(SEC_CI, "Contribution I \u2014 Explainable Black-Box Clustering", active=4, note="Let us move to the first contribution: explaining black-box clustering with Shapley "
            "values. This contribution answers RQ1, and I will present it in four steps: objectives, "
            "methodology, results, findings.")
    # C1 slides (4-chip contribution layout)
    contrib("Contribution I \u2014 Explainable Black-Box Clustering",
        ["Objectives", "Methodology", "Results", "Findings"],
        [lambda sl: add_bullets(sl, 0.4, 1.9, 12.5, 4.6, [
            ("Gap: Shapley explanation is dominant in supervised tasks, but clustering remains comparatively under-explained.", 0, True),
            ("Existing clustering-interpretability methods privilege local or global explanation, not both.", 1),
            ("They often fail to scale or preserve coherence across clusters.", 1),
            ("Objectives: build a pipeline yielding cluster-level explanation while preserving feature-level attribution.", 0, True),
            ("Preserve the semantics of the original feature space; justify why Shapley is better than LIME.", 1),
        ])],
        "Why start with clustering? Here the model creates its own structure, so cluster meaning must "
        "be inferred after the fact \u2014 which makes clustering the hardest and most natural test bed "
        "for attribution. The gap is that explainable clustering is fragmented: methods favour a local "
        "or a global explanation, rarely both; they often fail to scale; and they rarely preserve "
        "coherence across clusters. And because Shapley explanation is well established in supervised "
        "problems, its near absence from unsupervised clustering is striking. Our objectives are "
        "threefold: build a pipeline yielding a cluster-level explanation while preserving "
        "feature-level attribution; preserve the semantics of the original feature space rather than a "
        "transformed latent space; and justify why Shapley is the right concept rather than an ad-hoc "
        "surrogate such as LIME.", chip=0)

    contrib("Contribution I \u2014 Explainable Black-Box Clustering",
        ["Objectives", "Methodology", "Results", "Findings"],
        [lambda sl: add_bullets(sl, 0.4, 1.9, 12.5, 4.6, [
            ("RQ1 \u00b7 How can Shapley values explain black-box clustering faithfully at instance and cluster level?", 0, True),
            ("Objective O1 \u00b7 Build a pipeline that yields cluster-level explanation while keeping feature-level attribution.", 1),
            ("Objective O2 \u00b7 Preserve the semantics of the original feature space, not a reduced latent space.", 1),
            ("Objective O3 \u00b7 Justify Shapley over an ad-hoc surrogate such as LIME.", 1),
        ])],
        "RQ1 asks how Shapley values can explain black-box clustering faithfully at both instance "
        "and cluster level, and it decomposes into three objectives. O1 is a pipeline that produces a "
        "cluster-level explanation while preserving feature-level attribution, so the explanation is "
        "coherent rather than a single summary number. O2 requires that attribution stays in the "
        "original feature space \u2014 density, pH, acidity, sulfur dioxide, alcohol \u2014 because that "
        "is what makes it actionable for a domain expert. O3 is a justificatory objective: we must show "
        "why Shapley is the right allocation rule rather than an ad-hoc surrogate such as LIME. The next "
        "methodology slide shows how these three objectives are met.", chip=0)

    contrib("Contribution I \u2014 Explainable Black-Box Clustering",
        ["Objectives", "Methodology", "Results", "Findings"],
        [lambda sl: add_bullets(sl, 0.4, 1.9, 12.5, 4.6, [
            ("Player set N = F \u2014 each feature is a player.", 0, True),
            ("Value function v(S) = Silhouette( KMeans(X_S, k*) ) \u2014 how well data cluster using only features in S.", 0),
            ("A feature's Shapley value = its expected marginal contribution to clustering quality over all coalition orders.", 0),
            ("Why Silhouette: bounded, normalised, semantically intuitive.", 0),
            ("Direct evaluation for every coalition is intractable \u2192 we need a bridge.", 0, True),
        ])],
        "We frame clustering as a cooperative game where features are the players. The value function "
        "measures how well the data cluster when we use only the features in a given coalition: "
        "concretely, the value of coalition S is the Silhouette of K-Means run on the feature subset "
        "X_S, with a fixed k. I choose Silhouette because it is bounded, normalised and semantically "
        "intuitive; Davies-Bouldin or Calinski-Harabasz would also be defensible. A feature then "
        "receives a high Shapley value when its presence consistently improves separation across all "
        "coalitions, in expectation. The problem is that evaluating Silhouette for every subset is "
        "combinatorial and infeasible even for a modest number of features, which is exactly why we "
        "need a bridge between the unsupervised partition and a tractable attribution method.", chip=1)

    contrib("Contribution I \u2014 Explainable Black-Box Clustering",
        ["Objectives", "Methodology", "Results", "Findings"],
        [lambda sl: add_bullets(sl, 0.4, 1.9, 12.5, 4.6, [
            ("Once K-Means produces cluster labels, train a LightGBM multiclass surrogate to predict them from original features.", 0, True),
            ("Apply TreeSHAP to the surrogate \u2014 fast, exact tree-based attribution in the original semantic feature space.", 0),
            ("Direct TreeSHAP on K-Means is impossible (it explains tree models, not centroids).", 1),
            ("The surrogate preserves the chemistry/pollution vocabulary that makes the analysis actionable.", 1),
            ("Validity condition: surrogate fidelity is high (macro-F1 \u2248 0.82).", 1, True),
        ])],
        "The bridge is the heart of the method. TreeSHAP explains tree models; it cannot explain K-Means "
        "centroids. And explaining the PCA representation would move attribution away from the "
        "interpretable variables we care about. So we convert an unsupervised partition into a "
        "supervised task: once K-Means produces cluster labels, we train a LightGBM multiclass "
        "classifier to predict them from the original features, then apply TreeSHAP to that surrogate. "
        "This keeps attribution in the original semantic feature space, which is what makes the "
        "explanation actionable. The key validity condition is surrogate fidelity: if the surrogate does "
        "not reproduce the partition well, the attribution is not faithful. We require a macro-F1 of "
        "around 0.82 as the floor.", chip=1)

    contrib("Contribution I \u2014 Explainable Black-Box Clustering",
        ["Objectives", "Methodology", "Results", "Findings"],
        [lambda sl: add_bullets(sl, 0.4, 1.9, 12.5, 4.6, [
            ("Stage 1 \u2013 PCA: stabilise geometry + visual diagnostic. NOT the explanatory space.", 0),
            ("Stage 2 \u2013 K-Means++ with multi-criteria k selection (elbow, Silhouette, Davies\u2013Bouldin).", 0),
            ("Stage 3 \u2013 LightGBM surrogate trained on original features to predict cluster labels.", 0),
            ("Stage 4 \u2013 TreeSHAP attribution in the original feature space.", 0),
            ("Stage 5 \u2013 Aggregate into global importance, cluster-specific profiles, and local force plots.", 0),
            ("Complexity: dominated by PCA and repeated K-Means; TreeSHAP scales with tree count, not exponentially in features.", 1),
        ]),
        lambda sl: add_figure(sl, os.path.join(HERE, '_figs', 'c1_pipeline.png'), 1.6, 3.75, w=10.1)],
        "The pipeline has five stages. Stage one is feature standardisation. Stage two is PCA, which I "
        "want to stress is used only as a computational and visual diagnostic; it is deliberately not the "
        "explanatory space, because explaining principal components would defeat the purpose of "
        "actionable attribution. Stage three is K-Means-plus-plus with multi-criteria k selection, "
        "drawing on the elbow method, Silhouette and Davies-Bouldin. Stage four is the LightGBM "
        "surrogate trained on the original features to predict cluster labels. Stage five is TreeSHAP "
        "attribution in the original feature space, aggregated into global importance, cluster-specific "
        "profiles, and local force plots. Complexity is dominated by PCA and repeated K-Means; TreeSHAP "
        "scales with tree count and depth rather than exponentially in features, which is what makes the "
        "approach tractable.", chip=1)

    contrib("Contribution I \u2014 Explainable Black-Box Clustering",
        ["Objectives", "Methodology", "Results", "Findings"],
        [lambda sl: add_bullets(sl, 0.4, 1.9, 12.5, 4.6, [
            ("Multi-criteria evaluation across k \u2208 {2..10} using elbow, Silhouette, Davies\u2013Bouldin.", 0),
            ("We select k* = 3 \u2014 even though it is NOT geometrically optimal:", 0, True),
            ("k = 2: Silhouette 0.214, Davies\u2013Bouldin 1.775 (better separation).", 1),
            ("k = 3: Silhouette 0.144, Davies\u2013Bouldin 2.097 (weaker separation).", 1),
            ("Why: three clusters give a semantically richer oenological partition \u2192 more actionable.", 1, True),
            ("Note: the higher Silhouette \u2248 0.63 belongs to Beijing (C2), not this wine partition.", 1),
        ])],
        "This is an honest and important point. We evaluate k from two to ten, and deliberately select "
        "k equal to three even though it is not geometrically optimal. With k equal to two we get a "
        "stronger Silhouette of 0.214 and a Davies-Bouldin of 1.775; with k equal to three the Silhouette "
        "drops to 0.144 and the Davies-Bouldin rises to 2.097. On raw geometry, two is the better "
        "partition. We choose three nonetheless, because three clusters support three distinct, "
        "chemically meaningful oenological narratives, far more actionable than a dichotomous split. I "
        "select the partition a domain expert would find useful, not the one that maximises a separation "
        "index. And to avoid a common confusion: the 0.63 Silhouette belongs to the Beijing dataset in "
        "Contribution II, not this one.", chip=2)

    contrib("Contribution I \u2014 Explainable Black-Box Clustering",
        ["Objectives", "Methodology", "Results", "Findings"],
        [lambda sl: add_bullets(sl, 0.4, 1.9, 12.5, 4.6, [
            ("Global SHAP ranking (high \u2192 low):", 0, True),
            ("density \u2192 pH \u2192 fixed acidity \u2192 sulfur-dioxide \u2192 alcohol", 1),
            ("Dominant drivers relate to structure, preservation, and sensory balance.", 1),
            ("This is NOT an arbitrary classifier artefact \u2014 it recovers a chemically interpretable hierarchy.", 1, True),
        ])],
        "The global SHAP ranking is dominated by density, followed by pH, fixed acidity, "
        "sulfur-dioxide-related variables and alcohol. These are precisely the variables that an "
        "oenologist would point to as governing structure, preservation and sensory balance. The "
        "important point is that this is not an arbitrary classifier artefact. Because the surrogate was "
        "trained to reproduce the partition from the original variables, and because TreeSHAP attributes "
        "in that original space, the recovered hierarchy corresponds to a chemically interpretable "
        "structure. This is the strongest evidence that the pipeline is faithful: it recovers domain "
        "knowledge without being told to. It also illustrates the value of keeping attribution in the "
        "semantic feature space rather than in a reduced or latent space.", chip=2)

    contrib("Contribution I \u2014 Explainable Black-Box Clustering",
        ["Objectives", "Methodology", "Results", "Findings"],
        [lambda sl: add_bullets(sl, 0.4, 1.9, 12.5, 4.6, [
            ("Three clusters show distinct explanatory signatures.", 0, True),
            ("Cluster 0 \u2013 density + sulfur-dioxide-related variables.", 1),
            ("Cluster 1 \u2013 acidity and pH-related effects.", 1),
            ("Cluster 2 \u2013 a different balance of acidity, alcohol, and related chemical attributes.", 1),
            ("The same small set of variables recurs across clusters, with different relative weights within each.", 1),
        ])],
        "The cluster-specific profiles show that the solution is not only globally interpretable but "
        "internally differentiated. Each cluster exhibits a distinct signature. Cluster zero is driven "
        "primarily by density and sulfur-dioxide-related variables; cluster one by acidity and pH "
        "effects; cluster two by a different balance across acidity, alcohol and related chemical "
        "attributes. Crucially, the same small set of variables recurs across all three clusters, but "
        "with different relative weights within each. This is the actionable insight: a cluster is not "
        "just a label, it is a distinct, domain-meaningful combination of drivers. It also shows the "
        "method supports both a global reading, which variable matters most overall, and a local, "
        "per-cluster reading, where the practical value lies.", chip=2)

    contrib("Contribution I \u2014 Explainable Black-Box Clustering",
        ["Objectives", "Methodology", "Results", "Findings"],
        [lambda sl: add_table(sl, 0.4, 1.9, 12.5, 3.8,
             ["Criterion", "SHAP (cooperative)", "LIME (local surrogate)"],
             [["Basis", "Cooperative-game marginal contribution", "Local surrogate approximation"],
              ["Local / global", "Both", "Primarily local"],
              ["Theoretical guarantee", "Efficiency, symmetry, null player, additivity", "None equivalent"],
              ["Stability", "Higher when surrogate faithful", "Sensitive to perturbation design"],
              ["Cluster comparison", "Strong", "Limited"]],
             col_ratios=[3, 4.4, 4.4])],
        "This table summarises why we prefer Shapley over LIME. SHAP grounds attribution in a "
        "cooperative-game allocation rule and satisfies four axioms: efficiency, so explanatory mass is "
        "fully allocated; symmetry, so equal marginal contributors get equal credit; the null-player "
        "property, so zero marginal contributors get zero; and additivity, so explanations compose. "
        "LIME fits a local surrogate and has no equivalent guarantee; it is sensitive to perturbation "
        "design. SHAP supports both local and global explanation, whereas LIME is primarily local, and "
        "SHAP is much stronger for comparing clusters, which is central here. One caveat: in this "
        "surrogate pipeline, efficiency holds with respect to the LightGBM output, not directly to the "
        "Silhouette-based game. That is exactly why surrogate fidelity matters.", chip=2)

    contrib("Contribution I \u2014 Explainable Black-Box Clustering",
        ["Objectives", "Methodology", "Results", "Findings"],
        [lambda sl: add_bullets(sl, 0.4, 1.9, 12.5, 4.6, [
            ("Answers to the questions and objectives", 0, True),
            ("RQ1 answered in the affirmative \u2014 faithful, coherent cluster-level explanation from Shapley values.", 1),
            ("O1 met \u00b7 cluster-level explanation while preserving feature-level attribution.", 1),
            ("O2 met \u00b7 attribution returned to the original chemical variables.", 1),
            ("O3 met \u00b7 Shapley grounded in four axioms; LIME has no equivalent guarantee.", 1),
        ])],
        "To answer RQ1 directly: yes. Shapley values explain a black-box partition faithfully, provided "
        "the surrogate is high-fidelity, and coherently at cluster level. Each objective is met. O1 is met "
        "because the pipeline yields both a global reading \u2014 density, pH, acidity, sulfur dioxide, "
        "alcohol \u2014 and a per-cluster reading where the same drivers re-weigh differently. O2 is met "
        "because attribution is returned to the original chemical variables, which is what makes it "
        "actionable. O3 is met because Shapley satisfies efficiency, symmetry, the null-player property "
        "and additivity, whereas LIME has no equivalent guarantee. The honest caveat is that efficiency "
        "holds with respect to the surrogate output, not directly to the Silhouette-based game \u2014 "
        "which is why surrogate fidelity is the critical validity condition.", chip=2)

    contrib("Contribution I \u2014 Explainable Black-Box Clustering",
        ["Objectives", "Methodology", "Results", "Findings"],
        [lambda sl: add_bullets(sl, 0.4, 1.9, 12.5, 4.6, [
            ("Key findings", 0, True),
            ("Cluster-level explanation anchored to individual feature contributions.", 1),
            ("Explanations returned to the original chemical variables, not a latent space.", 1),
            ("Theoretically grounded case for Shapley over LIME (efficiency, symmetry, null player, additivity).", 1),
            ("Recovers an oenologically interpretable ranking \u2014 density, pH, acidity, sulfur dioxide, alcohol.", 1),
        ])],
        "C1 answers RQ1 in the affirmative. The key findings are fourfold. Shapley values explain the "
        "black-box partition at both cluster level and instance level, anchoring each cluster's meaning "
        "to individual feature contributions. They return attribution to the original chemical variables "
        "\u2014 density, pH, fixed acidity, sulfur dioxide, alcohol \u2014 rather than to a reduced or "
        "latent space, which is what makes the explanation actionable. Because the surrogate reproduces "
        "the partition from the original variables, the recovered ranking is chemically interpretable and "
        "faithful. And the approach is theoretically grounded: SHAP satisfies efficiency, symmetry, the "
        "null-player property and additivity, which LIME does not.", chip=3)

    contrib("Contribution I \u2014 Explainable Black-Box Clustering",
        ["Objectives", "Methodology", "Results", "Findings"],
        [lambda sl: add_bullets(sl, 0.4, 1.9, 12.5, 4.6, [
            ("Limitations", 0, True),
            ("Fidelity depends on the LightGBM surrogate \u2014 not a direct mechanism of the K-Means geometry.", 1),
            ("Confined to tabular data; no structured, graph, or image input.", 1),
            ("Single-level structure: cannot explain how importance reconfigures between a partition and its sub-partitions.", 1),
            ("The surrogate approximation compresses observation-level variation.", 1),
        ])],
        "Let me be clear about the scope. First, the explanation depends on the fidelity of the LightGBM "
        "surrogate; it is not a direct mechanism of the K-Means geometry. Second, the approach is confined "
        "to tabular data. Third, and most importantly for what follows, it is single-level: it cannot yet "
        "address hierarchical coherence, meaning it cannot explain how feature importance reconfigures "
        "between a partition and its sub-partitions. And the surrogate plus representative-instance "
        "reporting compresses some observation-level variation. These limits define the point of departure "
        "for Contribution II.", chip=3)

    contrib("Contribution I \u2014 Explainable Black-Box Clustering",
        ["Objectives", "Methodology", "Results", "Findings"],
        [lambda sl: add_bullets(sl, 0.4, 1.9, 12.5, 4.6, [
            ("Takeaways", 0, True),
            ("Shapley attribution is a single, principled lens for explaining an unsupervised partition.", 1),
            ("Keeping attribution in the original feature space is what makes it actionable.", 1),
            ("But real data are rarely single-level: broad regimes contain nested sub-groups.", 1, True),
            ("So the next question is whether this logic survives scale and hierarchy.", 0, True),
        ])],
        "The takeaway is that Shapley attribution works as a single, principled lens for explaining an "
        "unsupervised partition, and that keeping attribution in the original feature space is what makes "
        "it actionable. But real data are rarely single-level: broad regimes contain nested sub-groups, "
        "and a variable can be globally important yet locally uninformative. That flat limitation directly "
        "motivates Contribution II: can the explanation logic be scaled to multi-level, large-scale "
        "clustering without losing coherence?", chip=3)

    # ---- 19. SECTION: Contribution II ----
    section(SEC_CII, "Contribution II \u2014 Enhanced Multi-Level XAI for Large-Scale Clustering", active=5, note="This brings us to the second contribution: scaling the explanation logic to multi-level, "
            "large-scale clustering. It answers RQ2, and the central concern is not merely scale but "
            "multi-granularity.")
    contrib("Contribution II \u2014 Multi-Level XAI for Large-Scale Clustering",
        ["Objectives", "Methodology", "Results", "Findings"],
        [lambda sl: add_bullets(sl, 0.4, 1.9, 12.5, 4.6, [
            ("Gap: once clustering is multi-level, feature importance must stay interpretable within a cluster, across sub-clusters, and across the hierarchy.", 0, True),
            ("Large-scale data make exact explanation computationally burdensome.", 1),
            ("Flat explanation may be true yet incomplete \u2014 it cannot show how importance changes inside a cluster.", 1),
            ("Objectives: a genuinely multi-level workflow; a formal cross-level consistency argument (Prop. 6.1); validation on a structurally different large-scale dataset.", 0, True),
        ])],
        "C2 asks whether the C1 logic survives scale and hierarchy. Large real-world data contain "
        "structure at more than one granularity: broad regimes at the top, nested sub-groups within "
        "them. A variable can be globally important yet locally uninformative, or the reverse, so a flat "
        "explanation is true but incomplete \u2014 it cannot show how importance reconfigures as you "
        "zoom in. The chapter adds three things. First, a genuine multi-level workflow, not a rerun of "
        "the single-level pipeline. Second, a formal cross-level consistency argument, Proposition 6.1, "
        "so that differences across levels can be interpreted rather than dismissed as inconsistency. "
        "Third, validation on a structurally different large-scale dataset, Beijing air quality, testing "
        "whether the method generalises beyond the small, correlated wine dataset.", chip=0)

    contrib("Contribution II \u2014 Multi-Level XAI for Large-Scale Clustering",
        ["Objectives", "Methodology", "Results", "Findings"],
        [lambda sl: add_bullets(sl, 0.4, 1.9, 12.5, 4.6, [
            ("RQ2 \u00b7 How can this extend to large-scale, hierarchical clustering without losing tractability or consistency?", 0, True),
            ("Objective O1 \u00b7 A genuinely multi-level workflow, not a rerun of the single-level pipeline.", 1),
            ("Objective O2 \u00b7 A formal cross-level consistency argument (Proposition 6.1).", 1),
            ("Objective O3 \u00b7 Validation on a structurally different large-scale dataset.", 1),
        ])],
        "RQ2 asks whether the C1 logic survives scale and hierarchy without losing tractability or "
        "consistency, and it decomposes into three objectives. O1 is a genuinely multi-level workflow: "
        "coarse clustering at the top and nested sub-groups within, so that attribution stays interpretable "
        "inside a cluster, across sub-clusters, and across the hierarchy. O2 is a formal cross-level "
        "consistency argument, Proposition 6.1, so that differences across levels can be interpreted rather "
        "than dismissed as inconsistency. O3 is an empirical check on generality: validation on a "
        "structurally different large-scale dataset, Beijing air quality, rather than the small, "
        "chemically correlated wine dataset.", chip=0)

    contrib("Contribution II \u2014 Multi-Level XAI for Large-Scale Clustering",
        ["Objectives", "Methodology", "Results", "Findings"],
        [lambda sl: add_bullets(sl, 0.4, 1.9, 12.5, 4.6, [
            ("Recursive/nested: coarse clustering on the full dataset, then subdivide each cluster.", 0, True),
            ("For each level, train a level-specific surrogate and compute SHAP in the SAME original feature space.", 0),
            ("Cross-level aggregation is NOT a naive average \u2014 it respects cluster size and nesting structure.", 0, True),
            ("Parent-level attribution = an expectation over the explanatory structure of its descendants.", 1),
            ("The hierarchy is a pragmatic analytical device, not a claim of true ontological hierarchy.", 1),
        ]),
        lambda sl: add_figure(sl, os.path.join(HERE, '_figs', 'c2_hierarchy.png'), 1.6, 3.75, w=10.1)],
        "The multi-level architecture proceeds recursively: we learn a coarse clustering on the full "
        "dataset, then subdivide each cluster where appropriate, producing a nested structure. For each "
        "level we train a level-specific surrogate and compute SHAP values in the same original feature "
        "space. The cross-level aggregation is deliberately not a naive average: it respects cluster "
        "size and nesting structure. A parent-level attribution is an expectation over the explanatory "
        "structure of its descendants, so a coarse-level importance reflects the size-weighted "
        "contributions of the clusters it contains. One methodological point: the hierarchy here is a "
        "pragmatic analytical device. I am not claiming the data have a true, ontological hierarchy of "
        "levels; the nested structure is a computational and interpretive tool, and Proposition 6.1 is "
        "stated with that in mind.", chip=1)

    contrib("Contribution II \u2014 Multi-Level XAI for Large-Scale Clustering",
        ["Objectives", "Methodology", "Results", "Findings"],
        [lambda sl: add_bullets(sl, 0.4, 1.9, 12.5, 4.6, [
            ("Let \u03a6^(l,c)_j = E_{x~c}[|\u03c6_j^(l)(x)|] be the expected absolute SHAP importance of feature j at level l in cluster c.", 0),
            ("Let w_c' = |c'| / |c| be the relative size of child c' within parent c.", 0),
            ("For a strict nested hierarchy on a consistent feature space:", 0, True),
            ("\u03a6^(l,c)_j = \u03a3_{c'\u2208child(c)} w_c' \u00b7 \u03a6^(l+1,c')_j + \u03b5_j", 1, True),
            ("\u03b5_j is a residual from surrogate mismatch, vanishing under perfect fidelity.", 1),
            ("Derived via law of total expectation (children partition the parent).", 1),
            ("Does NOT imply explanations are identical across levels \u2014 it implies differences can be interpreted, not dismissed as inconsistency.", 1),
        ])],
        "Proposition 6.1 is the mathematical heart of the contribution. Let Phi at level l in cluster c "
        "for feature j be the expected absolute SHAP importance over examples in that cluster, and let "
        "w sub c-prime be the relative size of a child cluster within its parent. For a strict nested "
        "hierarchy on a consistent feature space, the parent's expected absolute importance equals the "
        "sum over its children of the child's relative size times the child's expected absolute "
        "importance, plus a residual epsilon that comes from surrogate mismatch and vanishes under "
        "perfect fidelity. The derivation uses the law of total expectation and the fact that children "
        "partition the parent. Crucially, it does not claim explanations are identical across levels; it "
        "says differences can be interpreted rather than dismissed as inconsistency. That makes a "
        "hierarchical explanation self-consistent and auditable.", chip=1)

    contrib("Contribution II \u2014 Multi-Level XAI for Large-Scale Clustering",
        ["Objectives", "Methodology", "Results", "Findings"],
        [lambda sl: add_bullets(sl, 0.4, 1.9, 12.5, 4.6, [
            ("Full dataset, k = 3 (strong convergence on multi-criteria evaluation).", 0, True),
            ("Silhouette \u2248 0.63 \u2014 materially stronger separation than wine.", 1),
            ("Davies\u2013Bouldin \u2248 0.55 \u2014 low between-cluster ambiguity.", 1),
            ("PCA projection (2 components) used only for visual inspection.", 1),
            ("Sensitivity: robust to modest variation in k, projection dim, surrogate depth; only low-ranked variables shift.", 1),
        ])],
        "On Beijing, the multi-criteria evaluation converges much more strongly than on wine, and this "
        "is worth emphasising. The full dataset clusters cleanly into k equal to three, with a "
        "Silhouette of about 0.63 and a Davies-Bouldin of about 0.55. Recall that this is where the 0.63 "
        "figure comes from \u2014 it belongs to Beijing, not to the wine partition, whose Silhouette "
        "is 0.144. The low Davies-Bouldin indicates little between-cluster ambiguity. The PCA "
        "projection, using two components, is used only for visual inspection. I also tested "
        "sensitivity: the conclusions are robust to modest variation in k, in the number of projection "
        "dimensions, and in the surrogate depth; only low-ranked variables shift. So the finding is not "
        "an artefact of a particular parameter choice.", chip=2)

    contrib("Contribution II \u2014 Multi-Level XAI for Large-Scale Clustering",
        ["Objectives", "Methodology", "Results", "Findings"],
        [lambda sl: add_bullets(sl, 0.4, 1.9, 12.5, 4.6, [
            ("Global SHAP ranking (high \u2192 low):", 0, True),
            ("temperature \u2192 dew point \u2192 pressure \u2192 CO \u2192 NO2 \u2192 PM10 \u2192 PM2.5", 1),
            ("It is NOT simply pollutant concentrations that matter \u2014 meteorological variables play a structurally central role.", 1),
            ("Temperature, dew point and pressure condition dispersion, trapping, and photochemical behaviour.", 1),
            ("This is the kind of insight flat descriptive summaries often fail to make explicit.", 1),
        ])],
        "The global ranking is analytically rich. The dominant features are temperature, dew point and "
        "pressure, followed by carbon monoxide, nitrogen dioxide, PM10 and PM2.5. It is striking that "
        "it is not simply the pollutant concentrations that dominate. Meteorological variables play a "
        "structurally central role, because temperature, dew point and pressure condition dispersion, "
        "trapping and photochemical behaviour. In other words, they determine the atmospheric conditions "
        "under which pollution accumulates. This is exactly the kind of insight that flat, descriptive "
        "summaries often fail to make explicit: a naive ranking of pollutant concentrations would "
        "overlook the fact that the weather is what sets the regime. It is a good illustration of why "
        "attribution in the original variable space is valuable.", chip=2)

    contrib("Contribution II \u2014 Multi-Level XAI for Large-Scale Clustering",
        ["Objectives", "Methodology", "Results", "Findings"],
        [lambda sl: add_bullets(sl, 0.4, 1.9, 12.5, 4.6, [
            ("Regime A \u2013 warm photochemical: ozone, temperature, dew point prominent (summer photochemical smog).", 0),
            ("Regime B \u2013 wintertime smog: CO, SO2, PM dominate; low wind speed suppresses dispersion.", 0),
            ("Regime C \u2013 comparatively clean air: favourable meteorology, weak pollutant pushes.", 0),
            ("The framework shows not only that these regimes exist, but which variable combinations define them.", 1, True),
        ])],
        "The force plots reveal three representative regimes, and this is where the interpretation "
        "becomes concrete. Regime A is a warm, photochemical regime, where ozone, temperature and dew "
        "point are prominent \u2014 this is characteristic of summer photochemical smog. Regime B is a "
        "wintertime smog regime, dominated by carbon monoxide, sulfur dioxide and particulate matter, "
        "with low wind speed suppressing dispersion. Regime C is a comparatively clean-air regime, "
        "associated with favourable meteorology and weak pollutant pushes. The interpretative value is "
        "not merely showing that these regimes exist \u2014 it is showing which combinations of "
        "variables define each one. That is the actionable insight: an air-quality analyst can see the "
        "specific drivers of a given pollution episode and act on them.", chip=2)

    contrib("Contribution II \u2014 Multi-Level XAI for Large-Scale Clustering",
        ["Objectives", "Methodology", "Results", "Findings"],
        [lambda sl: add_bullets(sl, 0.4, 1.9, 12.5, 4.6, [
            ("At the coarse level, temperature and dew point dominate \u2014 they differentiate broad atmospheric regimes.", 0),
            ("Within individual clusters, CO, SO2, PM10, wind speed, pressure, or ozone become more discriminative.", 0),
            ("This change is NOT contradictory \u2014 it is exactly what a multi-level explanation should reveal.", 0, True),
            ("Parent-level story = regime selection. Cluster-level story = variation within a regime.", 1),
            ("A variable can be globally important yet locally uninformative within a sub-cluster.", 1),
        ])],
        "This is the conceptual payoff of the chapter, and the reason a flat explanation is "
        "insufficient. At the coarse level, temperature and dew point dominate, because they "
        "differentiate broad atmospheric regimes. Within individual clusters, the discriminative "
        "variables shift to carbon monoxide, sulfur dioxide, PM10, wind speed, pressure or ozone. This "
        "change is not a contradiction \u2014 it is exactly what a multi-level explanation should "
        "reveal. The parent-level story is regime selection; the cluster-level story is variation within "
        "a regime. A variable can be globally important yet locally uninformative within a "
        "sub-cluster. Proposition 6.1 lets us interpret these shifts as meaningful structure rather than "
        "noise; a flat, single-level explanation would present them as conflicting, whereas our "
        "framework makes them coherent.", chip=2)

    contrib("Contribution II \u2014 Multi-Level XAI for Large-Scale Clustering",
        ["Objectives", "Methodology", "Results", "Findings"],
        [lambda sl: add_bullets(sl, 0.4, 1.9, 12.5, 4.6, [
            ("Wine: small, dense, chemically correlated. Beijing: large, noisy, temporally and meteorologically variable.", 0),
            ("The same explanatory logic remains productive in both \u2192 not tied to one domain-specific peculiarity.", 0, True),
            ("vs SHAP-based clustering literature: Beijing Silhouette \u2248 0.63 vs Gramegna & Giudici credit-risk 0.37.", 1),
            ("LIME comparator: weaker structural coherence, less stable local narratives for hierarchical reasoning.", 1),
        ])],
        "The same logic works on both a small, dense, chemically correlated dataset and a large, noisy, "
        "environmentally variable one, which supports the generality of the approach. It is not tied to "
        "a domain-specific peculiarity of wine chemistry. I also compared against the SHAP-based "
        "clustering literature: our Beijing partition achieves a Silhouette of about 0.63, whereas "
        "Gramegna and Giudici, working on credit-risk data, report around 0.37. On the LIME comparator, "
        "we observe weaker structural coherence and less stable local narratives for hierarchical "
        "reasoning. I want to be careful here: these are comparative observations rather than a "
        "comprehensive benchmark, but they support the claim that a cooperative-attribution approach "
        "produces more structured, more coherent explanations for hierarchical clustering.", chip=2)

    contrib("Contribution II \u2014 Multi-Level XAI for Large-Scale Clustering",
        ["Objectives", "Methodology", "Results", "Findings"],
        [lambda sl: add_bullets(sl, 0.4, 1.9, 12.5, 4.6, [
            ("Answers to the questions and objectives", 0, True),
            ("RQ2 answered in the affirmative, with bounds \u2014 coherence retained under scale and hierarchy.", 1),
            ("O1 met \u00b7 a multi-granular explanation that does not collapse into a single flat summary.", 1),
            ("O2 met \u00b7 Proposition 6.1 provides a formal cross-level consistency argument.", 1),
            ("O3 met \u00b7 validated on a structurally different large-scale dataset (Beijing air quality).", 1),
        ])],
        "To answer RQ2 directly: yes, with bounds. Shapley-based clustering explanation can scale to "
        "hierarchical, large-scale settings without losing interpretive coherence, provided the hierarchy "
        "is modelled explicitly and the approximation is transparent. O1 is met: we obtain a "
        "multi-granular explanation that does not collapse into a single flat summary. O2 is met: "
        "Proposition 6.1 gives a formal cross-level consistency argument, so regime-level and "
        "cluster-level stories are not at odds. O3 is met: the approach validates on Beijing air quality, "
        "a large, noisy, temporally variable dataset, which supports generality. The honest bound is "
        "that the model is still an explanation of a pre-computed partition; it does not yet influence "
        "learning itself. That is the bridge to Contribution III.", chip=2)

    contrib("Contribution II \u2014 Multi-Level XAI for Large-Scale Clustering",
        ["Objectives", "Methodology", "Results", "Findings"],
        [lambda sl: add_bullets(sl, 0.4, 1.9, 12.5, 4.6, [
            ("Key findings", 0, True),
            ("Scalable, multi-granular explanation that does not collapse into a single flat summary.", 1),
            ("Formal cross-level consistency argument (Proposition 6.1).", 1),
            ("Validated on a structurally different large-scale dataset (Beijing air quality).", 1),
            ("Differences across levels are interpretable, not dismissed as inconsistency.", 1),
        ])],
        "C2 answers RQ2 in the affirmative. The key findings are fourfold. The explanation scales to a "
        "multi-granular structure without collapsing into a single flat summary, because the hierarchy is "
        "modelled explicitly rather than flattened. Proposition 6.1 provides a formal cross-level "
        "consistency argument, so that differences across levels can be interpreted as meaningful structure "
        "rather than dismissed as noise. The approach generalises to a structurally different large-scale "
        "dataset, Beijing air quality, which is far noisier and more temporally variable than the small "
        "correlated wine data. And the method produces parent-level regime-selection stories alongside "
        "cluster-level within-regime variation.", chip=3)

    contrib("Contribution II \u2014 Multi-Level XAI for Large-Scale Clustering",
        ["Objectives", "Methodology", "Results", "Findings"],
        [lambda sl: add_bullets(sl, 0.4, 1.9, 12.5, 4.6, [
            ("Limitations", 0, True),
            ("Clustering remains static, even though the Beijing data are temporal.", 1),
            ("Surrogate-based SHAP plus representative-instance reporting compress observation-level variation.", 1),
            ("Confined to tabular data.", 1),
            ("Still an explanation of a pre-computed partition \u2014 it does not influence learning.", 1, True),
        ])],
        "The limitations are equally clear. The clustering remains static even though the Beijing data are "
        "temporal, so the dynamics are not exploited. The surrogate-based SHAP plus representative-instance "
        "reporting compresses observation-level variation, so some richness is lost. The whole approach is "
        "confined to tabular data. And, most importantly for what follows, the model is still an explanation "
        "of a pre-computed partition: it does not yet influence learning itself. This is the bridge to "
        "Contribution III.", chip=3)

    contrib("Contribution II \u2014 Multi-Level XAI for Large-Scale Clustering",
        ["Objectives", "Methodology", "Results", "Findings"],
        [lambda sl: add_bullets(sl, 0.4, 1.9, 12.5, 4.6, [
            ("Takeaways", 0, True),
            ("Shapley attribution stays coherent across granularity \u2014 when the hierarchy is explicit.", 1),
            ("Explanations become interpretable against scale, not just against a single flat partition.", 1),
            ("But the attribution is still post-hoc: it explains a partition that was already computed.", 1, True),
            ("So the next step is to let attribution shape the learning itself.", 0, True),
        ])],
        "The takeaway is that Shapley attribution stays coherent across granularity, provided the hierarchy "
        "is made explicit. Explanations become interpretable against scale, not just against a single flat "
        "partition, which is precisely where the literature was weakest. But the attribution is still "
        "post-hoc: it explains a partition that was already computed, and it never influences how the model "
        "learns. That is the limit that motivates Contribution III \u2014 can cooperative attribution move "
        "beyond explanation and enter the learning dynamics of a recommender?", chip=3)

    # ---- 20. SECTION: Contribution III ----
    section(SEC_CIII, "Contribution III \u2014 DyHuCoG: A Dynamic Hypergraph Cooperative Game", active=6, note="The third and principal contribution introduces DyHuCoG, a Dynamic Hypergraph Cooperative "
            "Game, where Shapley attribution becomes an in-training signal inside a hypergraph "
            "recommender rather than a post-hoc explanation. It answers RQ3 and RQ4.")
    contrib("Contribution III \u2014 DyHuCoG: A Dynamic Hypergraph Cooperative Game",
        ["Objectives", "Methodology", "Results", "Findings"],
        [lambda sl: add_bullets(sl, 0.4, 1.9, 12.5, 4.6, [
            ("Gap: graph and hypergraph recommenders treat message importance as either uniform or attention-weighted, without a principled marginal-contribution account.", 0, True),
            ("Diversity is often a secondary objective or a re-ranking heuristic.", 1),
            ("Interpretability is added after prediction, not integrated into the learning objective.", 1),
            ("Objectives: formulate recommendation as a cooperative game; embed preference-aware Monte Carlo Shapley into message passing; improve ranking, coverage, and diversity jointly.", 0, True),
        ])],
        "C3 is the flagship contribution and the conceptual shift of the whole thesis. The gap is that "
        "graph and hypergraph recommenders treat message importance as either uniform or "
        "attention-weighted, without a principled account of marginal contribution. Diversity is "
        "typically a secondary objective or a re-ranking heuristic, not part of the core learning "
        "problem; and interpretability is added after prediction rather than integrated into the "
        "objective. Our objectives are threefold: formulate recommendation as a cooperative game in "
        "which users, items and contexts are players; embed preference-aware Monte Carlo Shapley "
        "estimates into hypergraph message passing so that attribution directly shapes how information "
        "flows; and improve ranking accuracy, coverage and intra-list diversity jointly. The strongest "
        "claim is not that we explain a recommender; it is that attribution becomes an in-training "
        "signal.", chip=0)

    contrib("Contribution III \u2014 DyHuCoG: A Dynamic Hypergraph Cooperative Game",
        ["Objectives", "Methodology", "Results", "Findings"],
        [lambda sl: add_bullets(sl, 0.4, 1.9, 12.5, 4.6, [
            ("RQ3 \u00b7 Can cooperative attribution move beyond post-hoc and enter the learning dynamics of graph recommenders?", 0, True),
            ("RQ4 \u00b7 Can a recommender jointly optimise ranking accuracy, context and diversity when importance is estimated by a cooperative-game utility?", 1),
            ("Objective O1 \u00b7 Formulate recommendation as a cooperative game with users, items and contexts as players.", 1),
            ("Objective O2 \u00b7 Embed preference-aware Monte Carlo Shapley into hypergraph message passing.", 1),
            ("Objective O3 \u00b7 Improve ranking, coverage and diversity jointly.", 1),
        ])],
        "C3 answers RQ3 and RQ4 together, and it is the conceptual shift of the thesis. RQ3 asks whether "
        "cooperative attribution can move beyond post-hoc explanation and enter the learning dynamics of a "
        "graph recommender. RQ4 asks whether a recommender can jointly optimise ranking accuracy, context "
        "and diversity when importance is estimated by a cooperative-game utility. Both map to three "
        "objectives. O1 is to formulate recommendation as a cooperative game in which users, items and "
        "contexts are players. O2 is to embed preference-aware Monte Carlo Shapley estimates into "
        "hypergraph message passing, so attribution directly shapes how information flows. O3 is to improve "
        "ranking, coverage and intra-list diversity jointly \u2014 not at the expense of one another.", chip=0)

    contrib("Contribution III \u2014 DyHuCoG: A Dynamic Hypergraph Cooperative Game",
        ["Objectives", "Methodology", "Results", "Findings"],
        [lambda sl: add_bullets(sl, 0.4, 1.9, 12.5, 4.6, [
            ("Player set N = U \u222a I \u222a C (users, items, contexts).", 0, True),
            ("Hypergraph H = (V, E, W); V = U \u222a I \u222a C; W = dynamic edge weights from Shapley estimates.", 0),
            ("Coalition S \u2286 N represents entities participating in a recommendation episode.", 0),
            ("Coalition value v(S) measures the quality of the recommendation outcome achievable by S.", 0),
            ("Top-N task: produce a ranked list L_u balancing relevance, diversity, and contextual fit.", 0),
        ])],
        "We model recommendation as a cooperative game whose players are users, items and contexts. The "
        "hypergraph has vertices drawn from this same set, and hyperedges encode user-item-context "
        "interactions, with dynamic edge weights derived from Shapley estimates. A coalition represents "
        "the entities participating in a recommendation episode, and the coalition value measures the "
        "quality of the recommendation outcome achievable by that coalition. The task is top-N "
        "recommendation: we produce a ranked list for a user that balances relevance, diversity and "
        "contextual fit. This parallels the clustering formulation, but with a value function that is "
        "recommendation-oriented. That parallelism is deliberate and is one of the thesis-level "
        "contributions: the same cooperative-game logic applies to both tasks.", chip=1)

    contrib("Contribution III \u2014 DyHuCoG: A Dynamic Hypergraph Cooperative Game",
        ["Objectives", "Methodology", "Results", "Findings"],
        [lambda sl: add_bullets(sl, 0.4, 1.9, 12.5, 4.6, [
            ("v(S) = \u03b1 \u00b7 NDCG@20(S) + \u03b2 \u00b7 Diversity(S) + \u03b3 \u00b7 ContextScore(S), with \u03b1 + \u03b2 + \u03b3 = 1.", 0, True),
            ("The same trade-off the recommender must satisfy is the trade-off from which attribution is computed \u2014 explanatory game and predictive objective are aligned by design.", 1),
            ("Preference-weighted: v_pref(S) = v(S) + \u03bb_pref \u00b7 \u03a3_{(u,i)\u2208S} sim(u,i).", 0, True),
            ("\u03b1 = 0.60, \u03b2 = 0.25, \u03b3 = 0.15; \u03bb_pref = 0.20 \u2014 grid-searched, stable (<1.5% variance in NDCG@20).", 1),
            ("Coalition evaluation scoped to the interaction episode (a few dozen players), not the full catalogue.", 1),
        ])],
        "The coalition utility combines ranking quality, diversity and context: a weighted sum of NDCG "
        "at twenty, a diversity term and a context score, with weights summing to one. The key point is "
        "that this is the same trade-off the recommender must satisfy and also the trade-off from which "
        "attribution is computed, so the explanatory game and the predictive objective are aligned by "
        "design. We add a preference-consistency bonus proportional to the sum of user-item similarities "
        "in the coalition. Alpha, beta, gamma were grid-searched to 0.60, 0.25, 0.15, with the preference "
        "coefficient at 0.20 and NDCG varying by under one and a half percent. One honest boundary: "
        "coalition evaluation is scoped to the interaction episode, a few dozen players, not the full "
        "catalogue.", chip=1)

    contrib("Contribution III \u2014 DyHuCoG: A Dynamic Hypergraph Cooperative Game",
        ["Objectives", "Methodology", "Results", "Findings"],
        [lambda sl: add_bullets(sl, 0.4, 1.9, 12.5, 4.6, [
            ("Exact Shapley is combinatorial and infeasible for realistic systems.", 0),
            ("Monte Carlo estimator: \u03c6\u0302_j = (1/M) \u03a3_m [ v(S_m \u222a {j}) \u2212 v(S_m) ].", 0, True),
            ("Preference-aware: \u03c6\u0302_j^pref = (1/M) \u03a3_m [ v_pref(S_m \u222a {j}) \u2212 v_pref(S_m) ].", 0, True),
            ("Unbiased; variance = \u03c3\u00b2/M \u2192 MSE decays O(1/M), absolute error O(1/\u221aM).", 1),
            ("M = 50 selected: MSE \u2248 1.4\u00d710\u207b\u2075, ~99% accuracy on MovieLens-1M.", 1, True),
            ("Refreshed every 10 batches (~49 updates/epoch), smoothed by exponential moving average.", 1),
        ])],
        "Exact Shapley computation is combinatorial and infeasible for a realistic recommender, so we "
        "use a Monte Carlo estimator that averages, over M sampled coalitions, the difference between "
        "the value including a player and the value excluding it; the preference-aware variant applies "
        "the same estimator to the preference-weighted utility. It is unbiased, with variance decaying "
        "as sigma squared over M, so mean squared error goes as one over M and absolute error as one "
        "over root M. We chose M equal to fifty: mean squared error about one point four times ten to "
        "the minus five, and roughly ninety-nine percent accuracy. Because preferences evolve, estimates "
        "are refreshed every ten batches, about forty-nine updates per epoch, and smoothed with an "
        "exponential moving average, so attribution is adaptive without making training hypersensitive "
        "to a single estimate.", chip=1)

    contrib("Contribution III \u2014 DyHuCoG: A Dynamic Hypergraph Cooperative Game",
        ["Objectives", "Methodology", "Results", "Findings"],
        [lambda sl: add_bullets(sl, 0.4, 1.9, 12.5, 4.6, [
            ("Base propagation: e^(l+1) = \u03c3( D^-1/2 A D^-1/2 e^(l) ).", 0),
            ("Shapley-weighted: e_j^(l+1) = \u03c3( W^(l) e_j^(l) + \u03a3_{k\u2208N(j)} w_jk e_k^(l) ).", 0, True),
            ("Normalised weights: w_jk = \u03c6\u0302_jk / \u03a3_{k'\u2208N(j)} \u03c6\u0302_jk'.", 0, True),
            ("Clipped + exponentially smoothed before normalisation (stabilises sparse regimes).", 1),
            ("Attention gate: a_ui = \u03c3( W_a[ e_u, e_i, l_i ] ); y_ui = (1 + a_ui) \u27e8e_u, e_i\u27e9.", 1),
            ("Context-aware score: f(u,i,c) = y_ui + \u03bb_c \u27e8g(c_ui), e_cui\u27e9.", 1),
        ]),
        lambda sl: add_figure(sl, os.path.join(HERE, '_figs', 'c3_dyhucog.png'), 1.6, 3.75, w=10.1)],
        "The architecture is the decisive move. The base propagation is standard hypergraph message "
        "passing; the Shapley-weighted version weights each message by a normalised Shapley coefficient, "
        "dividing by the sum of coefficients over the neighbourhood. I clip and exponentially smooth the "
        "estimates before normalising, which stabilises the sparse regime. Layer fusion combines "
        "embeddings across layers with learned coefficients. On top of that, an interaction-level "
        "attention gate interpolates between the Shapley-weighted score and the standard inner-product "
        "score, acting as a stabiliser; a context-aware score then adds a context term. The essential "
        "point is that the model is told not only who is connected to whom, but how much each coalition "
        "is worth, and that worth directly governs how information propagates. That is what "
        "distinguishes it from a post-hoc explainer.", chip=1)

    contrib("Contribution III \u2014 DyHuCoG: A Dynamic Hypergraph Cooperative Game",
        ["Objectives", "Methodology", "Results", "Findings"],
        [lambda sl: add_bullets(sl, 0.4, 1.9, 12.5, 4.6, [
            ("L = L_rec + \u03bb_div L_div + \u03bb_ctx L_ctx + \u03bb_reg L_reg.", 0, True),
            ("L_rec \u2013 Bayesian Personalised Ranking (pairwise, implicit feedback).", 1),
            ("L_div \u2013 Intra-List Diversity regulariser: penalises redundant ranked lists.", 1),
            ("L_ctx \u2013 Context alignment: match context embedding to context-node representation.", 1),
            ("L_reg \u2013 L2 weight decay.", 1),
            ("The learning objective and coalition value are aligned: DyHuCoG trains to optimise the same balance that later determines attribution.", 1, True),
        ])],
        "The composite loss combines four terms. The recommendation loss is Bayesian Personalised "
        "Ranking, which is a pairwise objective appropriate for implicit feedback. The diversity loss is "
        "an intra-list diversity regulariser that penalises redundant ranked lists. The context loss "
        "aligns the context embedding with the context-node representation. And the regularisation term "
        "is weight decay. Negatives are drawn from a popularity-aware distribution, with periodic "
        "hard-negative refresh, and optimisation uses Adam. The conceptual point I want to drive home is "
        "alignment: the model is trained to optimise the same accuracy-diversity-context balance that "
        "later determines cooperative attribution. So the explanation is not a separate diagnostic; it is "
        "a direct read-out of the objective the model is already optimising. This is what gives the "
        "explanation its structural faithfulness.", chip=1)

    contrib("Contribution III \u2014 DyHuCoG: A Dynamic Hypergraph Cooperative Game",
        ["Objectives", "Methodology", "Results", "Findings"],
        [lambda sl: add_table(sl, 0.4, 1.9, 12.6, 3.9,
             ["Dataset", "Model", "NDCG@20", "Recall@20", "Coverage", "Diversity"],
             [["MovieLens-1M", "HPCF", "0.2528", "0.2098", "0.342", "0.461"],
              ["MovieLens-1M", "DyHuCoG", "0.2775", "0.2362", "0.397", "0.516"],
              ["Amazon-Book", "HPCF", "0.0270", "0.0359", "0.259", "0.535"],
              ["Amazon-Book", "DyHuCoG", "0.0306", "0.0417", "0.336", "0.602"]],
             col_ratios=[3.2, 2, 2.2, 2.2, 2.2, 2.2], hl_rows=[1, 3])],
        "This is the headline result, and I would ask you to read it carefully because it shows ranking "
        "and diversity improving together. On MovieLens-1M, the strongest baseline HPCF achieves NDCG at "
        "twenty of 0.2528 and recall of 0.2098; DyHuCoG improves these to 0.2775 and 0.2362, plus nine "
        "point seven seven percent in NDCG and plus twelve point five eight percent in recall. At the "
        "same time, coverage rises from 0.342 to 0.397 and intra-list diversity from 0.461 to 0.516. On "
        "Amazon-Book, which is much sparser, the relative gains are larger: plus thirteen point three "
        "three percent NDCG and plus sixteen point one six percent recall, with coverage up by nearly "
        "thirty percent. The sparser the data, the larger the gain \u2014 precisely the pattern you "
        "would expect if Shapley-driven weighting is most valuable when signal is weak.", chip=2)

    contrib("Contribution III \u2014 DyHuCoG: A Dynamic Hypergraph Cooperative Game",
        ["Objectives", "Methodology", "Results", "Findings"],
        [lambda sl: add_table(sl, 0.4, 1.9, 12.4, 3.6,
             ["Metric", "MovieLens-1M", "Amazon-Book"],
             [["NDCG@20", "+9.77%", "+13.33%"],
              ["Recall@20", "+12.58%", "+16.16%"],
              ["Coverage", "+16.1%", "+29.7%"],
              ["Intra-List Diversity", "+11.9%", "+12.5%"]],
             col_ratios=[4, 3.2, 3.2])],
        "This slide summarises the relative gains over the strongest baseline, and the key message is "
        "that DyHuCoG improves ranking accuracy, coverage and diversity simultaneously \u2014 it does "
        "not sacrifice one for the others. The largest relative gains appear on the sparser Amazon-Book "
        "dataset, consistent with the hypothesis that Shapley-guided weighting helps most precisely when "
        "interaction data are weak. Together these results are evidence that the accuracy-diversity "
        "trade-off, often treated as structurally fixed, is in fact negotiable if attribution is handled "
        "as a first-class part of the learning objective. NDCG up, recall up, coverage up, diversity up: "
        "that combination is the core empirical claim of the third contribution.", chip=2)

    contrib("Contribution III \u2014 DyHuCoG: A Dynamic Hypergraph Cooperative Game",
        ["Objectives", "Methodology", "Results", "Findings"],
        [lambda sl: add_bullets(sl, 0.4, 1.9, 12.5, 4.6, [
            ("MovieLens-1M: Coverage 0.342 \u2192 0.397 (+16.1%); ILD 0.461 \u2192 0.516 (+11.9%).", 0, True),
            ("Amazon-Book: Coverage 0.259 \u2192 0.336 (+29.7%); ILD 0.535 \u2192 0.602 (+12.5%).", 0, True),
            ("Reduced filter-bubble effect and greater discovery opportunity \u2014 while NDCG/Recall also improve, so accuracy is not traded off for diversity.", 1),
        ])],
        "Both levels of diversity improve. On MovieLens-1M, catalogue coverage rises from 0.342 to "
        "0.397, a sixteen percent improvement, and intra-list diversity from 0.461 to 0.516, an eleven "
        "point nine percent improvement. On Amazon-Book, coverage rises from 0.259 to 0.336, nearly "
        "thirty percent, and intra-list diversity from 0.535 to 0.602. The practical consequence is a "
        "reduced filter-bubble effect and greater discovery opportunity for the user. The point I want "
        "to stress: both diversity metrics improve while NDCG and recall also improve \u2014 so we are "
        "not trading accuracy for diversity. The model surfaces more of the catalogue, recommends less "
        "redundant lists, and ranks better at the same time. That is the strongest evidence that "
        "diversity is a genuine objective here, not a post-hoc re-ranking heuristic.", chip=2)

    contrib("Contribution III \u2014 DyHuCoG: A Dynamic Hypergraph Cooperative Game",
        ["Objectives", "Methodology", "Results", "Findings"],
        [lambda sl: add_table(sl, 0.4, 1.9, 12.5, 3.9,
             ["Variant", "ML-1M NDCG@20", "% Drop", "Amazon NDCG@20", "% Drop"],
             [["Full DyHuCoG", "0.2775", "\u2013", "0.0306", "\u2013"],
              ["w/o Shapley Value", "0.2647", "4.6%", "0.0287", "6.1%"],
              ["w/o Hypergraph", "0.2586", "6.8%", "0.0279", "8.9%"],
              ["w/o Attention", "0.2678", "3.5%", "0.0295", "3.5%"],
              ["w/o Context", "0.2547", "8.2%", "0.0272", "11.0%"],
              ["w/o Diversity", "0.2614", "5.8%", "0.0288", "5.8%"]],
             col_ratios=[3.6, 2.4, 1.8, 2.4, 1.8], hl_rows=[1, 5])],
        "This ablation isolates each component, and every component contributes. The full model achieves "
        "0.2775 NDCG at twenty on MovieLens-1M. Removing the Shapley value drops it to 0.2647, a four "
        "point six percent loss. Removing the hypergraph structure drops it to 0.2586, a six point eight "
        "percent loss. Removing the attention gate costs three point five percent. Removing context "
        "causes the largest single loss at eight point two percent, consistent with context providing the "
        "representational substrate on which Shapley weighting operates. Removing diversity costs five "
        "point eight percent. The same pattern holds on Amazon-Book, with context removal again largest. "
        "The four to six percent drop when Shapley is removed supports the argument that "
        "marginal-contribution estimation is load-bearing rather than decorative.", chip=2)

    contrib("Contribution III \u2014 DyHuCoG: A Dynamic Hypergraph Cooperative Game",
        ["Objectives", "Methodology", "Results", "Findings"],
        [lambda sl: add_bullets(sl, 0.4, 1.9, 12.5, 4.6, [
            ("Training: DyHuCoG ~2000 s vs HPCF ~1125 s on MovieLens-1M (\u2248 1.78\u00d7).", 0),
            ("Inference: 1.84 ms/query (ML-1M), 8.52 ms (Amazon) \u2014 suitable for real-time deployment.", 0),
            ("Memory: 4.4 vs 4.1 GB (ML-1M); 17.9 vs 16.8 GB (Amazon).", 0),
            ("Per-epoch cost: O((L+1)md) + O((M/f)m).", 0),
            ("Shapley convergence: M=50 \u2192 MSE 1.4\u00d710\u207b\u2075, ~99% accuracy; M=100 \u2192 MSE 3.5\u00d710\u207b\u2076 (diminishing returns).", 0, True),
        ])],
        "The attribution cost is proportionate, and I report it transparently because it bears on "
        "deployability. On MovieLens-1M, DyHuCoG takes about two thousand seconds to train versus "
        "roughly eleven hundred for HPCF, about one point seven eight times. Inference is one point eight "
        "four milliseconds per query on MovieLens and eight point five two on Amazon, comfortably within "
        "real-time requirements. Memory is modestly higher, four point four gigabytes versus four point "
        "one. The per-epoch cost scales with the number of layers, the embedding dimension and the Monte "
        "Carlo budget. On Shapley convergence, M equal to fifty gives ninety-nine percent accuracy, while "
        "M equal to one hundred gives ninety-nine point five with diminishing returns. So fifty is the "
        "right operating point, balancing accuracy against cost.", chip=2)

    contrib("Contribution III \u2014 DyHuCoG: A Dynamic Hypergraph Cooperative Game",
        ["Objectives", "Methodology", "Results", "Findings"],
        [lambda sl: add_bullets(sl, 0.4, 1.9, 12.5, 4.6, [
            ("Paired t-tests on per-user NDCG@20 (n = 6,040 users; df = 6,039).", 0),
            ("DyHuCoG outperforms every baseline with extremely small p-values after Holm\u2013Bonferroni correction.", 0),
            ("vs HPCF: t = 46.38, Cohen\u2019s d_z = 1.3345, p = 1.81\u00d710\u207b\u00b2\u2077\u2070.", 1, True),
            ("Wilcoxon signed-rank test also significant (p < 0.001).", 1),
            ("Effect sizes are large \u2014 improvements are substantively meaningful, not merely statistically visible.", 1),
        ])],
        "I also validated that the improvements are statistically significant and substantively "
        "meaningful. I ran paired t-tests on per-user NDCG at twenty across six thousand and forty "
        "users, comparing DyHuCoG against every baseline, with Holm-Bonferroni correction. Against HPCF, "
        "the t-statistic is 46.38, Cohen's d is 1.33, and the p-value is on the order of ten to the "
        "minus two hundred and seventy. The Wilcoxon signed-rank test is also significant at p less than "
        "0.001. The effect sizes are large, and that matters: with six thousand users, even a trivially "
        "small difference can appear significant, so reporting effect size shows the improvement is not "
        "merely a large-sample artefact. I should be scrupulous about scope: these tabulated paired tests "
        "apply only to MovieLens-1M; the Amazon-Book results remain descriptive.", chip=2)

    contrib("Contribution III \u2014 DyHuCoG: A Dynamic Hypergraph Cooperative Game",
        ["Objectives", "Methodology", "Results", "Findings"],
        [lambda sl: add_bullets(sl, 0.4, 1.9, 12.5, 4.6, [
            ("Cold-start (5 or fewer interactions): NDCG@20 \u2248 0.061 (user) and 0.057 (item), improving over HPCF by ~10%.", 0),
            ("Cross-dataset: MovieLens +9.9%, Amazon +14.8%, Yelp2018 +11.8%.", 0),
            ("Interpretability: a SHAP waterfall decomposes a recommendation into ranking, diversity, context and preference contributions.", 0, True),
            ("Popularity bias: Shapley measures marginal utility, not raw frequency \u2014 weak but informative interactions retain influence.", 1),
        ])],
        "DyHuCoG also improves the regimes where recommenders are at their most brittle: cold-start. For "
        "users with five or fewer interactions, NDCG at twenty is about 0.061, roughly ten percent over "
        "HPCF; for cold-start items about 0.057, a gain of roughly nine point six percent. Across "
        "datasets the improvement holds at plus nine point nine percent on MovieLens, plus fourteen point "
        "eight on Amazon, and plus eleven point eight on Yelp2018, which I use only as an auxiliary "
        "robustness benchmark. On interpretability, the model produces a SHAP waterfall that decomposes a "
        "recommendation into the same ranking, diversity, context and preference components used during "
        "training, so the explanation is structurally faithful rather than an external approximation. "
        "And Shapley measures marginal utility rather than raw frequency, so weak but informative "
        "interactions retain influence, mitigating popularity bias.", chip=2)

    contrib("Contribution III \u2014 DyHuCoG: A Dynamic Hypergraph Cooperative Game",
        ["Objectives", "Methodology", "Results", "Findings"],
        [lambda sl: add_bullets(sl, 0.4, 1.9, 12.5, 4.6, [
            ("Answers to the questions and objectives", 0, True),
            ("RQ3 answered in the affirmative \u2014 attribution becomes an in-training signal, not a post-hoc diagnostic.", 1),
            ("RQ4 answered in the affirmative \u2014 ranking, coverage and diversity improve together.", 1),
            ("O1 met \u00b7 recommendation formulated as a cooperative game over users, items and contexts.", 1),
            ("O2 met \u00b7 preference-aware Monte Carlo Shapley embedded in message passing.", 1),
            ("O3 met \u00b7 +9.9% (MovieLens) / +14.8% (Amazon) NDCG, with higher coverage and diversity.", 1),
        ])],
        "To answer RQ3 and RQ4 directly: yes. RQ3 is answered in the affirmative because attribution is no "
        "longer a post-hoc diagnostic but an in-training signal that directly shapes message passing. RQ4 "
        "is answered in the affirmative because ranking, coverage and diversity improve together rather "
        "than trading against one another. Each objective is met. O1: recommendation is formulated as a "
        "cooperative game with users, items and contexts as players. O2: preference-aware Monte Carlo "
        "Shapley is embedded into hypergraph message passing. O3: we gain about nine point nine percent "
        "on MovieLens and fourteen point eight percent on Amazon, with higher coverage and diversity. The "
        "strongest claim is that the explanation is structurally faithful, because it reads out the same "
        "components the model already optimises. Within the stated edges \u2014 computational overhead, "
        "dependence on context, and the early-2026 baseline set \u2014 the claim stands.", chip=2)

    contrib("Contribution III \u2014 DyHuCoG: A Dynamic Hypergraph Cooperative Game",
        ["Objectives", "Methodology", "Results", "Findings"],
        [lambda sl: add_bullets(sl, 0.4, 1.9, 12.5, 4.6, [
            ("Key findings", 0, True),
            ("Cooperative attribution used as an in-training signal \u2014 a stronger claim than the clustering chapters.", 1),
            ("The accuracy\u2013diversity\u2013context trade-off is not structurally fixed.", 1),
            ("NDCG, recall, coverage and diversity improve together on both datasets.", 1),
            ("Largest gains on the sparsest data (Amazon-Book), consistent with Shapley weighting helping when signal is weak.", 1),
        ])],
        "C3 answers RQ3 and RQ4, and makes the strongest claim of the thesis. Cooperative attribution can "
        "move beyond post-hoc analysis and enter the learning dynamics of a recommender, letting it balance "
        "accuracy, diversity and context more effectively. The key findings are that ranking, coverage and "
        "diversity improve together \u2014 NDCG, recall, coverage and intra-list diversity all rise on both "
        "datasets \u2014 so the accuracy-diversity trade-off is not structurally fixed. The largest gains "
        "appear on the sparsest data, Amazon-Book, precisely the regime where Shapley-driven weighting is "
        "most valuable, and the improvement is statistically large against the strongest baseline.", chip=3)

    contrib("Contribution III \u2014 DyHuCoG: A Dynamic Hypergraph Cooperative Game",
        ["Objectives", "Methodology", "Results", "Findings"],
        [lambda sl: add_bullets(sl, 0.4, 1.9, 12.5, 4.6, [
            ("Limitations", 0, True),
            ("Measurable computational overhead \u2014 roughly 1.78\u00d7 training time over HPCF.", 1),
            ("Depends on availability of meaningful context.", 1),
            ("Monte Carlo Shapley could be improved by variance reduction.", 1),
            ("Ablation is component-wise, so it does not test factorial interactions.", 1),
            ("Baselines finalised through early 2026, so superiority is claimed only against the tested baselines.", 1),
        ])],
        "The limitations must be stated honestly. There is a measurable computational overhead, roughly "
        "one point seven eight times the training time of the strongest baseline, HPCF. The method depends "
        "on the availability of meaningful context. The Monte Carlo Shapley estimator could be improved by "
        "variance reduction. The ablation is component-wise, so it does not test factorial interactions. "
        "And the baselines were finalised in early 2026, so I claim superiority only against the tested "
        "baselines. Within those bounds, the claim stands.", chip=3)

    contrib("Contribution III \u2014 DyHuCoG: A Dynamic Hypergraph Cooperative Game",
        ["Objectives", "Methodology", "Results", "Findings"],
        [lambda sl: add_bullets(sl, 0.4, 1.9, 12.5, 4.6, [
            ("Takeaways", 0, True),
            ("Attribution is a first-class part of the learning objective, not a post-hoc diagnostic.", 1),
            ("The explanation is a direct read-out of the objective the model already optimises.", 1),
            ("This makes the explanation structurally faithful rather than an external approximation.", 1, True),
            ("Built on a principled, axiomatic basis \u2014 aligned with trustworthy-AI expectations.", 1),
        ])],
        "The takeaway is the conceptual shift of the thesis. Attribution is a first-class part of the "
        "learning objective, not a post-hoc diagnostic, and the explanation is a direct read-out of the "
        "objective the model already optimises. That is what makes it structurally faithful rather than an "
        "external approximation. And because Shapley values rest on an axiomatic, normative basis, the "
        "perspective is well aligned with the transparency and accountability requirements of emerging AI "
        "regulation. This is the central claim worth defending: explanation should be part of the modelling "
        "logic itself.", chip=3)

    # ---- 21. SECTION: Conclusion ----
    section(SEC_CONCL, "Conclusion & Perspectives", active=7, note="Let me now bring everything together. I will present a synthesis of the three "
            "contributions, the published papers, the honest limitations, the future directions, and "
            "finally a clear statement of the thesis answer.")
    contrib("Conclusion & Perspectives",
        ["Synthesis", "Limitations", "Perspectives", "Conclusion"],
        [lambda sl: add_table(sl, 0.4, 1.9, 12.6, 3.9,
             ["Contribution", "Main idea", "Achievement", "Key finding"],
             [["C1", "Explain black-box clustering via Shapley", "PCA\u2013KMeans\u2013LightGBM\u2013TreeSHAP pipeline", "Faithful, chemistry-consistent cluster attribution (wine)"],
              ["C2", "Multi-level, large-scale clustering XAI", "Cross-level SHAP aggregation + Prop. 6.1", "Coherent under hierarchy; interprets differences"],
              ["C3", "DyHuCoG hypergraph cooperative game", "Preference-aware Shapley as in-training signal", "Accuracy + coverage + diversity improve together"]],
             col_ratios=[1.6, 4.2, 4.2, 5.0])],
        "Three contributions, one thread. The synthesis table shows how each contribution builds on the "
        "previous one. C1 makes hidden structure intelligible: it introduces a Shapley-based pipeline that "
        "explains black-box clustering faithfully, with faithful and chemistry-consistent cluster "
        "attribution on wine. C2 keeps that explanation coherent under scale and hierarchy: it introduces "
        "cross-level SHAP aggregation and Proposition 6.1, and shows that hierarchical differences can be "
        "interpreted rather than dismissed as inconsistency. C3 carries the same attribution logic inside "
        "the learning dynamics of a recommender, using preference-aware Shapley as an in-training signal, "
        "so that accuracy, coverage and diversity improve together. The common thread is that cooperative "
        "attribution is a single mechanism used for explanation, for interpretation and for optimisation \u2014 "
        "not three unrelated tools.", chip=0)

    contrib("Conclusion & Perspectives",
        ["Synthesis", "Limitations", "Perspectives", "Conclusion"],
        [lambda sl: add_table(sl, 0.4, 1.9, 12.6, 3.9,
             ["No.", "Title", "Venue", "Status"],
             [["I", "Shapley Values for Explaining the Black Box Nature of ML Model Clustering", "Procedia Computer Science 220, 806\u2013811", "Published, 2023"],
              ["II", "Game Theory Meets Explainable AI: An Enhanced Approach to Understanding Black Box Models Through Shapley Values", "IJACSA 16(7), 716\u2013725", "Published, 2025"],
              ["III", "DyHuCoG: A Dynamic Hypergraph Cooperative Game for Preference-aware Recommendation", "IJIES 19(2), 887\u2013902", "Published, 2026"]],
             col_ratios=[1, 6.5, 3.2, 2])],
        "The thesis synthesises three peer-reviewed publications, and it is worth mapping them to the "
        "chapters. Paper I, on Shapley values for explaining the black-box nature of machine-learning "
        "model clustering, is published in Procedia Computer Science in 2023 and corresponds to Chapter "
        "five. Paper II, on game theory meeting explainable AI for understanding black-box models, is "
        "published in IJACSA in 2025 and corresponds to Chapter six. Paper III, DyHuCoG, is published in "
        "IJIES in 2026 and corresponds to Chapter seven. What the thesis adds on top of the three papers "
        "is the multi-level formalisation, the cross-chapter comparisons, and the explicit mapping to the "
        "five research questions. In other words, the thesis is a cumulative argument, not merely a "
        "bound collection of three papers.", chip=0)

    contrib("Conclusion & Perspectives",
        ["Synthesis", "Limitations", "Perspectives", "Conclusion"],
        [lambda sl: add_bullets(sl, 0.4, 1.9, 12.5, 4.6, [
            ("Computational \u2013 exact Shapley is intractable; every contribution relies on approximation, surrogates, or restricted reporting.", 0),
            ("Methodological \u2013 clustering depends on surrogate fidelity; recommendation depends on stable approximate contributions and adequate context.", 0),
            ("Empirical \u2013 tabular clustering + benchmark recommendation; no multimodal, sequential, or online deployment; no dedicated human-subject actionability study.", 0),
            ("Claim scope \u2013 a coherent and productive perspective, not one fully unified framework eliminating all tension.", 0, True),
        ])],
        "Let me state the limitations honestly. Computationally, exact Shapley is intractable, so every "
        "contribution relies on approximation, surrogates, or restricted reporting. Methodologically, the "
        "clustering chapters depend on surrogate fidelity, and the recommendation results depend on "
        "stable approximate contributions and meaningful context. Empirically, the clustering work is "
        "confined to tabular data and the recommendation work is offline: no multimodal, sequential or "
        "online deployment, and no dedicated human-subject study of actionability. And on the claim "
        "itself, I want to be precise. I am not claiming one fully unified framework that eliminates all "
        "tension between accuracy and interpretability; I am claiming a shared perspective, a common "
        "attribution language, coherent and productive across the three tasks. That is the honest "
        "boundary of the thesis.", chip=1)

    contrib("Conclusion & Perspectives",
        ["Synthesis", "Limitations", "Perspectives", "Conclusion"],
        [lambda sl: add_bullets(sl, 0.4, 1.9, 12.5, 4.6, [
            ("Scalable cooperative attribution \u2013 lower-variance Shapley, learned proposal distributions, adaptive refresh policies.", 0),
            ("Online / streaming recommendation \u2013 truly incremental settings with evolving graphs and delayed feedback.", 0),
            ("Richer human-centred evaluation \u2013 do explanations measurably improve analyst judgement, user trust, intervention quality, or perceived fairness?", 0),
            ("Broader trustworthy-AI evaluation \u2013 exposure fairness, transparency requirements, governance-oriented auditing.", 0),
        ])],
        "Future work turns the limitations into a concrete agenda. First, scalable cooperative "
        "attribution: lower-variance Shapley estimators, learned proposal distributions, and adaptive "
        "refresh policies, to close the computational gap. Second, online and streaming recommendation: "
        "truly incremental settings with continuously evolving graphs and delayed feedback, which would "
        "also address the static-graph limitation. Third, richer human-centred evaluation: I would like "
        "to test whether these explanations measurably improve analyst judgement, user trust, "
        "intervention quality, or perceived fairness \u2014 which is the natural test of the "
        "\"actionable\" part of my definition. Fourth, broader trustworthy-AI and fairness evaluation: "
        "exposure fairness, transparency requirements, and governance-oriented auditing, which connects "
        "the work to the EU AI Act, the OECD principles and the GDPR.", chip=2)

    contrib("Conclusion & Perspectives",
        ["Synthesis", "Limitations", "Perspectives", "Conclusion"],
        [lambda sl: add_bullets(sl, 0.4, 1.9, 12.5, 4.6, [
            ("Thesis answer: cooperative game theory can function as a shared methodological perspective for actionable explanation across clustering and recommendation.", 0, True),
            ("Key outcomes:", 0),
            ("Shapley attribution as a common formal language for feature, interaction, and context importance allocation.", 1),
            ("Faithful clustering explanation, hierarchical explanatory coherence, and contribution-aware recommendation learning.", 1),
            ("Explanation as method, not commentary \u2014 from post-hoc description to in-training guidance.", 1),
            ("Aligned with trustworthy-AI requirements (EU AI Act, OECD principles, GDPR).", 1),
        ])],
        "To close: the thesis answer is that cooperative game theory can function as a shared "
        "methodological perspective for actionable explanation across clustering and recommendation. "
        "The key outcomes are fourfold. Shapley attribution provides a common formal language for "
        "allocating importance to features, interactions and contexts. It yields faithful clustering "
        "explanation, hierarchical explanatory coherence, and contribution-aware recommendation "
        "learning. And crucially, it moves explanation from commentary to method: from post-hoc "
        "description to in-training guidance. Because Shapley values rest on an axiomatic, normative "
        "basis, this perspective is well aligned with the transparency and accountability requirements "
        "of emerging AI regulation. I believe this is the central claim worth defending: that "
        "explanation should be part of the modelling logic itself, not a by-product attached after "
        "prediction.", chip=3)

    # ---- 22. QUESTIONS ----
    s = clone_slide(prs, QUES_ARCH)
    set_lines(s, 'Rectangle 5', [FOOTER_TEXT])
    # shift the question block text (Rectangle 27 already holds the thank-you)
    set_slide_notes(s, "Thank you very much for your attention. I am now happy to take your "
        "questions and comments; I will do my best to answer them directly, and I am happy to go "
        "deeper into any of the methodology, the mathematics, or the experimental details.")

    # ---- remove original colleague slides (keep only ours) ----
    for sid in orig_ids:
        rId = sid.get(qn('r:id'))
        prs.part.drop_rel(rId)
        prs.slides._sldIdLst.remove(sid)

    # ---- teal accent chrome + high-contrast pager ----
    for slide in prs.slides:
        recolor_accent_teal(slide)
    digit_blobs = collect_digit_blobs(prs)
    remap_pager(prs, digit_blobs)

    prs.save(OUT)
    print(f"Saved {len(prs.slides)} slides -> {OUT}")

if __name__ == "__main__":
    main()
