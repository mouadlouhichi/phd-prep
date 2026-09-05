#!/usr/bin/env python3
"""
Build Mouad Louhichi's PhD viva presentation by REUSING the colleague's actual
deck as a template (exact theme, master, layouts, and logo/crest/number assets),
so the design and framework are identical to
  example-phd-passes/Presentation1 (1) (1).pptx
but the content is Mouad's thesis.

The colleague's PPTX is read from origin/main (via git show) so we don't need it
checked out. We then (1) keep its theme/master/layouts, (2) remove all its
content slides, (3) add our own slides using the SAME layout placeholders, and
(4) replicate its branded chrome (Morocco crest + ENSIAS logo top corners,
ENSIAS footer bar, numbered section headers).

Output: viva/MOUAD_LOUHICHI_VIVA_40min.pptx
"""
import os, subprocess, tempfile, shutil
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

HERE = os.path.dirname(os.path.abspath(__file__))
REPO = os.path.dirname(HERE)
OUT = os.path.join(HERE, "MOUAD_LOUHICHI_VIVA_40min.pptx")
SRC = os.path.join(HERE, "colleague_template.pptx")
RAW_SRC = "origin/main:example-phd-passes/Presentation1 (1) (1).pptx"

BLUE   = RGBColor(0x44, 0x72, 0xC4)   # theme accent1 (Office default)
DARK   = RGBColor(0x1F, 0x33, 0x55)
GREY   = RGBColor(0x40, 0x40, 0x40)
LGREY  = RGBColor(0x8A, 0x8A, 0x8A)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
HDR_BG = RGBColor(0xE9, 0xEE, 0xF6)

# ---------------------------------------------------------------------------
# Get template from git (origin/main) if not already present
# ---------------------------------------------------------------------------
def ensure_template():
    if not os.path.exists(SRC):
        with open(SRC, "wb") as f:
            subprocess.run(["git", "-C", REPO, "show", RAW_SRC],
                           stdout=f, check=True)
        print(f"Extracted template -> {SRC}")

# ---------------------------------------------------------------------------
def set_fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color

def no_line(shape):
    shape.line.fill.background()

def add_text(slide, x, y, w, h, text, size, color=GREY, bold=False,
             align=PP_ALIGN.LEFT, font="Calibri", italic=False,
             anchor=MSO_ANCHOR.TOP, wrap=True, spacing=1.0):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = spacing
    r = p.add_run()
    r.text = text
    r.font.name = font; r.font.size = Pt(size); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color
    return box

def add_bullets(slide, x, y, w, h, items, size=17, color=GREY, accent=BLUE,
                space=10, anchor=MSO_ANCHOR.TOP, spacing=1.05):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    first = True
    for it in items:
        if isinstance(it, str): text, level, bold = it, 0, False
        elif len(it) == 2: text, level = it; bold = False
        else: text, level, bold = it
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = level
        p.space_after = Pt(space if level == 0 else space - 3)
        p.line_spacing = spacing
        marker = "\u25aa  " if level == 0 else "\u2013  "
        r0 = p.add_run(); r0.text = marker
        r0.font.name = "Calibri"; r0.font.size = Pt(size); r0.font.bold = True
        r0.font.color.rgb = accent
        r = p.add_run(); r.text = text
        r.font.name = "Calibri"; r.font.size = Pt(size); r.font.bold = bold
        r.font.color.rgb = color
    return box

def add_rect(slide, x, y, w, h, color=None, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.shadow.inherit = False
    if color is not None:
        set_fill(shp, color)
    else:
        shp.fill.background()
    if line is not None:
        shp.line.color.rgb = line
        shp.line.width = Pt(1)
    else:
        no_line(shp)
    return shp

# ---------------------------------------------------------------------------
# Branded chrome assets (swapped to Mouad where needed, keep ENSIAS look)
# ---------------------------------------------------------------------------
def add_chrome_footer(slide, page_num, prs, footer_img, pager_left, pager_right):
    """Replicate the colleague's bottom bar: ENSIAS logo left + footer text + pager icons right."""
    # Footer text bar
    add_rect(slide, Inches(0.83), Inches(6.92), Inches(12.50), Inches(0.57),
             color=WHITE)
    add_text(slide, Inches(0.83), Inches(7.02), Inches(9.5), Inches(0.35),
             "PhD Viva \u2013 Mouad LOUHICHI \u2013 Cooperative Game Theory & Shapley for XAI in Recommendation Systems",
             size=11, color=LGREY)
    # ENSIAS logo (bottom-left)
    if footer_img:
        pic = slide.shapes.add_picture(footer_img, Inches(0.14), Inches(7.00),
                                       Inches(0.48), Inches(0.41))
    # pager icons (bottom-right) using template assets
    if pager_left:
        slide.shapes.add_picture(pager_left, Inches(12.50), Inches(7.00),
                                 Inches(0.42), Inches(0.42))
    if pager_right:
        slide.shapes.add_picture(pager_right, Inches(12.81), Inches(7.00),
                                 Inches(0.42), Inches(0.42))
    # slide number
    add_text(slide, Inches(12.1), Inches(7.02), Inches(0.5), Inches(0.35),
             str(page_num), size=11, color=LGREY, align=PP_ALIGN.RIGHT)

def add_chrome_header_logos(slide, crest, ensias_logo):
    """Top-right corner ENSIAS logo (and small crest on some)."""
    if ensias_logo:
        slide.shapes.add_picture(ensias_logo, Inches(10.69), Inches(0.03),
                                 Inches(2.32), Inches(1.64))
    if crest:
        slide.shapes.add_picture(crest, Inches(3.83), Inches(-0.05),
                                 Inches(5.68), Inches(1.24))

# ---------------------------------------------------------------------------
def add_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text

# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------
def blank(prs):
    """A brand-new blank slide on the template's default layout."""
    return prs.slides.add_slide(prs.slide_layouts[6])  # 'Blank'

def title_slide(prs, assets, title, subtitle, notes, jury=None, show_logos=True):
    s = blank(prs)
    if show_logos:
        # top corner logos (Kingdom crest left, ENSIAS right) + center crest
        add_chrome_header_logos(s, assets['crest_title'], assets['ensias_title'])
        # Also the top-left corner two logos (title slide uses Picture 4/5)
        s.shapes.add_picture(assets['title_logo_left'], Inches(0.74), Inches(0.05),
                             Inches(2.21), Inches(1.61))
        s.shapes.add_picture(assets['title_logo_right'], Inches(10.69), Inches(0.03),
                             Inches(2.32), Inches(1.64))
    # Banner
    add_text(s, Inches(1.81), Inches(1.21), Inches(10.04), Inches(0.30),
             "Doctoral Studies Center in Information and Engineering Sciences and Technologies (ST2I)",
             size=14, color=RGBColor(0x40,0x40,0x40), italic=True)
    # Big title
    add_rect(s, Inches(0.74), Inches(1.73), Inches(12.31), Inches(0.85))
    add_text(s, Inches(0.9), Inches(1.83), Inches(12.0), Inches(0.70), title,
             size=26, color=DARK, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(1.65), Inches(2.75), Inches(10.04), Inches(0.37), subtitle,
             size=18, color=BLUE, align=PP_ALIGN.CENTER, bold=True)
    if notes:
        add_notes(s, notes)
    return s

def section_slide(prs, assets, number, big, small, notes):
    s = blank(prs)
    # Center block with big title (like Contr. I header)
    add_rect(s, Inches(0.98), Inches(2.11), Inches(11.38), Inches(3.28), color=HDR_BG)
    add_text(s, Inches(2.2), Inches(2.6), Inches(9.5), Inches(1.2), big,
             size=36, color=DARK, bold=True, align=PP_ALIGN.LEFT)
    if small:
        add_text(s, Inches(2.2), Inches(4.4), Inches(9.5), Inches(0.6), small,
                 size=18, color=BLUE, italic=True, align=PP_ALIGN.LEFT)
    # big number chip
    if number:
        s.shapes.add_picture(assets['number_chip'], Inches(1.07), Inches(2.20),
                             Inches(0.86), Inches(0.86))
    # footer
    add_chrome_footer(s, number or 0, prs, assets['footer'], assets['pager_left'],
                      assets['pager_right'])
    add_notes(s, notes)
    return s

def content_slide(prs, assets, title, bullets, notes, subtitle=None, page=None,
                  layout="bullets", header="Introduction"):
    s = blank(prs)
    # top chrome: small header block + tab rectangles (mimic colleague)
    add_rect(s, Inches(4.80), Inches(0.10), Inches(8.45), Inches(0.53), color=HDR_BG)
    add_text(s, Inches(4.9), Inches(0.15), Inches(8.3), Inches(0.45), title,
             size=20, color=DARK, bold=True, align=PP_ALIGN.CENTER)
    if subtitle:
        add_rect(s, Inches(0.23), Inches(0.79), Inches(6.42), Inches(0.40), color=HDR_BG)
        add_text(s, Inches(0.3), Inches(0.82), Inches(6.3), Inches(0.36), subtitle,
                 size=14, color=BLUE, bold=True)
    # body
    if layout == "bullets":
        add_bullets(s, Inches(0.7), Inches(1.55), Inches(11.9), Inches(5.0), bullets,
                    size=17)
    elif layout == "two_col":
        left, right = bullets
        add_bullets(s, Inches(0.7), Inches(1.55), Inches(6.0), Inches(5.0), left, size=16)
        add_bullets(s, Inches(6.85), Inches(1.55), Inches(5.9), Inches(5.0), right, size=16)
    add_chrome_footer(s, page or 0, prs, assets['footer'], assets['pager_left'],
                      assets['pager_right'])
    add_notes(s, notes)
    return s

def table_slide(prs, assets, title, headers, rows, notes, subtitle=None, page=None,
                col_widths=None, font_size=12, highlight_rows=None):
    s = blank(prs)
    add_rect(s, Inches(4.80), Inches(0.10), Inches(8.45), Inches(0.53), color=HDR_BG)
    add_text(s, Inches(4.9), Inches(0.15), Inches(8.3), Inches(0.45), title,
             size=20, color=DARK, bold=True, align=PP_ALIGN.CENTER)
    if subtitle:
        add_rect(s, Inches(0.23), Inches(0.79), Inches(6.42), Inches(0.40), color=HDR_BG)
        add_text(s, Inches(0.3), Inches(0.82), Inches(6.3), Inches(0.36), subtitle,
                 size=14, color=BLUE, bold=True)
    ncols = len(headers); nrows = len(rows)+1
    left = Inches(0.7); top = Inches(1.7)
    width = Inches(11.9); row_h = Inches(0.40)
    shape = s.shapes.add_table(nrows, ncols, left, top, width, row_h*nrows)
    table = shape.table
    # remove default table style banding by setting a theme-ish, then custom fills
    if col_widths:
        total = sum(col_widths)
        for j,w in enumerate(col_widths):
            table.columns[j].width = Emu(int(width*(w/total)))
    for j,h in enumerate(headers):
        c = table.cell(0,j); c.text = h
        p = c.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.runs[0]; r.font.size = Pt(font_size); r.font.bold=True; r.font.name="Calibri"
        r.font.color.rgb = WHITE
        c.fill.solid(); c.fill.fore_color.rgb = BLUE
        c.vertical_anchor = MSO_ANCHOR.MIDDLE
    for i,row in enumerate(rows, start=1):
        hl = highlight_rows is not None and (i-1) in highlight_rows
        for j,val in enumerate(row):
            c = table.cell(i,j); c.text = str(val)
            p = c.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.LEFT if j==0 else PP_ALIGN.CENTER
            r = p.runs[0]; r.font.size = Pt(font_size); r.font.name="Calibri"
            r.font.color.rgb = WHITE if hl else GREY
            r.font.bold = hl
            if hl:
                c.fill.solid(); c.fill.fore_color.rgb = RGBColor(0x2E,0x5A,0x9E)
            elif i%2==0:
                c.fill.solid(); c.fill.fore_color.rgb = RGBColor(0xEF,0xF3,0xFA)
            else:
                c.fill.solid(); c.fill.fore_color.rgb = WHITE
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            c.margin_left = Inches(0.06); c.margin_right = Inches(0.06)
            c.margin_top = Inches(0.02); c.margin_bottom = Inches(0.02)
    add_chrome_footer(s, page or 0, prs, assets['footer'], assets['pager_left'],
                      assets['pager_right'])
    add_notes(s, notes)
    return s

def card_slide(prs, assets, title, cards, notes, subtitle=None, page=None, top=2.2):
    s = blank(prs)
    add_rect(s, Inches(4.80), Inches(0.10), Inches(8.45), Inches(0.53), color=HDR_BG)
    add_text(s, Inches(4.9), Inches(0.15), Inches(8.3), Inches(0.45), title,
             size=20, color=DARK, bold=True, align=PP_ALIGN.CENTER)
    if subtitle:
        add_rect(s, Inches(0.23), Inches(0.79), Inches(6.42), Inches(0.40), color=HDR_BG)
        add_text(s, Inches(0.3), Inches(0.82), Inches(6.3), Inches(0.36), subtitle,
                 size=14, color=BLUE, bold=True)
    n = len(cards); gap = Inches(0.35); total = Inches(11.9)
    cw = Emu(int((total - gap*(n-1))/n)); x = Inches(0.7); ch = Inches(3.0)
    for big,label,sub in cards:
        card = add_rect(s, x, Inches(top), cw, ch, color=RGBColor(0xF4,0xF7,0xFC))
        add_rect(s, x, Inches(top), cw, Inches(0.12), color=BLUE)
        add_text(s, x+Inches(0.15), Inches(top)+Inches(0.5), cw-Inches(0.3), Inches(0.9),
                 big, size=30, color=DARK, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, x+Inches(0.15), Inches(top)+Inches(1.5), cw-Inches(0.3), Inches(0.6),
                 label, size=15, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, x+Inches(0.15), Inches(top)+Inches(2.1), cw-Inches(0.3), Inches(0.9),
                 sub, size=11, color=GREY, align=PP_ALIGN.CENTER)
        x += cw + gap
    add_chrome_footer(s, page or 0, prs, assets['footer'], assets['pager_left'],
                      assets['pager_right'])
    add_notes(s, notes)
    return s

# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------
def main():
    ensure_template()
    prs = Presentation(SRC)

    # Extract reusable assets from the template deck
    assets = {}
    # Pull image blobs via navigating template slides (they're already in package)
    # We'll re-extract by position/name using the template's presentation object
    from pptx import Presentation as Pres
    tpl = Pres(SRC)
    def pic(slide_idx, name):
        for sh in tpl.slides[slide_idx].shapes:
            if sh.shape_type == 13 and sh.name == name:
                return sh.image
        return None
    def any_pic(slide_idx):
        for sh in tpl.slides[slide_idx].shapes:
            if sh.shape_type == 13:
                return sh.image
        return None
    # save asset blobs to temp files
    tmp = tempfile.mkdtemp()
    def blob_path(img, key):
        p = os.path.join(tmp, key + "." + img.ext)
        with open(p,"wb") as f: f.write(img.blob)
        return p
    assets['crest_title']  = blob_path(pic(0,"Picture 3"), "crest")
    assets['ensias_title'] = blob_path(pic(0,"Picture 5"), "ensias")
    assets['title_logo_left']  = blob_path(pic(0,"Picture 4"), "leftlogo")
    assets['title_logo_right'] = blob_path(pic(0,"Picture 5"), "rightlogo")
    assets['footer'] = blob_path(pic(16,"Picture 2"), "footer")
    assets['pager_left']  = blob_path(pic(16,"Picture 10"), "pagerL")
    assets['pager_right'] = blob_path(pic(16,"Picture 7"), "pagerR")
    assets['section_icon'] = blob_path(pic(7,"Picture 6"), "numchip")
    # number chip: we want a generic numeral; 'section_icon' is '2'. Use it as-is
    assets['number_chip'] = assets['section_icon']

    # ---- remove all existing slides from template ----
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    for sld in slides:
        rId = sld.get(qn('r:id'))
        prs.part.drop_rel(rId)
        xml_slides.remove(sld)

    # ---- now add our slides (blank layout) ----
    page = 0
    def bump():
        nonlocal page
        page += 1
        return page

    # 1 TITLE
    title_slide(prs, assets,
        "Cooperative Game Theory for Explainable AI in Recommendation Systems: A Shapley Framework for Actionable Insight",
        "PhD Viva Presented by: Mouad LOUHICHI \u2014 Supervisor: Pr. Mohamed LAZAAR | ENSIAS, Mohammed V University, Rabat",
        "Good morning. Thank you, President and Professors, for examining my thesis. I am Mouad Louhichi. "
        "My thesis is titled Cooperative Game Theory for Explainable AI in Recommendation Systems, a Shapley "
        "Framework for Actionable Insight, supervised by Professor Mohamed Lazaar.")
    bump()

    # 2 OUTLINE (replicate the colored box layout)
    def outline_slide():
        s = blank(prs)
        add_rect(s, Inches(0.00), Inches(0.00), Inches(13.33), Inches(1.06), color=HDR_BG)
        add_text(s, Inches(0.43), Inches(0.28), Inches(6), Inches(0.5), "VIVA Outline",
                 size=26, color=DARK, bold=True)
        boxes = [
            ("Introduction", ["Recommenders","Black-box problem","Why XAI"]),
            ("Context & Problematic", ["Limitations","Research gap"]),
            ("Protocols", ["Datasets","Baseline","Metrics","Hardware"]),
            ("Contributions", ["C1 Clustering XAI","C2 Multi-level XAI","C3 DyHuCoG"]),
            ("Conclusion & Perspectives", ["Synthesis","Limitations","Future work"]),
        ]
        col_x = [0.43, 3.03, 5.63, 8.21, 10.79]
        for (head, items), x in zip(boxes, col_x):
            add_rect(s, Inches(x), Inches(2.45), Inches(2.17), Inches(3.03),
                     color=RGBColor(0xF4,0xF7,0xFC), line=RGBColor(0xD8,0xE0,0xEC))
            add_text(s, Inches(x+0.15), Inches(2.55), Inches(1.9), Inches(1.0), head,
                     size=15, color=DARK, bold=True)
            add_bullets(s, Inches(x+0.15), Inches(3.35), Inches(1.9), Inches(2.0),
                        items, size=12, color=GREY, space=6)
        add_chrome_footer(s, 2, prs, assets['footer'], assets['pager_left'], assets['pager_right'])
        add_notes(s, "I will walk you through five parts: introduction, context and problematic, the "
                     "experimental protocol, then my three contributions, and finally the conclusion and perspectives. "
                     "Each contribution follows the same structure: objectives, methodology, results, and findings.")
    outline_slide(); bump()

    # 3 SECTION: Introduction
    section_slide(prs, assets, None, "Introduction", "Recommenders \u00b7 The black-box problem \u00b7 Why XAI",
        "Let us begin with the introduction."); bump()

    content_slide(prs, assets, "Motivation: Three Questions",
        [("Ubiquity \u2013 How do opaque AI systems shape what billions of users see, buy, and watch daily?",0),
         ("The black box \u2013 Why do state-of-the-art recommenders and clustering pipelines remain opaque to users and designers?",0),
         ("Toward trust \u2013 How can transparency be built as part of the model, not bolted on afterwards?",0),
         ("",0),
         ("Core tension: as models gain expressive power they lose transparency needed for trustworthy deployment.",1)],
        "Three questions frame the thesis. Ubiquity: opaque systems mediate what people see. The black box: "
        "even strong recommenders are hard to interrogate. Trust: transparency should be built into the "
        "modelling logic. The core tension is between predictive power and interpretability.",
        subtitle="Where the thesis begins", page=bump())

    content_slide(prs, assets, "Actionable Insight \u2014 the Definition",
        [("Definition 1.1 (Actionable insight)",0,True),
         ("An explanation is actionable when it identifies at least one modifiable factor whose change is associated with a specifiable change in model output, and that factor is accessible to the decision-maker.",1),
         ("Accessibility means expressed in the semantic vocabulary of the task domain:",0),
         ("a physicochemical variable (wine), a pollution indicator (air quality), or a preference signal (recommendation).",1),
         ("not an opaque latent code.",1)],
        "We frame the whole thesis around actionable insight: an explanation must point to a modifiable driver "
        "expressed in the language of the domain, so it supports intervention rather than merely describing.",
        subtitle="Actionability, not just plausibility", page=bump())

    content_slide(prs, assets, "Research Context",
        [("Recommenders evolved from similarity filters to complex representation-learning systems on sparse, high-dimensional, dynamic data.",0),
         ("MF \u2192 neural CF \u2192 graph CNNs \u2192 hypergraph: each step improved ranking but intensified the interpretability deficit.",0),
         ("Why the deficit matters:",0,True),
         ("Undermines user trust.",1),("Constrains debugging and scientific learning.",1),
         ("Collides with regulatory expectations (EU AI Act, OECD principles, GDPR).",1)],
        "The context is the progression from simple to hypergraph recommenders. Each step raised expressiveness "
        "while lowering transparency. The deficit matters for trust, debugging, and regulation.",
        subtitle="From similarities to hypergraph models", page=bump())

    # SECTION: Context & Problematic
    section_slide(prs, assets, None, "Context & Problematic", "Three structuring limitations \u00b7 The research gap",
        "Now let me look more precisely at the problem this thesis addresses."); bump()

    content_slide(prs, assets, "Recommendation & Clustering Paradigms",
        [("Collaborative filtering \u2013 users who behaved similarly will value similar items (user-/item-based).",0),
         ("Content-based \u2013 recommends items sharing attributes with a user profile.",0),
         ("Hybrid \u2013 combines collaborative and content signals.",0),
         ("Matrix factorisation \u2013 latent factors R \u2248 PQ\u1d40, compact but opaque.",0),
         ("Graph-based \u2013 interaction graph with neighbourhood propagation (LightGCN, hypergraph).",0)],
        "Quick orientation across paradigms we build on. Each strengthens modelling but complicates "
        "interpretation. Matrix factorisation made latent dimensions opaque; graph models kept importance "
        "implicit; hypergraph models added higher-order relations but assumed uniform message importance.",
        subtitle="The paradigm landscape", page=bump())

    content_slide(prs, assets, "Limitations of Classical Recommenders & Unsupervised Models",
        [("Data sparsity & scalability \u2013 the user\u2013item matrix is overwhelmingly empty.",0),
         ("Cold-start \u2013 new users and items are structurally disadvantaged.",0),
         ("Popularity bias & lack of diversity \u2013 exposure begets interaction, begets exposure (filter bubble).",0),
         ("Absence of interpretability \u2013 the most fundamental limit: weak explanatory traction.",0,True),
         ("For clustering: methods privilege local OR global explanation, not both; they struggle to scale; explanations rarely preserve coherence across resolutions.",1)],
        "Four classical limitations. The most fundamental is the absence of interpretability. In clustering "
        "specifically, approaches tend to privilege local or global explanation, struggle to scale, and rarely "
        "remain coherent across resolutions.",
        subtitle="Why the gap exists", page=bump())

    content_slide(prs, assets, "Three Structuring Limitations (Problem Statement)",
        [("1. Lack of explainability \u2013 complex models remain hard to interpret faithfully and actionably.",0,True),
         ("2. Difficulty of scaling \u2013 local explanations do not transfer naturally to hierarchical structures or large datasets.",0,True),
         ("3. Weak integration into learning \u2013 most explanations stay post-hoc and do not shape model dynamics or the accuracy-diversity-context trade-off.",0,True),
         ("",0),
         ("Thesis gap: the literature still lacks a single cooperative-attribution framework that can explain clustering faithfully, remain coherent under hierarchy, and then operate as an in-training signal in recommendation.",1)],
        "Three structuring problems: lack of explainability, poor scalability of local explanations, and weak "
        "integration into the learning loop. The gap is that no framework carries a single cooperative-attribution "
        "logic across clustering, hierarchy, and in-training recommendation.",
        subtitle="Problem statement", page=bump())

    content_slide(prs, assets, "Research Questions (RQ1\u2013RQ5) and Overall Aim",
        [("Aim: develop, justify, and evaluate a cooperative-game-theoretic perspective for XAI in clustering and recommendation, using Shapley attribution as both an explanatory mechanism and an in-training signal.",0,True),
         ("RQ1 \u00b7 How can Shapley values explain black-box clustering faithfully at instance and cluster level?",0),
         ("RQ2 \u00b7 How can this extend to large-scale, hierarchical clustering without losing tractability or consistency?",0),
         ("RQ3 \u00b7 Can cooperative attribution move beyond post-hoc and enter the learning dynamics of graph recommenders?",0),
         ("RQ4 \u00b7 Can a recommender jointly optimise ranking accuracy, context, and diversity when importance is estimated by a cooperative-game utility?",0),
         ("RQ5 \u00b7 What emerges when clustering explanation and recommendation learning are two stages of one cooperative-game perspective?",0)],
        "The aim is one cooperative-game perspective for XAI across clustering and recommendation. Five research "
        "questions form the spine: RQ1-RQ2 clustering, RQ3-RQ4 recommendation, RQ5 at thesis level.",
        subtitle="Research aim and questions", page=bump())

    content_slide(prs, assets, "The Three Contributions",
        [("C1 \u2013 Explainable black-box clustering: PCA\u2013K-Means\u2013LightGBM\u2013TreeSHAP pipeline.",0,True),
         ("\u2192 Wine Quality: faithful instance- and cluster-level feature attribution.",1),
         ("C2 \u2013 Enhanced multi-level XAI for large-scale clustering with cross-level SHAP aggregation.",0,True),
         ("\u2192 Beijing Air Quality: hierarchical attribution consistency under scale.",1),
         ("C3 \u2013 DyHuCoG: Dynamic Hypergraph Cooperative Game for preference-aware recommendation.",0,True),
         ("\u2192 MovieLens-1M & Amazon-Book: preference-aware Monte Carlo Shapley as in-training signal.",1),
         ("",0),
         ("Thesis claim: cooperative game theory functions as a shared attribution perspective for explanation, optimisation, and intervention.",2,True)],
        "Three contributions, one thread. C1 establishes Shapley-based explanation for black-box clustering. "
        "C2 scales it to hierarchy and large data. C3 replaces post-hoc attribution with an in-training signal "
        "in a hypergraph recommender. Together they support the thesis claim.",
        subtitle="The research arc", page=bump())

    # SECTION: Protocols
    section_slide(prs, assets, None, "Experimental Protocol", "Datasets \u00b7 Baselines \u00b7 Metrics \u00b7 Hardware",
        "Before the contributions, let me briefly cover the experimental setup shared across all three."); bump()

    table_slide(prs, assets, "Datasets Used Throughout the Thesis",
        ["Dataset","Scale","Type","Role"],
        [["Wine Quality (Portugal, vinho verde)","4,898 \u00d7 11","Tabular, numeric features","C1 \u2013 single-level clustering"],
         ["Beijing Multi-Site Air Quality","383,585 \u00d7 11","Tabular, pollutant + meteorology","C2 \u2013 multi-level clustering"],
         ["MovieLens-1M","6,040 u / 3,706 i / 1,000,209 int","Implicit-feedback","C3 \u2013 DyHuCoG"],
         ["Amazon-Book","52,643 u / 91,599 i / 2,984,108 int","Implicit-feedback (very sparse)","C3 \u2013 DyHuCoG"]],
        "Two clustering and two recommendation datasets. Clustering datasets are chosen for semantically "
        "interpretable features; recommendation datasets are standard benchmarks. Density: MovieLens-1M 0.0447, "
        "Amazon-Book 0.0006.",
        subtitle="Four datasets, two regimes", page=bump())

    content_slide(prs, assets, "Data Splitting & Preprocessing",
        [("Clustering: five-fold cross-validation for surrogate/attribution stability.",0),
         ("Recommendation: user-level, temporal split \u2014 70% train / 10% val / 20% test.",0),
         ("Leave-one-out: latest test positive per user is the target, ranked against negatives.",0),
         ("Implicit conversion: MovieLens-1M ratings > 3 treated as positive.",0),
         ("Popularity-aware negative sampling: q(i) \u221d f_i^\u03b7 for harder training contrasts.",0),
         ("Reproducibility: seeds {42,43,44,45,46}; early-stopping patience 20.",0)],
        "Splitting is designed against leakage. Clustering uses cross-validation; recommendation uses a temporal "
        "per-user holdout with leave-one-out. MovieLens ratings above 3 become positives, negatives are "
        "popularity-sampled, and five seeds are used for reproducibility.",
        subtitle="Leakage control and reproducibility", page=bump())

    content_slide(prs, assets, "Baselines & Evaluation Metrics",
        [("Clustering benchmarks: LIME-based surrogate explanation pipeline.",0,True),
         ("Recommendation benchmarks: MF, NCF, LightGCN, RecDCL, HCCF, HPCF (strongest reference).",0,True),
         ("",0),
         ("Ranking: Precision@K, Recall@K, NDCG@20 (principal).",0),
         ("System diversity: Catalogue Coverage = |\u22c3 R_u| / |I|.",0),
         ("List diversity: Intra-List Diversity (ILD) \u2014 average pairwise dissimilarity; built into the utility.",0),
         ("Clustering quality: Silhouette, Davies\u2013Bouldin.",0)],
        "Baselines span classical, neural, graph, and hypergraph methods to isolate the contribution of "
        "cooperative attribution. Metrics capture ranking quality plus coverage and ILD for recommendation, and "
        "Silhouette / Davies-Bouldin for clustering. ILD is part of the DyHuCoG coalition utility.",
        subtitle="Comparing under one evaluative frame", page=bump())

    content_slide(prs, assets, "Hardware & Software",
        [("CPU: Intel Core i9-14900K, 24 cores \u2014 clustering, preprocessing, data loading.",0),
         ("GPU: NVIDIA GeForce RTX 4090, 24 GB (Ada Lovelace) \u2014 DyHuCoG training & inference.",0),
         ("RAM: 48 GB; Storage: 2 TB SSD.",0),
         ("Python 3.8; scikit-learn (PCA, K-Means), LightGBM (surrogate), SHAP (TreeSHAP), PyTorch 2.0.1 (GNN/HGNN), NumPy/SciPy/pandas.",0),
         ("Altair for interactive SHAP visualisation; metrics at K \u2208 {5, 10, 20}.",0)],
        "Clustering runs on standard workstation CPU; DyHuCoG needs a GPU but stays within academic resources. "
        "The stack is scikit-learn + LightGBM + SHAP for clustering, PyTorch for the recommendation model.",
        subtitle="Reproducible environment", page=bump())

    # SECTION: Contribution I
    section_slide(prs, assets, "1", "Contribution I \u2014 Explainable Black-Box Clustering",
        "PCA \u2192 K-Means \u2192 LightGBM \u2192 SHAP \u00b7 Wine Quality \u00b7 RQ1",
        "Let us move to the first contribution: explaining black-box clustering with Shapley values."); bump()

    content_slide(prs, assets, "C1 \u00b7 Research Gap & Objectives",
        [("Gap: Shapley explanation is dominant in supervised tasks, but clustering remains comparatively under-explained.",0,True),
         ("Existing clustering-interpretability methods privilege local or global explanation, not both.",1),
         ("They often fail to scale or preserve coherence across clusters.",1),
         ("",0),
         ("Objectives: build a pipeline yielding cluster-level explanation while preserving feature-level attribution; preserve the semantics of the original feature space; justify why Shapley is better than LIME.",0,True)],
        "Clustering is a natural starting point: the model creates its own structure, so cluster meaning must be "
        "inferred after the fact. The gap is that explainable clustering is fragmented. Our objectives are a "
        "pipeline that explains a partition while preserving original feature semantics, and a case for Shapley "
        "over LIME.",
        subtitle="RQ1", page=bump())

    content_slide(prs, assets, "C1 \u00b7 Cooperative-Game Formulation for Clustering",
        [("Player set N = F \u2014 each feature is a player.",0,True),
         ("Value function v(S) = Silhouette( KMeans(X_S, k*) ) \u2014 how well data cluster using only features in S.",0),
         ("A feature's Shapley value = its expected marginal contribution to clustering quality over all coalition orders.",0),
         ("Why Silhouette: bounded, normalised, semantically intuitive. Alternatives (DB, CH) possible.",0),
         ("Direct evaluation for every coalition is intractable \u2192 we need a bridge.",0)],
        "We frame clustering as a cooperative game where features are players. The value function is the "
        "Silhouette of clustering on a feature subset. A feature gets high attribution when its presence "
        "consistently improves separation. But exact evaluation is combinatorial, so we need a bridge.",
        subtitle="Features as players", page=bump())

    content_slide(prs, assets, "C1 \u00b7 The LightGBM Bridge for Tractable Attribution",
        [("Once K-Means produces cluster labels, train a LightGBM multiclass surrogate to predict those labels from original features.",0),
         ("Apply TreeSHAP to the surrogate \u2014 exact, fast tree-based attribution in the original semantic feature space.",0),
         ("Why the bridge matters:",0,True),
         ("Direct TreeSHAP on K-Means is impossible (it explains tree models, not centroids).",1),
         ("Explaining the PCA representation would move attribution away from interpretable variables.",1),
         ("The surrogate preserves the chemistry/pollution vocabulary that makes the analysis actionable.",1),
         ("Validity condition: surrogate fidelity is high (macro-F1 \u2248 0.82).",1,True)],
        "The bridge is the heart of the method. We convert an unsupervised partition into a supervised prediction "
        "task by training a LightGBM classifier to reproduce the cluster labels, then run TreeSHAP on it. This keeps "
        "attribution in the original feature space. It is valid only if the surrogate is faithful.",
        subtitle="Unsupervised partition \u2192 supervised explanation", page=bump())

    content_slide(prs, assets, "C1 \u00b7 Proposed Pipeline",
        [("Stage 1 \u2013 PCA: stabilise geometry + visual diagnostic. NOT the explanatory space.",0),
         ("Stage 2 \u2013 K-Means++ with multi-criteria k selection (elbow, Silhouette, Davies\u2013Bouldin).",0),
         ("Stage 3 \u2013 LightGBM surrogate trained on original features to predict cluster labels.",0),
         ("Stage 4 \u2013 TreeSHAP attribution in the original feature space.",0),
         ("Stage 5 \u2013 Aggregate into global importance, cluster-specific profiles, and local force plots.",0),
         ("",0),
         ("Complexity: dominated by PCA and repeated K-Means; TreeSHAP scales with tree count/depth, not exponentially in features.",1)],
        "The five-stage pipeline. PCA is a computational and visual aid, not the explanatory space. K-Means defines "
        "the partition, the surrogate restores tractable supervised prediction, and TreeSHAP returns explanation to "
        "the original variables.",
        subtitle="Five stages", page=bump())

    content_slide(prs, assets, "C1 \u00b7 Optimal Cluster Selection \u2014 a Deliberate Choice",
        [("Multi-criteria evaluation across k \u2208 {2..10} using elbow, Silhouette, Davies\u2013Bouldin.",0),
         ("We select k* = 3 \u2014 even though it is NOT geometrically optimal:",0,True),
         ("k = 2: Silhouette 0.214, Davies\u2013Bouldin 1.775 (better separation).",1),
         ("k = 3: Silhouette 0.144, Davies\u2013Bouldin 2.097 (weaker separation).",1),
         ("Why: three clusters give a semantically richer oenological partition \u2192 more discriminative, more actionable.",1,True),
         ("Note: the higher Silhouette \u2248 0.63 belongs to Beijing (C2), not the wine partition.",1)],
        "An honest, important point. We select three clusters even though two gives better raw geometry metrics. The "
        "justification is interpretability: three clusters support three distinct, chemically meaningful narratives. "
        "The 0.63 Silhouette belongs to Beijing, not the retained wine partition.",
        subtitle="Interpretability over geometry", page=bump())

    content_slide(prs, assets, "C1 \u00b7 Global Feature Importance (Wine)",
        [("Global SHAP ranking (high \u2192 low):",0,True),
         ("density \u2192 pH \u2192 fixed acidity \u2192 sulfur-dioxide \u2192 alcohol",1),
         ("Dominant drivers are intimately related to structure, preservation, and sensory balance.",1),
         ("This is NOT an arbitrary classifier artefact \u2014 it recovers a chemically interpretable hierarchy.",1)],
        "The global SHAP ranking is dominated by density, pH, fixed acidity, sulfur dioxide, alcohol \u2014 "
        "consistent with oenological knowledge. This shows the surrogate-based pipeline recovers a chemically "
        "interpretable hierarchy rather than fitting arbitrary artefacts.",
        subtitle="A chemistry-consistent hierarchy", page=bump())

    content_slide(prs, assets, "C1 \u00b7 Cluster-Specific Profiles (Three Signatures)",
        [("Three clusters show distinct explanatory signatures.",0,True),
         ("Cluster 0 \u2013 density + sulfur-dioxide-related variables.",1),
         ("Cluster 1 \u2013 acidity and pH-related effects.",1),
         ("Cluster 2 \u2013 a different balance of acidity, alcohol, and related chemical attributes.",1),
         ("The same small set of variables recurs across clusters, with different relative weights within each.",1)],
        "The cluster-specific profiles show the retained solution is not only globally interpretable but internally "
        "differentiated. Each cluster has its own chemical signature \u2014 exactly the actionable insight we seek.",
        subtitle="Cluster-level heterogeneity", page=bump())

    content_slide(prs, assets, "C1 \u00b7 SHAP vs. LIME \u2014 Why Cooperative Attribution",
        [("SHAP grounds attribution in a cooperative-game allocation rule; LIME fits a local surrogate.",0),
         ("Stability: SHAP more stable when surrogate faithful; LIME depends on perturbation design.",0),
         ("Guarantees: SHAP satisfies efficiency, symmetry, null player, additivity \u2014 LIME has no equivalent.",0),
         ("Local/global coherence: SHAP supports both; LIME primarily local.",0),
         ("Cluster comparison: SHAP strong; LIME limited.",0),
         ("Caveat: in the surrogate pipeline, efficiency holds w.r.t. the LightGBM output, not the Silhouette value directly.",1)],
        "Why Shapley over LIME? The four axioms give SHAP a normative basis: efficiency (completeness), symmetry "
        "(fairness), null player (no spurious credit), additivity (composition). LIME's local surrogates are "
        "perturbation-sensitive. One caveat: efficiency holds w.r.t. the surrogate, not the Silhouette value.",
        subtitle="Axiomatic attribution vs local surrogate", page=bump())

    content_slide(prs, assets, "C1 \u00b7 Findings & Limitations",
        [("Achieved: a cluster-level explanation anchored to individual feature contributions; explanations returned to original variables; theoretically-grounded case for Shapley over LIME.",0,True),
         ("Limitations: fidelity depends on the LightGBM surrogate; tabular data only; single-level structure \u2014 cannot yet address hierarchical coherence. That is C2\u2019s point of departure.",0,True)],
        "C1 answers RQ1: Shapley values can explain black-box clustering faithfully to a high-fidelity surrogate and "
        "coherently at cluster level. The LIME comparison is theoretical and literature-backed. Limits: surrogate "
        "dependence, tabular-only, single-level structure \u2014 which C2 tackles.",
        subtitle="RQ1 answered", page=bump())

    # SECTION: Contribution II
    section_slide(prs, assets, "2", "Contribution II \u2014 Enhanced Multi-Level XAI for Large-Scale Clustering",
        "Beijing Air Quality \u00b7 Hierarchical attribution consistency \u00b7 RQ2",
        "This brings us to the second contribution: scaling the explanation logic to multi-level, large-scale clustering."); bump()

    content_slide(prs, assets, "C2 \u00b7 Research Gap & Objectives",
        [("Gap: once clustering is multi-level, feature importance must stay interpretable within a cluster, across sub-clusters, and across the hierarchy as a whole.",0,True),
         ("Large-scale data make exact explanation computationally burdensome.",1),
         ("Flat explanation may be true yet incomplete \u2014 it cannot show how importance changes inside a cluster.",1),
         ("Objectives: a genuinely multi-level workflow (not a rerun of C1); a formal cross-level consistency argument (Prop. 6.1); validation on a structurally different large-scale dataset.",0,True)],
        "C2 asks whether the C1 logic survives scale and hierarchy. Large real-world data contain structure at more "
        "than one granularity. A flat explanation cannot show how importance reconfigures as you zoom in. Our "
        "objectives are a real multi-level workflow, a formal consistency argument, and validation on Beijing.",
        subtitle="RQ2", page=bump())

    content_slide(prs, assets, "C2 \u00b7 Multi-Level Clustering Architecture",
        [("Recursive/nested: coarse clustering on the full dataset, then subdivide each cluster where appropriate.",0),
         ("For each level, train a level-specific surrogate and compute SHAP in the SAME original feature space.",0),
         ("Cross-level aggregation is NOT a naive average \u2014 it respects cluster size and nesting structure.",0,True),
         ("Parent-level attribution = an expectation over the explanatory structure of its descendants.",1),
         ("The hierarchy is a pragmatic analytical device, not a claim of true ontological hierarchy.",1)],
        "The multi-level architecture proceeds recursively. A coarse clustering is learned, then each cluster is "
        "subdivided. Each level gets its own surrogate and SHAP values in the same feature space. The hierarchy is an "
        "analytical device, not a metaphysical claim.",
        subtitle="Nested clustering as an analytical tool", page=bump())

    content_slide(prs, assets, "C2 \u00b7 Formal Result: Hierarchical Attribution Consistency (Prop. 6.1)",
        [("Let \u03a6^(l,c)_j = E_{x~c}[ |\u03c6_j^(l)(x)| ] be the expected absolute SHAP importance of feature j at level l in cluster c.",0),
         ("Let w_c' = |c'| / |c| be the relative size of child c' within parent c.",0),
         ("For a strict nested hierarchy on a consistent feature space:",0,True),
         ("\u03a6^(l,c)_j = \u03a3_{c'\u2208child(c)} w_c' \u00b7 \u03a6^(l+1,c')_j + \u03b5_j",1,True),
         ("\u03b5_j is a residual from surrogate mismatch, vanishing under perfect fidelity.",1),
         ("Derived via law of total expectation (children partition the parent).",1),
         ("Does NOT imply explanations are identical across levels \u2014 it implies differences can be interpreted, not dismissed as inconsistency.",1)],
        "Proposition 6.1 is the chapter's original formalisation. A parent's expected absolute importance is the "
        "size-weighted average of its children's, up to a surrogate residual. It follows from the law of total "
        "expectation. It does not claim explanations are identical across levels; it lets us interpret differences.",
        subtitle="The cross-level consistency claim", page=bump())

    content_slide(prs, assets, "C2 \u00b7 Results: Cluster Quality (Beijing)",
        [("Full dataset, k = 3 (strong convergence on multi-criteria evaluation).",0,True),
         ("Silhouette \u2248 0.63 \u2014 materially stronger separation than wine.",1),
         ("Davies\u2013Bouldin \u2248 0.55 \u2014 low between-cluster ambiguity.",1),
         ("PCA projection (2 components) used only for visual inspection.",1),
         ("Sensitivity: robust to modest variation in k, projection dim, surrogate depth; only low-ranked variables shift.",1)],
        "On Beijing, the multi-criteria evaluation converges much more strongly: three clusters with Silhouette "
        "about 0.63 and Davies-Bouldin 0.55 \u2014 a markedly clearer separation. Conclusions are robust to reasonable "
        "parameter changes.",
        subtitle="Much clearer separation than wine", page=bump())

    content_slide(prs, assets, "C2 \u00b7 Global Feature Importance (Beijing)",
        [("Global SHAP ranking (high \u2192 low):",0,True),
         ("temperature \u2192 dew point \u2192 pressure \u2192 CO \u2192 NO2 \u2192 PM10 \u2192 PM2.5",1),
         ("It is NOT simply pollutant concentrations that matter \u2014 meteorological variables play a structurally central role.",1),
         ("Temperature, dew point, pressure condition dispersion, trapping, and photochemical behaviour.",1),
         ("This is the kind of insight flat descriptive summaries often fail to make explicit.",1)],
        "The global ranking is analytically rich: temperature, dew point, pressure dominate, then CO, NO2, PM10, "
        "PM2.5. Meteorological variables are structurally central because they condition dispersion and photochemistry.",
        subtitle="Meteorology conditions pollution regimes", page=bump())

    content_slide(prs, assets, "C2 \u00b7 Three Pollution Regimes (Force Plots)",
        [("Regime A \u2013 warm photochemical: ozone, temperature, dew point prominent (summer photochemical smog).",0),
         ("Regime B \u2013 wintertime smog: CO, SO2, PM dominate; low wind speed suppresses dispersion.",0),
         ("Regime C \u2013 comparatively clean air: favourable meteorology, weak pollutant pushes.",0),
         ("The framework shows not only that these regimes exist, but which variable combinations define them.",1)],
        "The force plots reveal three representative regimes. The interpretative value is showing not only that these "
        "regimes exist, but which variable combinations define each \u2014 the actionable insight.",
        subtitle="Three interpretable air-quality regimes", page=bump())

    content_slide(prs, assets, "C2 \u00b7 The Multi-Level Insight",
        [("At the coarse level, temperature and dew point dominate \u2014 they differentiate broad atmospheric regimes.",0),
         ("Within individual clusters, CO, SO2, PM10, wind speed, pressure, or ozone become more discriminative.",0),
         ("This change is NOT contradictory \u2014 it is exactly what a multi-level explanation should reveal.",0,True),
         ("Parent-level story = regime selection. Cluster-level story = variation within a regime.",1),
         ("A variable can be globally important yet locally uninformative within a sub-cluster.",1)],
        "This is the conceptual payoff. Coarse level = regime selection (temperature, dew point); within clusters the "
        "discriminative variables shift to CO, SO2, PM10, wind speed. That is not inconsistency \u2014 it is what "
        "multi-level explanation is for.",
        subtitle="Why flat explanation is insufficient", page=bump())

    content_slide(prs, assets, "C2 \u00b7 Cross-Dataset Generalisation",
        [("Wine: small, dense, chemically correlated. Beijing: large, noisy, temporally and meteorologically variable.",0),
         ("The same explanatory logic remains productive in both \u2192 not tied to one domain-specific peculiarity.",0,True),
         ("vs SHAP-based clustering literature: Beijing Silhouette \u2248 0.63 vs Gramegna & Giudici credit-risk 0.37.",1),
         ("LIME comparator: weaker structural coherence, less stable local narratives for hierarchical reasoning.",1)],
        "The same logic works on both a small chemical dataset and a large noisy environmental one, supporting "
        "generality. Compared with prior SHAP-based clustering work, our Beijing partition is comparatively well "
        "separated.",
        subtitle="Generalisable beyond one domain", page=bump())

    content_slide(prs, assets, "C2 \u00b7 Findings & Limitations",
        [("Achieved: scalable, multi-granular explanation without collapsing into a single flat summary; formal consistency argument; validation on a structurally different large-scale dataset.",0,True),
         ("Limitations: still static clustering (despite temporal Beijing data); surrogate-based SHAP + representative-instance reporting compress observation-level variation; tabular only.",0,True)],
        "C2 answers RQ2: Shapley-based clustering explanation can scale to hierarchical, large-scale settings without "
        "losing interpretive coherence, provided the hierarchy is modelled explicitly and approximation is transparent. "
        "The main limit is that clustering is still static \u2014 the contrast with C3.",
        subtitle="RQ2 answered", page=bump())

    # SECTION: Contribution III
    section_slide(prs, assets, "3", "Contribution III \u2014 DyHuCoG: A Dynamic Hypergraph Cooperative Game",
        "Preference-aware Shapley in hypergraph message passing \u00b7 RQ3 & RQ4",
        "The third and principal contribution introduces DyHuCoG, where Shapley attribution becomes an in-training "
        "signal inside a hypergraph recommender."); bump()

    content_slide(prs, assets, "C3 \u00b7 Research Gap & Objectives",
        [("Gap: graph and hypergraph recommenders treat message importance as either uniform or attention-weighted, without a principled marginal-contribution account.",0,True),
         ("Diversity is often a secondary objective or a re-ranking heuristic.",1),
         ("Interpretability is added after prediction, not integrated into the learning objective.",1),
         ("Objectives: formulate recommendation as a cooperative game; embed preference-aware Monte Carlo Shapley into message passing; improve ranking, coverage, and intra-list diversity jointly.",0,True)],
        "C3 is the flagship. The gap: hypergraph recommenders assume importance is implicit, diversity is secondary, "
        "and interpretability is bolted on. Our objectives: formulate recommendation as a cooperative game, inject "
        "preference-aware Shapley into message passing, and jointly improve ranking, coverage, and diversity.",
        subtitle="RQ3 & RQ4", page=bump())

    content_slide(prs, assets, "C3 \u00b7 Problem Formulation \u2014 Recommendation as a Cooperative Game",
        [("Player set N = U \u222a I \u222a C (users, items, contexts).",0,True),
         ("Hypergraph H = (V, E, W); V = U \u222a I \u222a C; W = dynamic edge weights from Shapley estimates.",0),
         ("Coalition S \u2286 N represents entities participating in a recommendation episode.",0),
         ("Coalition value v(S) measures the quality of the recommendation outcome achievable by S.",0),
         ("Top-N task: produce a ranked list L_u balancing relevance, diversity, and contextual fit.",0)],
        "We model recommendation as a cooperative game whose players are users, items, and contexts. The hypergraph "
        "encodes user-item-context interactions as hyperedges. A coalition value measures how good the "
        "recommendation is for an episode. This parallels the clustering formulation with a recommendation-oriented "
        "value function.",
        subtitle="Users, items, contexts as players", page=bump())

    content_slide(prs, assets, "C3 \u00b7 Multi-Objective Coalition Utility",
        [("v(S) = \u03b1 \u00b7 NDCG@20(S) + \u03b2 \u00b7 Diversity(S) + \u03b3 \u00b7 ContextScore(S)",0,True),
         ("with \u03b1 + \u03b2 + \u03b3 = 1. The same trade-off the recommender must satisfy is the trade-off from which attribution is computed \u2014 explanatory game and predictive objective are aligned by design.",1),
         ("Preference-weighted: v_pref(S) = v(S) + \u03bb_pref \u00b7 \u03a3_{(u,i)\u2208S} sim(u,i)",0,True),
         ("\u03b1 = 0.60, \u03b2 = 0.25, \u03b3 = 0.15; \u03bb_pref = 0.20 \u2014 grid-searched, stable (<1.5% variance in NDCG@20).",1),
         ("Coalition evaluation scoped to the interaction episode (a few dozen players), not the full catalogue.",1)],
        "The coalition utility combines ranking quality, diversity, and context. The same trade-off the model must "
        "satisfy is the trade-off from which attribution is computed, so the explanatory and predictive objectives are "
        "aligned by design. We add a preference bonus for user-item agreement.",
        subtitle="The utility the game is built on", page=bump())

    content_slide(prs, assets, "C3 \u00b7 Preference-Aware Monte Carlo Shapley Estimation",
        [("Exact Shapley is combinatorial and infeasible for realistic systems.",0),
         ("Monte Carlo estimator: \u03c6\u0302_j = (1/M) \u03a3_m [ v(S_m \u222a {j}) \u2212 v(S_m) ]",0,True),
         ("Preference-aware: \u03c6\u0302_j^pref = (1/M) \u03a3_m [ v_pref(S_m \u222a {j}) \u2212 v_pref(S_m) ]",0,True),
         ("Unbiased; variance = \u03c3\u00b2/M \u2192 MSE decays O(1/M), absolute error O(1/\u221aM).",1),
         ("M = 50 selected: MSE \u2248 1.4\u00d710\u207b\u2075, ~99% accuracy on MovieLens-1M.",1,True),
         ("Refreshed every 10 batches (~49 updates/epoch), smoothed by exponential moving average.",1)],
        "Exact computation is infeasible, so we use a Monte Carlo estimator and its preference-aware variant. It is "
        "unbiased with variance decaying as 1/M. We choose M = 50 as the practical balance, refreshing every ten "
        "batches with exponential smoothing so attribution is adaptive but training stays stable.",
        subtitle="Dynamic attribution, not a static diagnostic", page=bump())

    content_slide(prs, assets, "C3 \u00b7 Architecture: Shapley-Weighted Hypergraph Message Passing",
        [("Base propagation: e^(l+1) = \u03c3( D^-1/2 A D^-1/2 e^(l) )",0),
         ("Shapley-weighted: e_j^(l+1) = \u03c3( W^(l) e_j^(l) + \u03a3_{k\u2208N(j)} w_jk e_k^(l) )",0,True),
         ("Normalised weights: w_jk = \u03c6\u0302_jk / \u03a3_{k'\u2208N(j)} \u03c6\u0302_jk'",0,True),
         ("Clipped + exponentially smoothed before normalisation (stabilises sparse regimes).",1),
         ("Attention gate: a_ui = \u03c3( W_a[ e_u, e_i, l_i ] ); y_ui = (1 + a_ui) \u27e8e_u, e_i\u27e9.",0),
         ("Context-aware score: f(u,i,c) = y_ui + \u03bb_c \u27e8g(c_ui), e_cui\u27e9.",0)],
        "The architecture is the decisive move. Propagation is standard hypergraph message passing, but messages are "
        "weighted by normalised Shapley coefficients, so the model is told not only who is connected to whom but how "
        "much each coalition is worth. An attention gate stabilises by interpolating between Shapley-weighted and "
        "uniform propagation.",
        subtitle="Relational structure filtered through cooperative importance", page=bump())

    content_slide(prs, assets, "C3 \u00b7 Multi-Objective Training",
        [("L = L_rec + \u03bb_div L_div + \u03bb_ctx L_ctx + \u03bb_reg L_reg",0,True),
         ("L_rec \u2013 Bayesian Personalised Ranking (pairwise, implicit feedback).",1),
         ("L_div \u2013 Intra-List Diversity regulariser: penalises redundant ranked lists.",1),
         ("L_ctx \u2013 Context alignment: match context embedding to context-node representation.",1),
         ("L_reg \u2013 L2 weight decay.",1),
         ("The learning objective and coalition value are aligned: DyHuCoG trains to optimise the same balance that later determines attribution.",1,True)],
        "The composite loss combines BPR ranking with explicit diversity and context regularisation plus weight decay. "
        "The conceptual point is alignment: the model is trained to optimise the same balance that later defines "
        "cooperative attribution.",
        subtitle="Accuracy, diversity, context trained together", page=bump())

    table_slide(prs, assets, "C3 \u00b7 Main Results (MovieLens-1M & Amazon-Book)",
        ["Dataset","Model","NDCG@20","Recall@20","Coverage","Diversity"],
        [["MovieLens-1M","HPCF","0.2528","0.2098","0.342","0.461"],
         ["MovieLens-1M","DyHuCoG","0.2775","0.2362","0.397","0.516"],
         ["Amazon-Book","HPCF","0.0270","0.0359","0.259","0.535"],
         ["Amazon-Book","DyHuCoG","0.0306","0.0417","0.336","0.602"]],
        "The headline result. On MovieLens-1M, DyHuCoG beats the strongest baseline HPCF by +9.77% in NDCG@20 and "
        "+12.58% in Recall@20, and improves coverage and diversity. On Amazon-Book the gains are larger: +13.33% "
        "NDCG@20 and +16.16% Recall@20, coverage up 29.7%.",
        subtitle="vs the strongest baseline HPCF", page=bump())

    card_slide(prs, assets, "C3 \u00b7 Relative Gains over HPCF",
        [("+9.77%","NDCG@20","MovieLens-1M"),("+12.58%","Recall@20","MovieLens-1M"),
         ("+13.33%","NDCG@20","Amazon-Book"),("+16.16%","Recall@20","Amazon-Book")],
        "The most important message: DyHuCoG improves ranking accuracy AND coverage AND diversity simultaneously, with "
        "the largest relative gains on the sparser dataset. This is evidence the accuracy-diversity trade-off is not "
        "structurally fixed.",
        subtitle="Accuracy + coverage + diversity together", page=bump())

    content_slide(prs, assets, "C3 \u00b7 Coverage & Intra-List Diversity",
        [("MovieLens-1M: Coverage 0.342 \u2192 0.397 (+16.1%); ILD 0.461 \u2192 0.516 (+11.9%).",0,True),
         ("Amazon-Book: Coverage 0.259 \u2192 0.336 (+29.7%); ILD 0.535 \u2192 0.602 (+12.5%).",0,True),
         ("Reduced filter-bubble effect and greater discovery opportunity \u2014 while NDCG/Recall also improve, so accuracy is not traded off for diversity.",1)],
        "Coverage and intra-list diversity both improve. More of the catalogue is surfaced and items are less "
        "redundant, while ranking quality also improves \u2014 we are not trading accuracy for diversity.",
        subtitle="Both system- and list-level diversity improve", page=bump())

    table_slide(prs, assets, "C3 \u00b7 Ablation Study (Component-wise)",
        ["Variant","ML-1M NDCG@20","% Drop","Amazon NDCG@20","% Drop"],
        [["Full DyHuCoG","0.2775","\u2013","0.0306","\u2013"],
         ["w/o Shapley Value","0.2647","4.6%","0.0287","6.1%"],
         ["w/o Hypergraph","0.2586","6.8%","0.0279","8.9%"],
         ["w/o Attention","0.2678","3.5%","0.0295","3.5%"],
         ["w/o Context","0.2547","8.2%","0.0272","11.0%"],
         ["w/o Diversity","0.2614","5.8%","0.0288","5.8%"]],
        "Every component contributes. Removing context causes the largest loss, because context provides the "
        "representational substrate on which Shapley weighting operates. Removing Shapley weighting degrades "
        "performance noticeably, supporting the argument that marginal-contribution estimation is not decorative.",
        subtitle="Each block is essential", page=bump(), highlight_rows=[1])

    content_slide(prs, assets, "C3 \u00b7 Computational Efficiency & Shapley Convergence",
        [("Training: DyHuCoG ~2000 s vs HPCF ~1125 s on MovieLens-1M (\u2248 1.78\u00d7).",0),
         ("Inference: 1.84 ms/query (ML-1M), 8.52 ms (Amazon) \u2014 suitable for real-time deployment.",0),
         ("Memory: 4.4 vs 4.1 GB (ML-1M); 17.9 vs 16.8 GB (Amazon).",0),
         ("Per-epoch cost: O((L+1)md) + O((M/f)m).",0),
         ("Shapley convergence: M=50 \u2192 MSE 1.4\u00d710\u207b\u2075, ~99% accuracy; M=100 \u2192 MSE 3.5\u00d710\u207b\u2076 (diminishing returns).",0,True)],
        "The cost is bounded. Training is about 1.78 times HPCF, inference stays under a couple of milliseconds, and "
        "memory is modestly higher. The Monte Carlo budget at M = 50 gives ~99% accuracy at acceptable cost.",
        subtitle="The attribution cost is proportionate", page=bump())

    content_slide(prs, assets, "C3 \u00b7 Statistical Validation (MovieLens-1M)",
        [("Paired t-tests on per-user NDCG@20 (n = 6,040 users; df = 6,039).",0),
         ("DyHuCoG outperforms every baseline with extremely small p-values after Holm\u2013Bonferroni correction.",0),
         ("vs HPCF: t = 46.38, Cohen\u2019s d_z = 1.3345, p = 1.81\u00d710\u207b\u00b2\u2077\u2070.",1,True),
         ("Wilcoxon signed-rank test also significant (p < 0.001).",1),
         ("Effect sizes are large \u2014 improvements are substantively meaningful, not merely statistically visible.",1)],
        "The improvements are statistically significant and large. Paired tests against HPCF give t = 46.38 with "
        "Cohen's d_z = 1.33, all surviving Holm-Bonferroni. Effect sizes are large, which matters because large user "
        "counts can make tiny differences appear significant.",
        subtitle="Significant AND substantively meaningful", page=bump())

    content_slide(prs, assets, "C3 \u00b7 Cold-Start Robustness & Interpretability",
        [("Cold-start (5 or fewer interactions): NDCG@20 \u2248 0.061 (user) and 0.057 (item), improving over HPCF by ~10%.",0),
         ("Cross-dataset: MovieLens +9.9%, Amazon +14.8%, Yelp2018 +11.8%.",0),
         ("Interpretability: a SHAP waterfall decomposes a recommendation into ranking, diversity, context, and preference contributions \u2014 structurally tied to the components that drove training.",0,True),
         ("Popularity bias: Shapley measures marginal utility, not raw frequency \u2014 weak but informative interactions retain influence.",1)],
        "DyHuCoG also improves the regimes where recommenders are usually most brittle: cold-start users and items. "
        "Interpretability is direct: a waterfall plot decomposes a recommendation into the same utility components "
        "used during training, so the explanation is structurally faithful.",
        subtitle="Robust in sparse regimes, and interpretable", page=bump())

    content_slide(prs, assets, "C3 \u00b7 Findings & Limitations",
        [("Achieved: cooperative attribution used as an in-training signal \u2014 a stronger claim than the clustering chapters; the accuracy\u2013diversity trade-off is not structurally fixed.",0,True),
         ("Limitations: measurable computational overhead; depends on availability of meaningful context; Monte Carlo Shapley could be improved by variance reduction; ablation is component-wise; baselines finalised through early 2026.",0,True)],
        "C3 answers RQ3 and RQ4: cooperative attribution can move into the learning dynamics of a recommender, and "
        "doing so lets the model balance accuracy, diversity, and context more effectively. Limits: overhead, context "
        "dependence, no factorial ablation, baseline set capped at early 2026.",
        subtitle="RQ3 & RQ4 answered", page=bump())

    # SECTION: Conclusion
    section_slide(prs, assets, None, "Conclusion & Perspectives", "Synthesis \u00b7 Limitations \u00b7 Future work",
        "Let me now bring everything together."); bump()

    table_slide(prs, assets, "Thesis Synthesis",
        ["Contribution","Main idea","Achievement","Key finding"],
        [["C1","Explain black-box clustering via Shapley","PCA\u2013KMeans\u2013LightGBM\u2013TreeSHAP pipeline","Faithful, chemistry-consistent cluster attribution (wine)"],
         ["C2","Multi-level, large-scale clustering XAI","Cross-level SHAP aggregation + Prop. 6.1","Coherent under hierarchy; interprets differences, not inconsistency"],
         ["C3","DyHuCoG hypergraph cooperative game","Preference-aware Shapley as in-training signal","Accuracy + coverage + diversity improve together (ML & Amazon)"]],
        "Three contributions, one thread. C1 makes hidden structure intelligible; C2 keeps it coherent under scale and "
        "hierarchy; C3 carries the same attribution logic inside the learning dynamics of a recommender.",
        subtitle="The common thread", page=bump())

    table_slide(prs, assets, "Published Research Papers",
        ["No.","Title","Venue","Status"],
        [["I","Shapley Values for Explaining the Black Box Nature of Machine Learning Model Clustering","Procedia Computer Science 220, 806\u2013811","Published, 2023"],
         ["II","Game Theory Meets Explainable AI: An Enhanced Approach to Understanding Black Box Models Through Shapley Values","IJACSA 16(7), 716\u2013725","Published, 2025"],
         ["III","DyHuCoG: A Dynamic Hypergraph Cooperative Game for Preference-aware Recommendation","IJIES 19(2), 887\u2013902","Published, 2026"]],
        "The thesis synthesises three peer-reviewed publications. Paper I is Chapter 5, Paper II is Chapter 6, Paper "
        "III is Chapter 7. The thesis adds the multi-level formalisation and thesis-level synthesis.",
        subtitle="Three publications \u2192 Chapters 5\u20137", page=bump())

    content_slide(prs, assets, "Thesis-Level Limitations",
        [("Computational \u2013 exact Shapley is intractable; every contribution relies on approximation, surrogates, or restricted reporting.",0),
         ("Methodological \u2013 clustering depends on surrogate fidelity; recommendation depends on stable approximate contributions and adequate context.",0),
         ("Empirical \u2013 tabular clustering + benchmark recommendation; no multimodal, sequential, or online deployment; no dedicated human-subject actionability study.",0),
         ("Claim scope \u2013 a coherent and productive perspective, not one fully unified framework eliminating all tension.",0)],
        "Honestly stated limits. Computationally we approximate everywhere. Methodologically, clustering depends on "
        "surrogate fidelity and recommendation on stable approximate contributions. Empirically, the work is tabular "
        "and offline. The claim is a shared perspective, not one grand unified framework.",
        subtitle="Scope is deliberate", page=bump())

    content_slide(prs, assets, "Future Research Directions",
        [("Scalable cooperative attribution \u2013 lower-variance Shapley, learned proposal distributions, adaptive refresh policies.",0),
         ("Online / streaming recommendation \u2013 truly incremental settings with continuously evolving graphs and delayed feedback.",0),
         ("Richer human-centred evaluation \u2013 do explanations measurably improve analyst judgement, user trust, intervention quality, or perceived fairness?",0),
         ("Broader trustworthy-AI evaluation \u2013 exposure fairness, transparency requirements, governance-oriented auditing.",0)],
        "Future work turns the limitations into an agenda: more scalable cooperative attribution, then online and "
        "streaming settings, richer human-centred evaluation of actionability, and broader trustworthy-AI and "
        "fairness auditing.",
        subtitle="Where the work goes next", page=bump())

    content_slide(prs, assets, "Thesis Answer & Key Outcomes",
        [("Thesis answer: cooperative game theory can function as a shared methodological perspective for actionable explanation across clustering and recommendation.",0,True),
         ("Key outcomes: Shapley attribution as a common formal language for feature, interaction, and context importance allocation.",1),
         ("Faithful clustering explanation, hierarchical explanatory coherence, and contribution-aware recommendation learning.",1),
         ("Explanation as method, not commentary \u2014 from post-hoc description to in-training guidance.",1),
         ("Aligned with trustworthy-AI requirements (EU AI Act, OECD principles, GDPR).",1)],
        "The thesis answer: cooperative game theory is a shared perspective for explanation, optimisation, and "
        "intervention. This moves explanation from commentary to method, aligning with the transparency and "
        "accountability requirements of emerging AI regulation.",
        subtitle="The central claim, stated", page=bump())

    content_slide(prs, assets, "Questions & Discussion",
        [("Thank you for your attention.",0,True),("I welcome your questions and comments.",0)],
        "Thank you very much for your attention. I am now happy to take your questions.",
        subtitle="Jury discussion", page=bump())

    prs.save(OUT)
    print(f"Saved {len(prs.slides)} slides -> {OUT}")

if __name__ == "__main__":
    main()
