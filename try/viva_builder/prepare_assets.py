"""
One-off asset extraction for build.py.

  python3 prepare_assets.py <template.pptx> <original_viva.pptx> <out_dir>

Creates <out_dir>/fonts  (embedded template fonts: Roca Two Bold, Nunito SemiBold/Bold/ExtraBold, as OTF/TTF)
        <out_dir>/assets (the template's decorative PNG accents, keyed by original media name)
        <out_dir>/figs   (logos + figures pulled from the original VIVA deck)
"""
import io
import os
import re
import struct
import sys
import zipfile

from PIL import Image
from pptx import Presentation

FONT_NAMES = {21: "RocaTwo-Bold.otf", 22: "Nunito-SemiBold.ttf", 23: "Nunito-Bold.ttf", 24: "Nunito-ExtraBold.ttf"}
ACCENT_MEDIA = ["image1.png", "image5.png", "image7.png", "image9.png", "image11.png", "image14.png", "image16.png",
                "image18.png", "image21.png", "image23.png", "image26.png", "image28.png", "image31.png", "image33.png",
                "image35.png", "image38.png", "image40.png", "image43.png", "image46.png", "image49.png", "image51.png"]


def extract_fonts(template, out):
    os.makedirs(out, exist_ok=True)
    z = zipfile.ZipFile(template)
    for name in z.namelist():
        m = re.match(r"ppt/fonts/font(\d+)\.fntdata$", name)
        if not m:
            continue
        data = z.read(name)
        # EOT wrapper: FontDataSize is the second uint32 of the header; payload is the tail of the file
        font_size = struct.unpack("<I", data[4:8])[0]
        payload = data[-font_size:]
        target = FONT_NAMES.get(int(m.group(1)), f"font{m.group(1)}.ttf")
        open(os.path.join(out, target), "wb").write(payload)
        print("font", target, len(payload))
    # generic fallback used for width estimation of math glyphs
    dv = "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf"
    if os.path.exists(dv):
        import shutil
        shutil.copy(dv, os.path.join(out, "DejaVuSans.ttf"))


def extract_accents(template, out):
    os.makedirs(out, exist_ok=True)
    z = zipfile.ZipFile(template)
    for media in ACCENT_MEDIA:
        open(os.path.join(out, media), "wb").write(z.read("ppt/media/" + media))
    print("accents", len(ACCENT_MEDIA))


def extract_figs(viva, out):
    os.makedirs(out, exist_ok=True)
    p = Presentation(viva)
    wanted = {1: {89: "um5_logo", 90: "ensias_logo", 96: "ministry_logo"}, 26: {731: "c1_pipeline"}, 55: {1538: "dyhucog_arch"}}
    for i, s in enumerate(p.slides, 1):
        if i not in wanted:
            continue
        for sh in s.shapes:
            if sh.shape_type == 13 and sh.shape_id in wanted[i]:
                im = Image.open(io.BytesIO(sh.image.blob))
                base = wanted[i][sh.shape_id]
                if base == "ensias_logo":
                    im.convert("RGB").save(os.path.join(out, base + ".jpg"), quality=92)
                else:
                    im.convert("RGBA").save(os.path.join(out, base + ".png"))
                print("fig", base, im.size)


if __name__ == "__main__":
    template, viva, out = sys.argv[1:4]
    extract_fonts(template, os.path.join(out, "fonts"))
    extract_accents(template, os.path.join(out, "assets"))
    extract_figs(viva, os.path.join(out, "figs"))
