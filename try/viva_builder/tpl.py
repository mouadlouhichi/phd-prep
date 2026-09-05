"""
Design-system helpers for building the VIVA deck on top of the
"Beige Green Modern Illustrative Playful Thesis Defense Presentation" template.

Everything is drawn with native PowerPoint shapes using the template's own
palette, fonts (embedded Roca Two Bold / Nunito) and decorative PNG accents.
"""
import copy
import os
import re

from lxml import etree
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from fontTools.ttLib import TTFont

# --------------------------------------------------------------------------
# Geometry
# --------------------------------------------------------------------------
IN = 914400
SLIDE_W = 18288000
SLIDE_H = 10287000
MARGIN_L = 1028700            # 1.125 in  (template left margin)
MARGIN_R = SLIDE_W - 1028700  # 17259300 (template right edge)
CONTENT_W = MARGIN_R - MARGIN_L
HEADER_Y = 990600             # template header row


def emu(inches):
    return int(round(inches * IN))


# --------------------------------------------------------------------------
# Palette (from the template)
# --------------------------------------------------------------------------
BG = "FEF8F3"
INK = "111111"
GREEN = "124944"
YELLOW = "ECC665"
ORANGE = "DF8330"
WHITE = "FFFFFF"
# derived tints used for tables / cards
TINT = "F6ECDF"
TINT2 = "FBF4EB"
RULE = "E3D6C4"
GREEN_SOFT = "E6EEEA"

# --------------------------------------------------------------------------
# Fonts (typeface names exactly as embedded in the template)
# --------------------------------------------------------------------------
HERE = os.path.dirname(os.path.abspath(__file__))
FONT_DIR = os.environ.get("VIVA_FONT_DIR", "/tmp/viva_build/fonts")

FONTS = {
    "title": ("Roca Two Bold", "RocaTwo-Bold.otf"),
    "body": ("Nunito Semi-Bold", "Nunito-SemiBold.ttf"),
    "bold": ("Nunito Bold", "Nunito-Bold.ttf"),
    "xbold": ("Nunito Ultra-Bold", "Nunito-ExtraBold.ttf"),
    "math": ("Cambria Math", "DejaVuSans.ttf"),
}


class _Metrics:
    def __init__(self, path):
        self.f = TTFont(path)
        self.cmap = self.f.getBestCmap()
        self.hmtx = self.f["hmtx"]
        self.upm = self.f["head"].unitsPerEm

    def has(self, ch):
        return ord(ch) in self.cmap

    def width_units(self, text):
        w = 0
        for ch in text:
            g = self.cmap.get(ord(ch))
            if g is None:
                g = self.cmap.get(ord("n"))
            w += self.hmtx[g][0]
        return w


_METRICS = {}


def metrics(key):
    if key not in _METRICS:
        fn = FONTS[key][1]
        path = os.path.join(FONT_DIR, fn)
        if not os.path.exists(path):
            path = "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf"
        _METRICS[key] = _Metrics(path)
    return _METRICS[key]


def text_width_pt(text, key, size_pt, spc_pt=0.0):
    """Approximate rendered width in points."""
    m = metrics(key)
    w = 0.0
    for ch in text:
        mm = m if m.has(ch) else metrics("math")
        w += mm.width_units(ch) / mm.upm * size_pt
    return w * 1.03 + spc_pt * len(text)


def runs_width_pt(runs):
    w = 0.0
    for r in runs:
        f = 0.67 if r.baseline else 1.0
        spc = r.spc if r.spc is not None else (-0.05 * r.size if r.font == "title" else 0)
        w += text_width_pt(r.text, r.font, r.size * f, spc)
    return w


def wrap_count(text, key, size_pt, width_pt, spc_pt=0.0):
    """Greedy word-wrap: returns number of lines needed."""
    if width_pt <= 0:
        return 1
    lines = 1
    cur = 0.0
    space = text_width_pt(" ", key, size_pt, spc_pt)
    for word in text.split(" "):
        w = text_width_pt(word, key, size_pt, spc_pt)
        if cur == 0:
            cur = w
        elif cur + space + w <= width_pt:
            cur += space + w
        else:
            lines += 1
            cur = w
        while cur > width_pt:  # very long token
            lines += 1
            cur -= width_pt
    return lines


# --------------------------------------------------------------------------
# Rich text model
# --------------------------------------------------------------------------
class Run:
    def __init__(self, text, font="body", size=22, color=INK, spc=None, bold=True, italic=False, baseline=None):
        self.text = text
        self.font = font
        self.size = size
        self.color = color
        self.spc = spc
        self.bold = bold
        self.italic = italic
        self.baseline = baseline


class Para:
    def __init__(self, runs, align="l", lnspc=None, spcbef=0, level=0, bullet=None, bullet_color=None, marl=None, indent=None):
        if isinstance(runs, str):
            runs = [Run(runs)]
        elif isinstance(runs, Run):
            runs = [runs]
        self.runs = runs
        self.align = align
        self.lnspc = lnspc          # points (exact) ; None -> 1.36 x size
        self.spcbef = spcbef        # points
        self.level = level
        self.bullet = bullet        # bullet char or None
        self.bullet_color = bullet_color
        self.marl = marl
        self.indent = indent

    @property
    def size(self):
        return max(r.size for r in self.runs) if self.runs else 22

    @property
    def text(self):
        return "".join(r.text for r in self.runs)

    def line_pitch(self):
        return self.lnspc if self.lnspc else self.size * 1.36


NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"


def _sub(parent, tag, **attrs):
    el = etree.SubElement(parent, qn(tag))
    for k, v in attrs.items():
        el.set(k, str(v))
    return el


GREEK = set(range(0x0391, 0x03CA))
ARROWS = set(range(0x2190, 0x2196))
FALLBACK_SANS = "Calibri"        # theme minor font: Greek letters + arrows
FALLBACK_MATH = "Cambria Math"   # set-theory symbols etc.


def _fallback_for(cluster, m):
    """Return None (main font ok) or a fallback typeface for a grapheme cluster."""
    import unicodedata
    if all(m.has(ch) for ch in cluster):
        return None
    if len(cluster) > 1 and any(unicodedata.combining(ch) for ch in cluster[1:]):
        return FALLBACK_MATH
    o = ord(cluster[0])
    if o in GREEK or o in ARROWS:
        return FALLBACK_SANS
    return FALLBACK_MATH


def _split_math(text, key):
    """Split text into (segment, typeface-or-None) chunks; glyphs missing in the
    embedded Nunito/Roca subsets are routed to an Office-bundled fallback font.
    Whitespace is kept inside the current run so no run starts with a space."""
    import unicodedata
    m = metrics(key)
    clusters = []
    for ch in text:
        if clusters and unicodedata.combining(ch):
            clusters[-1] += ch
        else:
            clusters.append(ch)
    out = []
    cur, cur_face = "", None
    for cl in clusters:
        if cl.isspace():
            cur += cl
            continue
        face = _fallback_for(cl, m)
        if cur.strip() == "" or face == cur_face:
            cur += cl
            cur_face = face
        else:
            out.append((cur, cur_face))
            cur, cur_face = cl, face
    if cur:
        out.append((cur, cur_face))
    return out


def _add_run(p, run):
    for seg, face in _split_math(run.text, run.font):
        if not seg:
            continue
        r = _sub(p, "a:r")
        rpr = _sub(r, "a:rPr", lang="en-US", sz=int(round(run.size * 100)))
        if run.bold:
            rpr.set("b", "true")
        if run.italic:
            rpr.set("i", "true")
        spc = run.spc
        if spc is None and run.font == "title":
            spc = -0.05 * run.size
        if spc:
            rpr.set("spc", str(int(round(spc * 100))))
        if run.baseline:
            rpr.set("baseline", str(run.baseline))
        sf = _sub(rpr, "a:solidFill")
        _sub(sf, "a:srgbClr", val=run.color)
        typeface = face or FONTS[run.font][0]
        for tag in ("a:latin", "a:ea", "a:cs", "a:sym"):
            _sub(rpr, tag, typeface=typeface)
        t = _sub(r, "a:t")
        t.text = seg


def fill_text_frame(txBody, paras, anchor="t", insets=(0, 0, 0, 0), wrap=True):
    """Replace the content of an <p:txBody> with our paragraphs."""
    for child in list(txBody):
        txBody.remove(child)
    bp = _sub(txBody, "a:bodyPr", anchor=anchor, rtlCol="false",
              lIns=insets[0], tIns=insets[1], rIns=insets[2], bIns=insets[3],
              wrap="square" if wrap else "none")
    _sub(txBody, "a:lstStyle")
    if not paras:
        p = _sub(txBody, "a:p")
        _sub(p, "a:endParaRPr", lang="en-US", sz=1800)
    for para in paras:
        p = _sub(txBody, "a:p")
        ppr = _sub(p, "a:pPr", algn=para.align)
        marl = para.marl
        indent = para.indent
        if para.bullet:
            if marl is None:
                marl = emu(0.30) if para.level == 0 else emu(0.72)
            if indent is None:
                indent = -emu(0.30) if para.level == 0 else -emu(0.28)
        if marl is not None:
            ppr.set("marL", str(int(marl)))
        if indent is not None:
            ppr.set("indent", str(int(indent)))
        ls = _sub(ppr, "a:lnSpc")
        _sub(ls, "a:spcPts", val=int(round(para.line_pitch() * 100)))
        sb = _sub(ppr, "a:spcBef")
        _sub(sb, "a:spcPts", val=int(round(para.spcbef * 100)))
        if para.bullet:
            bc = _sub(ppr, "a:buClr")
            _sub(bc, "a:srgbClr", val=para.bullet_color or para.runs[0].color)
            _sub(ppr, "a:buSzPct", val=100000)
            _sub(ppr, "a:buFont", typeface="Arial")
            _sub(ppr, "a:buChar", char=para.bullet)
        else:
            _sub(ppr, "a:buNone")
        for run in para.runs:
            _add_run(p, run)
        if not para.runs:
            _sub(p, "a:endParaRPr", lang="en-US", sz=int(round(para.size * 100)))


def paras_height_pt(paras, width_emu, insets=(0, 0, 0, 0)):
    """Estimated height (points) of paragraphs laid out in a box of width_emu."""
    total = 0.0
    inner_w = (width_emu - insets[0] - insets[2]) / IN * 72
    for para in paras:
        marl = para.marl
        if marl is None:
            marl = (emu(0.30) if para.level == 0 else emu(0.72)) if para.bullet else 0
        w = inner_w - marl / IN * 72
        n = wrap_runs(para.runs, w) if para.text.strip() else 1
        total += n * para.line_pitch() + para.spcbef
    total += (insets[1] + insets[3]) / IN * 72
    return total


def wrap_runs(runs, width_pt):
    """Greedy word-wrap over rich runs (bold / sub-superscript aware)."""
    if width_pt <= 0:
        return 1
    # tokenise into words carrying their run style
    words = []  # list of (width, breakable_before)
    cur_w = 0.0
    started = False
    for r in runs:
        f = 0.67 if r.baseline else 1.0
        spc = r.spc if r.spc is not None else (-0.05 * r.size if r.font == "title" else 0)
        parts = re.split(r"(\s+)", r.text)
        for part in parts:
            if part == "":
                continue
            if part.isspace():
                if started:
                    words.append(cur_w)
                cur_w = 0.0
                started = False
                continue
            cur_w += text_width_pt(part, r.font, r.size * f, spc)
            started = True
    if started:
        words.append(cur_w)
    space = text_width_pt(" ", runs[0].font, runs[0].size) if runs else 4
    lines, cur = 1, 0.0
    for w in words:
        if cur == 0:
            cur = w
        elif cur + space + w <= width_pt:
            cur += space + w
        else:
            lines += 1
            cur = w
        while cur > width_pt:
            lines += 1
            cur -= width_pt
    return lines


# --------------------------------------------------------------------------
# Deck / slide primitives
# --------------------------------------------------------------------------
class Deck:
    def __init__(self, template_path):
        self.prs = Presentation(template_path)
        self.n_template = len(self.prs.slides)
        self.blank = [l for l in self.prs.slide_layouts if l.name == "Blank"][0]
        self.slides = []

    def new_slide(self, bg=BG):
        s = self.prs.slides.add_slide(self.blank)
        # remove anything cloned from the layout
        for shp in list(s.shapes):
            shp._element.getparent().remove(shp._element)
        cSld = s._element.find(qn("p:cSld"))
        bgel = etree.SubElement(cSld, qn("p:bg"))
        cSld.remove(bgel)
        cSld.insert(0, bgel)
        bgPr = _sub(bgel, "p:bgPr")
        sf = _sub(bgPr, "a:solidFill")
        _sub(sf, "a:srgbClr", val=bg)
        _sub(bgPr, "a:effectLst")
        self.slides.append(s)
        return s

    def finalize(self, out_path):
        sldIdLst = self.prs.slides._sldIdLst
        for sldId in list(sldIdLst)[: self.n_template]:
            self.prs.part.drop_rel(sldId.rId)
            sldIdLst.remove(sldId)
        self.prs.save(out_path)


def _strip_style(shape):
    st = shape._element.find(qn("p:style"))
    if st is not None:
        shape._element.remove(st)


def rrect(slide, x, y, w, h, fill=GREEN, line=None, line_w=2.25, radius=None, shape=MSO_SHAPE.ROUNDED_RECTANGLE, name=None):
    shp = slide.shapes.add_shape(shape, Emu(int(x)), Emu(int(y)), Emu(int(w)), Emu(int(h)))
    _strip_style(shp)
    if shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        if radius is None:
            radius = min(emu(0.32), min(w, h) / 2)
        shp.adjustments[0] = max(0.0, min(0.5, radius / float(min(w, h))))
    if fill:
        shp.fill.solid()
        shp.fill.fore_color.rgb = RGBColor.from_string(fill)
    else:
        shp.fill.background()
    if line:
        shp.line.color.rgb = RGBColor.from_string(line)
        shp.line.width = Pt(line_w)
    else:
        shp.line.fill.background()
    if name:
        shp.name = name
    # empty text frame -> keep clean
    tf = shp.text_frame
    fill_text_frame(shp._element.find(qn("p:txBody")), [], anchor="ctr")
    return shp


def rect(slide, x, y, w, h, fill=GREEN, line=None, line_w=2.25):
    return rrect(slide, x, y, w, h, fill=fill, line=line, line_w=line_w, shape=MSO_SHAPE.RECTANGLE)


def ellipse(slide, x, y, w, h, fill=YELLOW, line=None):
    return rrect(slide, x, y, w, h, fill=fill, line=line, shape=MSO_SHAPE.OVAL)


def textbox(slide, x, y, w, h, paras, anchor="t", insets=(0, 0, 0, 0), wrap=True, name=None):
    tb = slide.shapes.add_textbox(Emu(int(x)), Emu(int(y)), Emu(int(w)), Emu(int(h)))
    fill_text_frame(tb._element.find(qn("p:txBody")), paras, anchor=anchor, insets=insets, wrap=wrap)
    if name:
        tb.name = name
    return tb


def shape_text(shp, paras, anchor="ctr", insets=(emu(0.1), emu(0.05), emu(0.1), emu(0.05))):
    fill_text_frame(shp._element.find(qn("p:txBody")), paras, anchor=anchor, insets=insets)


def picture(slide, path, x, y, w=None, h=None, rot=0):
    kw = {}
    if w is not None:
        kw["width"] = Emu(int(w))
    if h is not None:
        kw["height"] = Emu(int(h))
    pic = slide.shapes.add_picture(path, Emu(int(x)), Emu(int(y)), **kw)
    if rot:
        pic.rotation = rot
    return pic


def line(slide, x1, y1, x2, y2, color=GREEN, w=2.25, arrow=False, dash=None):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Emu(int(x1)), Emu(int(y1)), Emu(int(x2)), Emu(int(y2)))
    _strip_style(c)
    c.line.color.rgb = RGBColor.from_string(color)
    c.line.width = Pt(w)
    ln = c._element.spPr.find(qn("a:ln"))
    if dash:
        _sub(ln, "a:prstDash", val=dash)
    if arrow:
        _sub(ln, "a:tailEnd", type="triangle", w="med", len="med")
    return c


def pill(slide, x, y, text, active=True, h=None, size=20, font="bold", pad=0.32, min_w=None, fill=None, text_color=None, line_color=None, w=None):
    """Rounded 'pill' label. Returns (shape, width)."""
    if h is None:
        h = emu(0.46)
    tw = text_width_pt(text, font, size) / 72 * IN
    width = int(tw + 2 * emu(pad))
    if min_w:
        width = max(width, int(min_w))
    if w:
        width = int(w)
    if fill is None:
        fill = GREEN if active else None
    if text_color is None:
        text_color = WHITE if active else GREEN
    if line_color is None:
        line_color = None if active else GREEN
    shp = rrect(slide, x, y, width, h, fill=fill, line=line_color, line_w=2.25, radius=h / 2)
    shape_text(shp, [Para([Run(text, font=font, size=size, color=text_color)], align="ctr", lnspc=size * 1.2)], anchor="ctr", insets=(emu(0.05), 0, emu(0.05), 0))
    return shp, width


def badge(slide, x, y, d, text, fill=YELLOW, color=INK, size=26, font="xbold"):
    shp = ellipse(slide, x, y, d, d, fill=fill)
    shape_text(shp, [Para([Run(text, font=font, size=size, color=color)], align="ctr", lnspc=size * 1.1)], anchor="ctr", insets=(0, 0, 0, 0))
    return shp


# --------------------------------------------------------------------------
# Tables
# --------------------------------------------------------------------------
NO_STYLE_GUID = "{2D5ABB26-0587-4C30-8999-92F81FD0307C}"


def _cell_border(cell, color=RULE, w_pt=1.0, sides="LRTB"):
    tcPr = cell._tc.get_or_add_tcPr()
    for side in "LRTB":
        tag = "a:ln" + side
        for el in tcPr.findall(qn(tag)):
            tcPr.remove(el)
    # insertion order: lnL lnR lnT lnB ... before fill elements
    fills = [c for c in tcPr if etree.QName(c).localname in ("noFill", "solidFill", "gradFill", "blipFill", "pattFill", "grpFill")]
    idx = tcPr.index(fills[0]) if fills else len(tcPr)
    for side in "LRTB":
        ln = etree.Element(qn("a:ln" + side), w=str(int(w_pt * 12700)), cap="flat", cmpd="sng", algn="ctr")
        if side in sides:
            sf = _sub(ln, "a:solidFill")
            _sub(sf, "a:srgbClr", val=color)
        else:
            _sub(ln, "a:noFill")
        _sub(ln, "a:prstDash", val="solid")
        tcPr.insert(idx, ln)
        idx += 1


def table(slide, x, y, w, rows, col_widths=None, header=True, size=17, row_h=None, head_fill=GREEN, head_color=WHITE, band=(WHITE, TINT), align=None, first_col_bold=True, pad=0.09, line_pitch=None, head_size=None, header_align=None):
    """rows: list of lists of str (or Para lists). Returns (graphicFrame, total_height_emu)."""
    nrows, ncols = len(rows), len(rows[0])
    if col_widths is None:
        col_widths = [w / ncols] * ncols
    else:
        tot = float(sum(col_widths))
        col_widths = [w * c / tot for c in col_widths]
    if head_size is None:
        head_size = size
    pitch = line_pitch or size * 1.3
    # compute row heights from wrapped text
    heights = []
    for ri, row in enumerate(rows):
        is_head = header and ri == 0
        sz = head_size if is_head else size
        maxlines = 1
        for ci, val in enumerate(row):
            txt = val if isinstance(val, str) else "".join(p.text for p in val)
            key = "bold" if (is_head or (first_col_bold and ci == 0)) else "body"
            avail = (col_widths[ci] - 2 * emu(pad)) / IN * 72
            n = 1
            for part in txt.split("\n"):
                n_part = wrap_count(part, key, sz, avail)
                n += n_part
            n -= 1
            maxlines = max(maxlines, n)
        h = maxlines * (head_size if is_head else size) * 1.3 + 2 * emu(pad) / IN * 72 + 4
        hh = emu(h / 72)
        if row_h:
            hh = max(hh, int(row_h))
        heights.append(hh)
    total_h = sum(heights)
    gf = slide.shapes.add_table(nrows, ncols, Emu(int(x)), Emu(int(y)), Emu(int(w)), Emu(int(total_h)))
    tbl = gf.table
    tblPr = tbl._tbl.tblPr
    for attr in ("firstRow", "bandRow", "firstCol", "lastRow", "lastCol", "bandCol"):
        tblPr.set(attr, "0")
    sid = tblPr.find(qn("a:tableStyleId"))
    if sid is None:
        sid = _sub(tblPr, "a:tableStyleId")
    sid.text = NO_STYLE_GUID
    for ci, cw in enumerate(col_widths):
        tbl.columns[ci].width = Emu(int(cw))
    for ri, row in enumerate(rows):
        tbl.rows[ri].height = Emu(int(heights[ri]))
        is_head = header and ri == 0
        for ci, val in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.margin_left = cell.margin_right = Emu(emu(pad))
            cell.margin_top = cell.margin_bottom = Emu(emu(0.05))
            cell.vertical_anchor = 1  # top? (MSO_ANCHOR.TOP=1) ; middle=3
            from pptx.enum.text import MSO_ANCHOR
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            if is_head:
                fillc, color, key = head_fill, head_color, "bold"
                sz = head_size
            else:
                fillc = band[(ri - (1 if header else 0)) % 2]
                color = INK
                key = "bold" if (first_col_bold and ci == 0) else "body"
                sz = size
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor.from_string(fillc)
            al = "l"
            if align:
                al = align[ci]
            if is_head and header_align:
                al = header_align[ci]
            if isinstance(val, str):
                paras = [Para([Run(part, font=key, size=sz, color=color)], align=al, lnspc=sz * 1.3) for part in val.split("\n")]
            else:
                paras = val
            fill_text_frame(cell._tc.get_or_add_txBody(), paras, anchor="ctr", insets=(0, 0, 0, 0))
            _cell_border(cell, color=RULE if not is_head else head_fill, w_pt=1.0)
    return gf, total_h


# --------------------------------------------------------------------------
# Decorative accents (template media, copied to assets/)
# --------------------------------------------------------------------------
ASSETS = os.environ.get("VIVA_ASSET_DIR", "/tmp/viva_build/assets")
ACCENTS = {
    "arrow_black": "image1.png",
    "arrow_white": "image21.png",
    "asterisk_y": "image5.png",
    "asterisk_o": "image26.png",
    "sparkle_o": "image9.png",
    "sparkle_y": "image11.png",
    "comb_y": "image14.png",
    "flower_o": "image16.png",
    "flower_y": "image40.png",
    "trident_y": "image18.png",
    "dots_y": "image23.png",
    "half_y": "image28.png",
    "uparrow_g": "image31.png",
    "fan_y": "image33.png",
    "candy": "image35.png",
    "v_o": "image38.png",
    "squiggle_g": "image43.png",
    "crown_g": "image46.png",
    "rainbow_o": "image49.png",
    "grid": "image51.png",
    "m_g": "image7.png",
}


def accent(slide, key, x, y, w, rot=0):
    from PIL import Image
    path = os.path.join(ASSETS, ACCENTS[key])
    im = Image.open(path)
    ratio = im.height / im.width
    return picture(slide, path, x, y, w=w, h=w * ratio, rot=rot)
