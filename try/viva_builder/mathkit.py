"""
Native PowerPoint equations for the VIVA deck.

A small LaTeX subset is converted to OMML (Office Math Markup) so every formula is a
*real* PowerPoint equation: typeset by PowerPoint in Cambria Math, editable with the
built-in equation editor and rendered crisply at any zoom.  Exactly like PowerPoint
itself, each equation shape is written as

    <mc:AlternateContent>
      <mc:Choice Requires="a14">  ... text box with <a14:m><m:oMathPara> ...   (PowerPoint 2010+)
      <mc:Fallback>               ... same box filled with a rendered PNG        (other viewers)

The PNG fallback (and all size measurements used for layout) come from matplotlib's
mathtext engine with the STIX fonts, which is metrically close to Cambria Math.

Pipeline:  LaTeX  --latex2mathml-->  MathML  --this module-->  OMML
           LaTeX  --matplotlib.mathtext-->  (width, height, depth) + fallback PNG
"""
import hashlib
import os
import re
import tempfile

import matplotlib

matplotlib.use("Agg")
matplotlib.rcParams["mathtext.fontset"] = "stix"
matplotlib.rcParams["mathtext.default"] = "it"
matplotlib.rcParams["font.family"] = "STIXGeneral"      # \text{} in the fallback pictures

EQ_MODE = os.environ.get("VIVA_EQ_MODE", "native")     # native = OMML + picture fallback | picture = pictures only

from latex2mathml.converter import convert as _latex_to_mathml  # noqa: E402
from lxml import etree  # noqa: E402
from matplotlib.figure import Figure  # noqa: E402
from matplotlib.font_manager import FontProperties  # noqa: E402
from matplotlib.mathtext import MathTextParser  # noqa: E402
from pptx.oxml.ns import qn  # noqa: E402

NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
NS_P = "http://schemas.openxmlformats.org/presentationml/2006/main"
NS_R = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
NS_M = "http://schemas.openxmlformats.org/officeDocument/2006/math"
NS_A14 = "http://schemas.microsoft.com/office/drawing/2010/main"
NS_MC = "http://schemas.openxmlformats.org/markup-compatibility/2006"
XML_SPACE = "{http://www.w3.org/XML/1998/namespace}space"

MATH_FONT = "Cambria Math"          # ships with every Office since 2007
PT = 12700                          # EMU per point
FALLBACK_DPI = 220

_NARY = {"∑": "∑", "∏": "∏", "∫": "∫", "⋃": "⋃", "⋂": "⋂"}
_ACCENTS = {"^": "\u0302", "ˆ": "\u0302", "¯": "\u0305", "‾": "\u0305", "~": "\u0303", "˜": "\u0303",
            "→": "\u20D7", "˙": "\u0307", "¨": "\u0308", "ˇ": "\u030C"}
_SPACES = {"negativethinmathspace": "", "negativemediummathspace": "", "negativethickmathspace": "",
           "0.167em": "\u2009", "0.222em": "\u2005", "0.278em": "\u2005", "1em": "\u2003", "2em": "\u2003\u2003"}
_CHAR_FIX = {"\u29F5": "\u2216",    # latex2mathml emits U+29F5 for \setminus; use SET MINUS
             "\u007E": "\u223C",    # \sim
             "\u00A0": " "}

# Mathematical Alphanumeric Symbols -> (base letter, m:scr, m:sty).  Office stores the base letter and
# the script/style flags; the same glyph is produced by Cambria Math.
_ALPHA_BLOCKS = [  # (start of A..Z a..z run, scr, sty)
    (0x1D400, None, "b"), (0x1D434, None, "i"), (0x1D468, None, "bi"),
    (0x1D49C, "script", "p"), (0x1D4D0, "script", "b"),
    (0x1D504, "fraktur", "p"), (0x1D56C, "fraktur", "b"),
    (0x1D538, "double-struck", "p"),
    (0x1D5A0, "sans-serif", "p"), (0x1D5D4, "sans-serif", "b"),
    (0x1D670, "monospace", "p"),
]
_LETTERLIKE = {  # BMP letter-like symbols that are the "holes" of the plane-1 blocks
    "\u2102": ("C", "double-struck"), "\u210D": ("H", "double-struck"), "\u2115": ("N", "double-struck"),
    "\u2119": ("P", "double-struck"), "\u211A": ("Q", "double-struck"), "\u211D": ("R", "double-struck"),
    "\u2124": ("Z", "double-struck"),
    "\u212C": ("B", "script"), "\u2130": ("E", "script"), "\u2131": ("F", "script"), "\u210B": ("H", "script"),
    "\u2110": ("I", "script"), "\u2112": ("L", "script"), "\u2133": ("M", "script"), "\u211B": ("R", "script"),
    "\u212F": ("e", "script"), "\u210A": ("g", "script"), "\u2134": ("o", "script"),
    "\u212D": ("C", "fraktur"), "\u210C": ("H", "fraktur"), "\u2111": ("I", "fraktur"), "\u211C": ("R", "fraktur"),
    "\u2128": ("Z", "fraktur"),
}


def _decode_alpha(ch):
    """Return (letter, scr, sty) for a math-alphanumeric char, else None."""
    if ch in _LETTERLIKE:
        base, scr = _LETTERLIKE[ch]
        return base, scr, "p"
    o = ord(ch)
    for start, scr, sty in _ALPHA_BLOCKS:
        if start <= o < start + 52:
            k = o - start
            base = chr(ord("A") + k) if k < 26 else chr(ord("a") + k - 26)
            return base, scr, sty
    return None


def _m(tag):
    return "{%s}%s" % (NS_M, tag)


def _local(tag):
    return tag.rsplit("}", 1)[-1]


# ----------------------------------------------------------------------------
# Intermediate representation (lists of nodes) built from MathML
# ----------------------------------------------------------------------------
class _R:
    """A math run.  style: None = math italic (variables), 'p' = plain/upright, 'b', 'bi',
    'nor' = normal text; scr = None | script | fraktur | double-struck | sans-serif | monospace."""

    def __init__(self, text, style=None, scr=None):
        self.text, self.style, self.scr = text, style, scr


class _N:
    """Structural node: kind in sSub sSup sSubSup f nary d acc rad, args = dict of lists."""

    def __init__(self, kind, **args):
        self.kind, self.args = kind, args


def _merge(nodes):
    out = []
    for n in nodes:
        if isinstance(n, _R) and out and isinstance(out[-1], _R) and out[-1].style == n.style and out[-1].scr == n.scr:
            out[-1] = _R(out[-1].text + n.text, n.style, n.scr)
        else:
            out.append(n)
    return out


def _is_prime_sup(node):
    return _local(node.tag) == "msup" and len(node) == 2 and _local(node[1].tag) == "mi" and (node[1].text or "") in ("′", "″", "'")


def _is_nary_base(node):
    """<mo>∑</mo> or <mrow><mo>∑</mo></mrow> (the latter = {\\sum}: limits at the side)."""
    t = _local(node.tag)
    if t == "mo" and (node.text or "").strip() in _NARY:
        return "undOvr"
    if t == "mrow" and len(node) == 1 and _local(node[0].tag) == "mo" and (node[0].text or "").strip() in _NARY:
        return "subSup"
    return None


def _nary_char(node):
    return (node.text or node[0].text).strip()


class _Builder:
    def conv_children(self, node):
        kids = list(node)
        out, i = [], 0
        while i < len(kids):
            k = kids[i]
            tag = _local(k.tag)
            base = k[0] if tag in ("msub", "msup", "msubsup", "munder", "mover", "munderover") and len(k) else None
            lim = _is_nary_base(base) if base is not None and tag != "mover" else None
            if lim or _is_nary_base(k):
                arg = kids[i + 1] if i + 1 < len(kids) else None
                out.append(self.nary(k, arg, lim))
                i += 2 if arg is not None else 1
                continue
            out += self.conv(k)
            i += 1
        return _merge(out)

    def nary(self, node, arg, lim):
        tag = _local(node.tag)
        sub = sup = []
        if tag in ("msub", "munder"):
            base, sub = node[0], self.conv(node[1])
        elif tag in ("msup",):
            base, sup = node[0], self.conv(node[1])
        elif tag in ("msubsup", "munderover"):
            base, sub, sup = node[0], self.conv(node[1]), self.conv(node[2])
        else:
            base = node
        e = self.conv(arg) if arg is not None else []
        return _N("nary", chr=_nary_char(base), sub=sub, sup=sup, e=e, lim=lim or "undOvr")

    def conv(self, node):
        tag = _local(node.tag)
        text = node.text or ""
        for a, b in _CHAR_FIX.items():
            text = text.replace(a, b)
        if tag in ("math", "mstyle", "mpadded", "mphantom", "semantics"):
            return self.conv_children(node)
        if tag == "mrow":
            kids = list(node)
            if len(kids) >= 2 and _local(kids[0].tag) == "mo" and kids[0].get("form") == "prefix" and kids[0].get("fence") == "true" \
                    and _local(kids[-1].tag) == "mo" and kids[-1].get("form") == "postfix":
                inner = etree.Element("mrow")
                for c in kids[1:-1]:
                    inner.append(c)
                beg = (kids[0].text or "").strip()
                end = (kids[-1].text or "").strip()
                return [_N("d", beg=beg, end=end, e=self.conv_children(inner))]
            return self.conv_children(node)
        if tag == "mi":
            if not text.strip():
                return []
            dec = _decode_alpha(text) if len(text) == 1 else None
            if dec:
                base, scr, sty = dec
                return [_R(base, None if sty == "i" else sty, scr)]
            if not any(ch.isalpha() for ch in text):      # ⟨ ∅ … : symbols, never italic
                return [_R(text, "p")]
            upright = node.get("mathvariant") == "normal" or (len(text) > 1 and text.isalpha())
            return [_R(text, "p" if upright else None)]
        if tag == "mo":
            if not text.strip():
                return []
            return [_R(text.strip() if text.strip().isalpha() else text, "p")]   # operators / \operatorname{}
        if tag == "mn":
            return [_R(text, "p")]
        if tag == "mtext":
            return [_R(text, "nor")] if text else []
        if tag == "mspace":
            return [_R(_SPACES.get(node.get("width", ""), " "), "p")] if _SPACES.get(node.get("width", ""), " ") else []
        if tag in ("msub", "msup", "msubsup"):
            if _is_prime_sup(node):
                return self.conv(node[0]) + [_R("′")]
            e = self.conv(node[0])
            if tag == "msub":
                return [_N("sSub", e=e, sub=self.conv(node[1]))]
            if tag == "msup":
                return [_N("sSup", e=e, sup=self.conv(node[1]))]
            return [_N("sSubSup", e=e, sub=self.conv(node[1]), sup=self.conv(node[2]))]
        if tag == "mfrac":
            return [_N("f", num=self.conv(node[0]), den=self.conv(node[1]))]
        if tag == "mover":
            ch = (node[1].text or "").strip() if len(node) > 1 else ""
            if node.get("accent") == "true" or ch in _ACCENTS:
                return [_N("acc", chr=_ACCENTS.get(ch, ch or "\u0302"), e=self.conv(node[0]))]
            return [_N("limUpp", e=self.conv(node[0]), lim=self.conv(node[1]))]
        if tag == "munder":
            return [_N("limLow", e=self.conv(node[0]), lim=self.conv(node[1]))]
        if tag == "msqrt":
            return [_N("rad", e=self.conv_children(node), deg=[])]
        if tag == "mroot":
            return [_N("rad", e=self.conv(node[0]), deg=self.conv(node[1]))]
        if tag == "mfenced":
            return [_N("d", beg=node.get("open", "("), end=node.get("close", ")"), e=self.conv_children(node))]
        raise ValueError("unsupported MathML element in equation: <%s>" % tag)


def latex_to_ir(tex):
    mml = _latex_to_mathml(tex)
    root = etree.fromstring(mml.encode("utf8"))
    return _Builder().conv(root)


# ----------------------------------------------------------------------------
# IR -> OMML
# ----------------------------------------------------------------------------
class _Emitter:
    def __init__(self, size_pt, color):
        self.sz = str(int(round(size_pt * 100)))
        self.color = color

    def rpr(self, italic=True, bold=False):
        el = etree.Element(qn("a:rPr"), lang="en-US", sz=self.sz)
        if bold:
            el.set("b", "1")
        if italic:
            el.set("i", "1")
        sf = etree.SubElement(el, qn("a:solidFill"))
        etree.SubElement(sf, qn("a:srgbClr"), val=self.color)
        etree.SubElement(el, qn("a:latin"), typeface=MATH_FONT, panose="02040503050406030204", pitchFamily="18", charset="0")
        return el

    def ctrl(self, parent):
        c = etree.SubElement(parent, _m("ctrlPr"))
        c.append(self.rpr())
        return c

    def run(self, r):
        el = etree.Element(_m("r"))
        if r.style or r.scr:
            rp = etree.SubElement(el, _m("rPr"))
            if r.style == "nor":
                etree.SubElement(rp, _m("nor"))
            else:
                if r.scr:
                    etree.SubElement(rp, _m("scr"), {_m("val"): r.scr})
                if r.style:
                    etree.SubElement(rp, _m("sty"), {_m("val"): r.style})
        el.append(self.rpr(italic=r.style in (None, "i", "bi"), bold=r.style in ("b", "bi")))
        t = etree.SubElement(el, _m("t"))
        t.text = r.text
        t.set(XML_SPACE, "preserve")
        return el

    def arg(self, parent, tag, nodes):
        el = etree.SubElement(parent, _m(tag))
        for x in self.emit(nodes):
            el.append(x)
        return el

    def emit(self, nodes):
        out = []
        for n in _merge(nodes):
            if isinstance(n, _R):
                out.append(self.run(n))
                continue
            k, a = n.kind, n.args
            el = etree.Element(_m(k))
            if k == "sSub":
                self.ctrl(etree.SubElement(el, _m("sSubPr")))
                self.arg(el, "e", a["e"]); self.arg(el, "sub", a["sub"])
            elif k == "sSup":
                self.ctrl(etree.SubElement(el, _m("sSupPr")))
                self.arg(el, "e", a["e"]); self.arg(el, "sup", a["sup"])
            elif k == "sSubSup":
                self.ctrl(etree.SubElement(el, _m("sSubSupPr")))
                self.arg(el, "e", a["e"]); self.arg(el, "sub", a["sub"]); self.arg(el, "sup", a["sup"])
            elif k == "f":
                self.ctrl(etree.SubElement(el, _m("fPr")))
                self.arg(el, "num", a["num"]); self.arg(el, "den", a["den"])
            elif k == "nary":
                pr = etree.SubElement(el, _m("naryPr"))
                etree.SubElement(pr, _m("chr"), {_m("val"): a["chr"]})
                etree.SubElement(pr, _m("limLoc"), {_m("val"): a["lim"]})
                if not a["sub"]:
                    etree.SubElement(pr, _m("subHide"), {_m("val"): "1"})
                if not a["sup"]:
                    etree.SubElement(pr, _m("supHide"), {_m("val"): "1"})
                self.ctrl(pr)
                self.arg(el, "sub", a["sub"]); self.arg(el, "sup", a["sup"]); self.arg(el, "e", a["e"])
            elif k == "d":
                pr = etree.SubElement(el, _m("dPr"))
                if a["beg"] != "(":
                    etree.SubElement(pr, _m("begChr"), {_m("val"): a["beg"]})
                if a["end"] != ")":
                    etree.SubElement(pr, _m("endChr"), {_m("val"): a["end"]})
                self.ctrl(pr)
                self.arg(el, "e", a["e"])
            elif k == "acc":
                pr = etree.SubElement(el, _m("accPr"))
                etree.SubElement(pr, _m("chr"), {_m("val"): a["chr"]})
                self.ctrl(pr)
                self.arg(el, "e", a["e"])
            elif k == "rad":
                pr = etree.SubElement(el, _m("radPr"))
                if not a["deg"]:
                    etree.SubElement(pr, _m("degHide"), {_m("val"): "1"})
                self.ctrl(pr)
                self.arg(el, "deg", a["deg"]); self.arg(el, "e", a["e"])
            elif k in ("limLow", "limUpp"):
                self.ctrl(etree.SubElement(el, _m(k + "Pr")))
                self.arg(el, "e", a["e"]); self.arg(el, "lim", a["lim"])
            else:
                raise ValueError(k)
            out.append(el)
        return out


def omml_paragraph(tex, size_pt, color, align="ctr"):
    """<a:p> holding one display equation (a14:m / m:oMathPara)."""
    ir = latex_to_ir(tex)
    em = _Emitter(size_pt, color)
    p = etree.Element(qn("a:p"))
    etree.SubElement(p, qn("a:pPr"), algn=align)
    a14m = etree.SubElement(p, "{%s}m" % NS_A14, nsmap={"a14": NS_A14, "m": NS_M})
    para = etree.SubElement(a14m, _m("oMathPara"))
    ppr2 = etree.SubElement(para, _m("oMathParaPr"))
    etree.SubElement(ppr2, _m("jc"), {_m("val"): {"ctr": "centerGroup", "l": "left", "r": "right"}[align]})
    om = etree.SubElement(para, _m("oMath"))
    for el in em.emit(ir):
        om.append(el)
    etree.SubElement(p, qn("a:endParaRPr"), lang="en-US", sz=str(int(round(size_pt * 100))))
    return p


# ----------------------------------------------------------------------------
# Measurement + fallback rendering (matplotlib mathtext, STIX)
# ----------------------------------------------------------------------------
_parser = MathTextParser("path")
_MEASURE = {}


def measure(tex, size_pt):
    """(width, height above baseline, depth below baseline) in points."""
    key = (tex, round(size_pt, 2))
    if key not in _MEASURE:
        w, h, d, _, _ = _parser.parse("$" + tex + "$", dpi=72, prop=FontProperties(size=size_pt))
        _MEASURE[key] = (float(w), float(h), float(d))
    return _MEASURE[key]


def box_size(tex, size_pt, headroom=1.22, min_lines=1.5):
    """Generous (w, h) in points for a box holding this equation.  Cambria Math (PowerPoint) runs
    ~15 % wider than the STIX fonts used for measuring, so allow headroom; never tighter than 1.5 lines."""
    w, h, d = measure(tex, size_pt)
    return w * 1.22 + size_pt * 0.6, max((h + d) * headroom + size_pt * 0.3, size_pt * min_lines)


def fit_size(tex, size_pt, max_w_pt, max_h_pt=None, min_size=12):
    """Largest size <= size_pt whose box fits the given width / height."""
    s = size_pt
    while s > min_size:
        w, h = box_size(tex, s)
        if w <= max_w_pt and (max_h_pt is None or h <= max_h_pt):
            break
        s -= 0.5
    return s


_PNG_CACHE = {}


def render_png(tex, size_pt, color, w_pt, h_pt, align="ctr", out_dir=None):
    """Transparent PNG of the equation laid out in a w x h (points) canvas; returns the path."""
    out_dir = out_dir or os.path.join(tempfile.gettempdir(), "viva_eq_png")
    os.makedirs(out_dir, exist_ok=True)
    key = hashlib.sha1(("%s|%s|%s|%.1f|%.1f|%s" % (tex, size_pt, color, w_pt, h_pt, align)).encode()).hexdigest()[:16]
    path = os.path.join(out_dir, key + ".png")
    if not os.path.exists(path):
        fig = Figure(figsize=(w_pt / 72.0, h_pt / 72.0), dpi=FALLBACK_DPI)
        x, ha = {"ctr": (0.5, "center"), "l": (0.005, "left"), "r": (0.995, "right")}[align]
        fig.text(x, 0.5, "$" + tex + "$", fontsize=size_pt, color="#" + color, ha=ha, va="center")
        fig.savefig(path, dpi=FALLBACK_DPI, transparent=True)
    return path


# ----------------------------------------------------------------------------
# Slide integration
# ----------------------------------------------------------------------------
def add_equation(slide, x, y, w, h, tex, size_pt=24, color="111111", align="ctr", anchor="ctr", name="Equation"):
    """Add a native equation text box (x, y, w, h in EMU) with a PNG fallback, PowerPoint-style."""
    from pptx.util import Emu
    if EQ_MODE == "picture":
        png = render_png(tex, size_pt, color, w / PT, h / PT, align)
        pic = slide.shapes.add_picture(png, Emu(int(x)), Emu(int(y)), Emu(int(w)), Emu(int(h)))
        pic.name = name
        return pic
    tb = slide.shapes.add_textbox(Emu(int(x)), Emu(int(y)), Emu(int(w)), Emu(int(h)))
    tb.name = name
    txBody = tb._element.find(qn("p:txBody"))
    for child in list(txBody):
        txBody.remove(child)
    etree.SubElement(txBody, qn("a:bodyPr"), wrap="none", rtlCol="0", anchor=anchor, lIns="0", tIns="0", rIns="0", bIns="0")
    etree.SubElement(txBody, qn("a:lstStyle"))
    txBody.append(omml_paragraph(tex, size_pt, color, align))

    # fallback: same box, filled with the rendered picture
    png = render_png(tex, size_pt, color, w / PT, h / PT, align)
    _, rId = slide.part.get_or_add_image_part(png)
    sp = tb._element
    fb_sp = etree.fromstring(
        '<p:sp xmlns:p="%s" xmlns:a="%s" xmlns:r="%s">'
        '<p:nvSpPr><p:cNvPr id="%d" name="%s"/><p:cNvSpPr txBox="1"/><p:nvPr/></p:nvSpPr>'
        '<p:spPr><a:xfrm><a:off x="%d" y="%d"/><a:ext cx="%d" cy="%d"/></a:xfrm>'
        '<a:prstGeom prst="rect"><a:avLst/></a:prstGeom>'
        '<a:blipFill><a:blip r:embed="%s"/><a:stretch><a:fillRect/></a:stretch></a:blipFill></p:spPr>'
        '<p:txBody><a:bodyPr/><a:lstStyle/><a:p><a:endParaRPr lang="en-US"/></a:p></p:txBody></p:sp>'
        % (NS_P, NS_A, NS_R, tb.shape_id, name, int(x), int(y), int(w), int(h), rId))

    parent = sp.getparent()
    idx = parent.index(sp)
    parent.remove(sp)
    ac = etree.Element("{%s}AlternateContent" % NS_MC, nsmap={"mc": NS_MC})
    choice = etree.SubElement(ac, "{%s}Choice" % NS_MC, nsmap={"a14": NS_A14})
    choice.set("Requires", "a14")
    choice.append(sp)
    fb = etree.SubElement(ac, "{%s}Fallback" % NS_MC)
    fb.append(fb_sp)
    parent.insert(idx, ac)
    return tb


def omml_xml(tex, size_pt=24, color="111111"):
    """Debug helper: pretty OMML for a formula."""
    return etree.tostring(omml_paragraph(tex, size_pt, color), pretty_print=True).decode()
